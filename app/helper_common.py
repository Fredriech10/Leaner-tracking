import base64
import mimetypes
import os
import re
import uuid
import zipfile
from io import BytesIO
from xml.etree import ElementTree as ET

from markupsafe import escape


INTERACTIVE_LEARNING_DIR = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "INTERACTIVE LEARNING",
)
LESSON_ASSET_DIR = os.path.join("static", "uploads", "lesson_assets")
MAX_LESSON_IMAGE_DIMENSION = 1600
DATA_URI_IMAGE_RE = re.compile(r"data:(image/[a-zA-Z0-9.+-]+);base64,([A-Za-z0-9+/=\r\n]+)")


def safe_int(value, default=0):
    try:
        return int(value)
    except (TypeError, ValueError):
        return default


def save_lesson_asset(raw_image, mime_type=None):
    os.makedirs(LESSON_ASSET_DIR, exist_ok=True)
    try:
        from PIL import Image

        image = Image.open(BytesIO(raw_image))
        image.thumbnail((MAX_LESSON_IMAGE_DIMENSION, MAX_LESSON_IMAGE_DIMENSION))
        filename = f"{uuid.uuid4().hex}.webp"
        filepath = os.path.join(LESSON_ASSET_DIR, filename)
        if image.mode not in ("RGB", "RGBA"):
            image = image.convert("RGBA" if "A" in image.getbands() else "RGB")
        save_kwargs = {"format": "WEBP", "quality": 82, "method": 6}
        if image.mode == "RGBA":
            save_kwargs["lossless"] = False
        image.save(filepath, **save_kwargs)
    except Exception:
        ext = mimetypes.guess_extension(mime_type or "") or ".png"
        if ext == ".jpe":
            ext = ".jpg"
        filename = f"{uuid.uuid4().hex}{ext}"
        filepath = os.path.join(LESSON_ASSET_DIR, filename)
        with open(filepath, "wb") as asset_file:
            asset_file.write(raw_image)
    return f"/static/uploads/lesson_assets/{filename}"


def save_data_uri_image(data_uri):
    match = DATA_URI_IMAGE_RE.fullmatch(data_uri.strip())
    if not match:
        return data_uri
    mime_type, encoded = match.groups()
    raw_image = base64.b64decode(encoded)
    return save_lesson_asset(raw_image, mime_type)


def externalize_data_uri_images(html_or_uri):
    if not html_or_uri:
        return html_or_uri
    return DATA_URI_IMAGE_RE.sub(lambda match: save_data_uri_image(match.group(0)), html_or_uri)


def parse_module_names(raw_value):
    raw_value = (raw_value or "").replace("\r", "\n")
    parts = []
    for chunk in raw_value.replace(";", ",").split(","):
        for line in chunk.split("\n"):
            cleaned = line.strip()
            if cleaned:
                parts.append(cleaned)
    seen = set()
    ordered = []
    for item in parts:
        key = item.lower()
        if key in seen:
            continue
        seen.add(key)
        ordered.append(item)
    return ordered


def normalize_question_bank_group_text(question_text):
    text = " ".join((question_text or "").strip().split())
    if not text:
        return ""
    text = re.sub(r"\s*\((?:case|scenario|batch)\s+[^)]*\)\s*$", "", text, flags=re.IGNORECASE)
    return text.strip().lower()


def normalize_relative_path(path):
    return path.replace("/", os.sep).replace("\\", os.sep).strip()


def get_interactive_learning_files():
    files = []
    if not os.path.isdir(INTERACTIVE_LEARNING_DIR):
        return files

    for root, _, filenames in os.walk(INTERACTIVE_LEARNING_DIR):
        for filename in filenames:
            lower = filename.lower()
            if lower.startswith("~$"):
                continue
            if lower.endswith((".ppt", ".pptx", ".pps", ".ppsx")):
                full_path = os.path.join(root, filename)
                rel_path = os.path.relpath(full_path, INTERACTIVE_LEARNING_DIR)
                files.append({
                    "relative_path": rel_path.replace("\\", "/"),
                    "display_name": os.path.splitext(filename)[0],
                    "folder": os.path.dirname(rel_path).replace("\\", "/"),
                })

    files.sort(key=lambda item: (item["folder"].lower(), item["display_name"].lower()))
    return files


def resolve_interactive_learning_path(relative_path):
    if not relative_path:
        return None

    normalized = os.path.normpath(os.path.join(INTERACTIVE_LEARNING_DIR, normalize_relative_path(relative_path)))
    base_dir = os.path.normpath(INTERACTIVE_LEARNING_DIR)

    if not normalized.startswith(base_dir):
        return None
    if not os.path.isfile(normalized):
        return None
    return normalized


def extract_pptx_slides(file_path):
    lower = file_path.lower()
    if not lower.endswith((".pptx", ".ppsx")):
        return []

    namespace = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    slides = []

    with zipfile.ZipFile(file_path) as archive:
        slide_names = [
            name for name in archive.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        ]
        slide_names.sort(key=lambda item: int(re.search(r"slide(\d+)\.xml$", item).group(1)))

        for index, slide_name in enumerate(slide_names, start=1):
            root = ET.fromstring(archive.read(slide_name))
            text_runs = [
                (node.text or "").strip()
                for node in root.findall(".//a:t", namespace)
                if (node.text or "").strip()
            ]
            slides.append({
                "number": index,
                "title": text_runs[0] if text_runs else f"Slide {index}",
                "text_runs": text_runs,
                "body": text_runs[1:] if len(text_runs) > 1 else [],
            })

    return slides


def pptx_to_content_slide_html(uploaded_file):
    from pptx import Presentation

    prs = Presentation(uploaded_file)
    slide_w = float(prs.slide_width)
    slide_h = float(prs.slide_height)
    html_slides = []

    for slide in prs.slides:
        boxes = []
        for shape in slide.shapes:
            left = max(0, (float(shape.left) / slide_w) * 100)
            top = max(0, (float(shape.top) / slide_h) * 100)
            width = max(5, (float(shape.width) / slide_w) * 100)
            height = max(5, (float(shape.height) / slide_h) * 100)
            style = f"left:{left:.2f}%;top:{top:.2f}%;width:{width:.2f}%;height:{height:.2f}%;"

            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if not text:
                    continue
                paragraphs = "".join(
                    f"<p>{escape(line)}</p>"
                    for line in text.splitlines()
                    if line.strip()
                )
                boxes.append(f'<div class="slide-box text-box" style="{style}">{paragraphs}</div>')
                continue

            if hasattr(shape, "image"):
                image = shape.image
                ext = (image.ext or "png").lower()
                mime = mimetypes.types_map.get(f".{ext}", "image/png")
                encoded = base64.b64encode(image.blob).decode("ascii")
                src = f"data:{mime};base64,{encoded}"
                boxes.append(f'<div class="slide-box image-box" style="{style}"><img src="{src}" alt=""></div>')

        if boxes:
            html_slides.append("".join(boxes))

    return html_slides
