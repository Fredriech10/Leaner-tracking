from flask import Flask, request, redirect, url_for, render_template, session, flash, jsonify
from markupsafe import escape
from datetime import datetime, timedelta
import threading
import time
import os
import sqlite3
import random
import csv
import json
import tempfile
import zipfile
import base64
import mimetypes
import re
import uuid
from collections import defaultdict
from pathlib import Path
from flask import send_file
import pandas as pd
import io
from io import BytesIO
from xml.etree import ElementTree as ET
from werkzeug.utils import secure_filename

TIMEOUT = 60
LESSON_SLIDE_TYPES = ("content_slide", "title_slide", "heading_slide")
DATA_URI_IMAGE_RE = re.compile(r"data:(image/[a-zA-Z0-9.+-]+);base64,([A-Za-z0-9+/=\r\n]+)")
LESSON_ASSET_DIR = os.path.join("static", "uploads", "lesson_assets")
MAX_LESSON_IMAGE_DIMENSION = 1600

def save_lesson_asset(raw_image, mime_type=None):
    os.makedirs(LESSON_ASSET_DIR, exist_ok=True)
    try:
        from PIL import Image

        image = Image.open(BytesIO(raw_image))
        image.thumbnail((MAX_LESSON_IMAGE_DIMENSION, MAX_LESSON_IMAGE_DIMENSION))
        filename = f"{uuid.uuid4().hex}.webp"
        filepath = os.path.join(LESSON_ASSET_DIR, filename)
        if image.mode not in ("RGB", "RGBA"):
            image = image.convert("RGBA" if "A" in image.getbands() else "RGB")
        save_kwargs = {"format": "WEBP", "quality": 82, "method": 6}
        if image.mode == "RGBA":
            save_kwargs["lossless"] = False
        image.save(filepath, **save_kwargs)
    except Exception:
        ext = mimetypes.guess_extension(mime_type or "") or ".png"
        if ext == ".jpe":
            ext = ".jpg"
        filename = f"{uuid.uuid4().hex}{ext}"
        filepath = os.path.join(LESSON_ASSET_DIR, filename)
        with open(filepath, "wb") as asset_file:
            asset_file.write(raw_image)
    return f"/static/uploads/lesson_assets/{filename}"

def save_data_uri_image(data_uri):
    match = DATA_URI_IMAGE_RE.fullmatch(data_uri.strip())
    if not match:
        return data_uri
    mime_type, encoded = match.groups()
    raw_image = base64.b64decode(encoded)
    return save_lesson_asset(raw_image, mime_type)

def externalize_data_uri_images(html_or_uri):
    if not html_or_uri:
        return html_or_uri
    return DATA_URI_IMAGE_RE.sub(lambda match: save_data_uri_image(match.group(0)), html_or_uri)

def safe_int(value, default=0):
    try:
        return int(value)
    except (TypeError, ValueError):
        return default

def update_active_user(username):
    """Update the last seen time for a user in the active_users dict"""
    with lock:
        active_users[username] = datetime.now()

def get_last_21_days():
    days = []
    current = datetime.now().date()
    term_days = get_all_term_days()  # empty set if no terms configured
    limit = 0

    while len(days) < 21 and limit < 365:
        s = current.strftime("%Y-%m-%d")
        if current.weekday() < 5 and (not term_days or s in term_days):
            days.append(s)
        current -= timedelta(days=1)
        limit += 1

    return list(reversed(days))

def get_last_7_days():
    days = []
    current = datetime.now().date()
    term_days = get_all_term_days()  # empty set if no terms configured
    limit = 0

    while len(days) < 7 and limit < 365:
        s = current.strftime("%Y-%m-%d")
        if current.weekday() < 5 and (not term_days or s in term_days):
            days.append(s)
        current -= timedelta(days=1)
        limit += 1

    return list(reversed(days))

# 🔹 Database setup
def get_current_year_attendance_days():
    days = []
    current = datetime(datetime.now().year, 1, 1).date()
    today = datetime.now().date()

    while current <= today:
        s = current.strftime("%Y-%m-%d")
        if current.weekday() < 5:
            days.append(s)
        current += timedelta(days=1)

    return days

DB_NAME = "school.db"
MARKING_DB_NAME = "marking_experiment.db"
INTERACTIVE_LEARNING_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "INTERACTIVE LEARNING")

def get_db():
    return sqlite3.connect(DB_NAME)


def get_marking_db():
    return sqlite3.connect(MARKING_DB_NAME)


def init_marking_db():
    conn = get_marking_db()
    cursor = conn.cursor()
    cursor.execute("""
    CREATE TABLE IF NOT EXISTS marking_setups (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        title TEXT,
        created_by TEXT,
        created_at TEXT,
        updated_at TEXT,
        question_paper_filename TEXT,
        question_paper_blob BLOB,
        memo_filename TEXT,
        memo_blob BLOB,
        json_script_filename TEXT,
        json_script_blob BLOB,
        llm_model TEXT,
        generation_warnings TEXT,
        generation_error TEXT,
        notes TEXT
    )
    """)
    try:
        cursor.execute("PRAGMA table_info(marking_setups)")
        columns = [col[1] for col in cursor.fetchall()]
        for column_name, definition in [
            ("memo_filename", "TEXT"),
            ("memo_blob", "BLOB"),
            ("llm_model", "TEXT"),
            ("generation_warnings", "TEXT"),
            ("generation_error", "TEXT"),
        ]:
            if column_name not in columns:
                cursor.execute(f"ALTER TABLE marking_setups ADD COLUMN {column_name} {definition}")
    except Exception as exc:
        print(f"Note: marking_setups migration: {exc}")
    conn.commit()
    conn.close()


def get_days_in_active_term():
    term_range = get_active_term_range()
    if not term_range:
        return []
    start_date, end_date = term_range
    days = []
    current = datetime.strptime(start_date, "%Y-%m-%d").date()
    end = datetime.strptime(end_date, "%Y-%m-%d").date()
    while current <= end:
        if current.weekday() < 5:
            days.append(current.strftime("%Y-%m-%d"))
        current += timedelta(days=1)
    return days

def get_last_7_days():
    days = []
    current = datetime.now().date()
    term_days = get_all_term_days()  # empty set if no terms configured
    limit = 0

    while len(days) < 7 and limit < 365:
        s = current.strftime("%Y-%m-%d")
        if current.weekday() < 5 and (not term_days or s in term_days):
            days.append(s)
        current -= timedelta(days=1)
        limit += 1

    return list(reversed(days))

# 🔹 Database setup
DB_NAME = "school.db"
INTERACTIVE_LEARNING_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "INTERACTIVE LEARNING")

def get_db():
    return sqlite3.connect(DB_NAME)


def normalize_relative_path(path):
    return path.replace("/", os.sep).replace("\\", os.sep).strip()


def get_interactive_learning_files():
    files = []
    if not os.path.isdir(INTERACTIVE_LEARNING_DIR):
        return files

    for root, _, filenames in os.walk(INTERACTIVE_LEARNING_DIR):
        for filename in filenames:
            lower = filename.lower()
            if lower.startswith("~$"):
                continue
            if lower.endswith((".ppt", ".pptx", ".pps", ".ppsx")):
                full_path = os.path.join(root, filename)
                rel_path = os.path.relpath(full_path, INTERACTIVE_LEARNING_DIR)
                files.append({
                    "relative_path": rel_path.replace("\\", "/"),
                    "display_name": os.path.splitext(filename)[0],
                    "folder": os.path.dirname(rel_path).replace("\\", "/"),
                })

    files.sort(key=lambda item: (item["folder"].lower(), item["display_name"].lower()))
    return files


def resolve_interactive_learning_path(relative_path):
    if not relative_path:
        return None

    normalized = os.path.normpath(os.path.join(INTERACTIVE_LEARNING_DIR, normalize_relative_path(relative_path)))
    base_dir = os.path.normpath(INTERACTIVE_LEARNING_DIR)

    if not normalized.startswith(base_dir):
        return None
    if not os.path.isfile(normalized):
        return None
    return normalized


def extract_pptx_slides(file_path):
    lower = file_path.lower()
    if not lower.endswith((".pptx", ".ppsx")):
        return []

    namespace = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    slides = []

    with zipfile.ZipFile(file_path) as archive:
        slide_names = [
            name for name in archive.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        ]
        slide_names.sort(key=lambda item: int(re.search(r"slide(\d+)\.xml$", item).group(1)))

        for index, slide_name in enumerate(slide_names, start=1):
            root = ET.fromstring(archive.read(slide_name))
            text_runs = [
                (node.text or "").strip()
                for node in root.findall(".//a:t", namespace)
                if (node.text or "").strip()
            ]
            slides.append({
                "number": index,
                "title": text_runs[0] if text_runs else f"Slide {index}",
                "text_runs": text_runs,
                "body": text_runs[1:] if len(text_runs) > 1 else [],
            })

    return slides


def get_term_dates():
    """Return all 4 terms as a dict {1: {start, end}, ...}."""
    conn = get_db()
    cursor = conn.cursor()
    try:
        cursor.execute("SELECT term, start_date, end_date FROM term_dates ORDER BY term")
        rows = cursor.fetchall()
    except Exception:
        rows = []
    conn.close()
    terms = {i: {"start": None, "end": None} for i in range(1, 5)}
    for term, start, end in rows:
        terms[term] = {"start": start, "end": end}
    return terms


def get_active_term_range():
    """Return (start, end) of the currently active term, or None."""
    today = datetime.now().date().isoformat()
    for t in get_term_dates().values():
        if t["start"] and t["end"] and t["start"] <= today <= t["end"]:
            return t["start"], t["end"]
    return None


def get_all_term_days():
    """Return a set of all weekdays that fall within any configured term (up to today)."""
    today = datetime.now().date().isoformat()
    days = set()
    for t in get_term_dates().values():
        if t["start"] and t["end"]:
            current = datetime.strptime(t["start"], "%Y-%m-%d").date()
            end = min(datetime.strptime(t["end"], "%Y-%m-%d").date(),
                      datetime.now().date())
            while current <= end:
                if current.weekday() < 5:
                    days.add(current.strftime("%Y-%m-%d"))
                current += timedelta(days=1)
    return days


def init_db():
    conn = get_db()
    cursor = conn.cursor()

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS results (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        username TEXT,
        subject TEXT,
        task TEXT,
        score INTEGER,
        feedback TEXT,
        timestamp TEXT
    )
    """)
    cursor.execute("""
    CREATE TABLE IF NOT EXISTS weaknesses (
        username TEXT,
        skill TEXT,
        count INTEGER,
        PRIMARY KEY (username, skill)
    )
    """)

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS users (
        username TEXT PRIMARY KEY,
        role TEXT,
        last_active TEXT,
        full_name TEXT,
        group_name TEXT,
        teacher_username TEXT
    )
    """)

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS group_teachers (
        group_name TEXT PRIMARY KEY,
        teacher_username TEXT
    )
    """)

    # Migration: add direct teacher assignment to users if missing
    try:
        cursor.execute("PRAGMA table_info(users)")
        user_cols = [c[1] for c in cursor.fetchall()]
        if "teacher_username" not in user_cols:
            cursor.execute("ALTER TABLE users ADD COLUMN teacher_username TEXT")
            print("Migration: added teacher_username column to users")
    except Exception as e:
        print(f"Note: users migration check: {e}")

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS login_history (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        username TEXT,
        login_time TEXT,
        date TEXT
    )
    """)

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS attendance_override (
        username TEXT,
        date TEXT,
        status TEXT,  -- 'present', 'late' or 'absent'
        PRIMARY KEY (username, date)
    )
    """)

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS communication_threads (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        student_username TEXT NOT NULL,
        teacher_username TEXT,
        group_name TEXT,
        topic TEXT NOT NULL,
        subject_line TEXT,
        attendance_date TEXT,
        status TEXT DEFAULT 'open',
        chat_session_id TEXT,
        created_at TEXT,
        updated_at TEXT,
        teacher_read_at TEXT,
        student_read_at TEXT
    )
    """)

    try:
        cursor.execute("PRAGMA table_info(communication_threads)")
        thread_cols = [c[1] for c in cursor.fetchall()]
        if "teacher_read_at" not in thread_cols:
            cursor.execute("ALTER TABLE communication_threads ADD COLUMN teacher_read_at TEXT")
        if "student_read_at" not in thread_cols:
            cursor.execute("ALTER TABLE communication_threads ADD COLUMN student_read_at TEXT")
    except Exception as e:
        print(f"Note: communication_threads migration check: {e}")

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS communication_messages (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        thread_id INTEGER NOT NULL,
        sender_username TEXT NOT NULL,
        sender_role TEXT,
        message TEXT NOT NULL,
        created_at TEXT,
        FOREIGN KEY(thread_id) REFERENCES communication_threads(id)
    )
    """)

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS teacher_quick_actions (
        username TEXT NOT NULL,
        action_key TEXT NOT NULL,
        PRIMARY KEY (username, action_key)
    )
    """)

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS excluded_dates (
        date TEXT,
        group_name TEXT,  -- NULL for global exclusions, specific group name for group-specific
        reason TEXT,
        created_by TEXT,
        created_at TEXT,
        PRIMARY KEY (date, group_name)
    )
    """)

    # Migration: Handle old table structure
    try:
        cursor.execute("PRAGMA table_info(excluded_dates)")
        columns = cursor.fetchall()
        column_names = [col[1] for col in columns]
        
        if 'group_name' not in column_names:
            # Old table structure exists, migrate it
            print("Migrating excluded_dates table...")
            
            # Create new table with proper structure
            cursor.execute("""
            CREATE TABLE excluded_dates_new (
                date TEXT,
                group_name TEXT,
                reason TEXT,
                created_by TEXT,
                created_at TEXT,
                PRIMARY KEY (date, group_name)
            )
            """)
            
            # Copy old data (assuming old structure was: date, reason, created_by, created_at)
            cursor.execute("""
            INSERT INTO excluded_dates_new (date, group_name, reason, created_by, created_at)
            SELECT date, NULL, reason, created_by, created_at FROM excluded_dates
            """)
            
            # Replace old table
            cursor.execute("DROP TABLE excluded_dates")
            cursor.execute("ALTER TABLE excluded_dates_new RENAME TO excluded_dates")
            print("Migration completed successfully")
            
    except Exception as e:
        print(f"Note: Migration check completed with: {e}")
        # Continue anyway - table should exist now

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS activities (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        username TEXT,
        action TEXT,
        timestamp TEXT
    )
    """)

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS subjects (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        name TEXT UNIQUE,
        created_by TEXT,
        created_at TEXT
    )
    """)

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS theory_tests (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        title TEXT,
        subject TEXT,
        background_image TEXT,
        background_fit TEXT DEFAULT 'cover',
        assign_date TEXT,  -- YYYY-MM-DD
        time_limit INTEGER,
        allow_multiple INTEGER DEFAULT 0,
        max_attempts INTEGER DEFAULT 1,
        show_answers INTEGER DEFAULT 1,
        created_by TEXT,
        created_at TEXT,
        is_active INTEGER DEFAULT 0
    )
    """)

    # Migration: add assign_date if missing (idempotent)
    try:
        cursor.execute("PRAGMA table_info(theory_tests)")
        cols = [c[1] for c in cursor.fetchall()]
        if "assign_date" not in cols:
            cursor.execute("ALTER TABLE theory_tests ADD COLUMN assign_date TEXT")
            print("Migration: added assign_date column to theory_tests")
    except Exception as e:
        print(f"Note: theory_tests assign_date migration: {e}")

    # Migration: remove group_name column from theory_tests if present,
    # and add allow_multiple, max_attempts, show_answers if missing
    try:
        cursor.execute("PRAGMA table_info(theory_tests)")
        cols = [c[1] for c in cursor.fetchall()]
        if "group_name" in cols:
            cursor.execute("""
                CREATE TABLE theory_tests_new (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    title TEXT, subject TEXT, time_limit INTEGER,
                    allow_multiple INTEGER DEFAULT 0,
                    max_attempts INTEGER DEFAULT 1,
                    show_answers INTEGER DEFAULT 1,
                    created_by TEXT, created_at TEXT, is_active INTEGER DEFAULT 0
                )
            """)
            cursor.execute("INSERT INTO theory_tests_new (id,title,subject,time_limit,created_by,created_at,is_active) SELECT id,title,subject,time_limit,created_by,created_at,is_active FROM theory_tests")
            cursor.execute("DROP TABLE theory_tests")
            cursor.execute("ALTER TABLE theory_tests_new RENAME TO theory_tests")
            print("Migration: removed group_name from theory_tests")
        else:
            if "allow_multiple" not in cols:
                cursor.execute("ALTER TABLE theory_tests ADD COLUMN allow_multiple INTEGER DEFAULT 0")
            if "max_attempts" not in cols:
                cursor.execute("ALTER TABLE theory_tests ADD COLUMN max_attempts INTEGER DEFAULT 1")
            if "show_answers" not in cols:
                cursor.execute("ALTER TABLE theory_tests ADD COLUMN show_answers INTEGER DEFAULT 1")
            if "background_image" not in cols:
                cursor.execute("ALTER TABLE theory_tests ADD COLUMN background_image TEXT")
            if "background_fit" not in cols:
                cursor.execute("ALTER TABLE theory_tests ADD COLUMN background_fit TEXT DEFAULT 'cover'")
    except Exception as e:
        print(f"Note: theory_tests migration: {e}")

    try:
        cursor.execute("PRAGMA table_info(theory_tests)")
        cols = [c[1] for c in cursor.fetchall()]
        if "background_image" not in cols:
            cursor.execute("ALTER TABLE theory_tests ADD COLUMN background_image TEXT")
        if "background_fit" not in cols:
            cursor.execute("ALTER TABLE theory_tests ADD COLUMN background_fit TEXT DEFAULT 'cover'")
    except Exception as e:
        print(f"Note: theory_tests background migration: {e}")

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS theory_test_groups (
        test_id INTEGER,
        group_name TEXT,
        PRIMARY KEY (test_id, group_name),
        FOREIGN KEY (test_id) REFERENCES theory_tests (id)
    )
    """)

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS theory_questions (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        test_id INTEGER,
        question_text TEXT,
        question_type TEXT,
        marks INTEGER DEFAULT 1,
        order_index INTEGER DEFAULT 0,
        FOREIGN KEY (test_id) REFERENCES theory_tests (id)
    )
    """)

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS theory_options (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        question_id INTEGER,
        option_text TEXT,
        is_correct INTEGER DEFAULT 0,
        match_pair TEXT,
        FOREIGN KEY (question_id) REFERENCES theory_questions (id)
    )
    """)

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS theory_submissions (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        test_id INTEGER,
        username TEXT,
        score INTEGER,
        total INTEGER,
        percentage INTEGER,
        submitted_at TEXT,
        FOREIGN KEY (test_id) REFERENCES theory_tests (id)
    )
    """)

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS theory_answers (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        submission_id INTEGER,
        question_id INTEGER,
        answer_text TEXT,
        is_correct INTEGER DEFAULT 0,
        marks_awarded INTEGER DEFAULT 0,
        FOREIGN KEY (submission_id) REFERENCES theory_submissions (id),
        FOREIGN KEY (question_id) REFERENCES theory_questions (id)
    )
    """)

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS question_bank_questions (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        question_text TEXT,
        question_type TEXT,
        marks INTEGER DEFAULT 1,
        subject TEXT,
        modules TEXT,
        created_by TEXT,
        created_at TEXT
    )
    """)

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS question_bank_options (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        bank_question_id INTEGER,
        option_text TEXT,
        is_correct INTEGER DEFAULT 0,
        match_pair TEXT,
        FOREIGN KEY (bank_question_id) REFERENCES question_bank_questions (id)
    )
    """)

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS theory_progress (
        test_id INTEGER,
        username TEXT,
        current_slide INTEGER DEFAULT 0,
        max_slide INTEGER DEFAULT 0,
        time_spent_seconds INTEGER DEFAULT 0,
        completed INTEGER DEFAULT 0,
        updated_at TEXT,
        PRIMARY KEY (test_id, username),
        FOREIGN KEY (test_id) REFERENCES theory_tests (id)
    )
    """)

    try:
        cursor.execute("PRAGMA table_info(theory_progress)")
        progress_cols = [c[1] for c in cursor.fetchall()]
        if "max_slide" not in progress_cols:
            cursor.execute("ALTER TABLE theory_progress ADD COLUMN max_slide INTEGER DEFAULT 0")
        if "time_spent_seconds" not in progress_cols:
            cursor.execute("ALTER TABLE theory_progress ADD COLUMN time_spent_seconds INTEGER DEFAULT 0")
    except Exception as e:
        print(f"Note: theory_progress migration: {e}")

    try:
        cursor.execute("PRAGMA table_info(theory_submissions)")
        sub_cols = [c[1] for c in cursor.fetchall()]
        if "time_spent_seconds" not in sub_cols:
            cursor.execute("ALTER TABLE theory_submissions ADD COLUMN time_spent_seconds INTEGER DEFAULT 0")
        if "submission_type" not in sub_cols:
            cursor.execute("ALTER TABLE theory_submissions ADD COLUMN submission_type TEXT DEFAULT 'test'")
    except Exception as e:
        print(f"Note: theory_submissions time/type migration: {e}")



    cursor.execute("""
    CREATE TABLE IF NOT EXISTS theory_lesson_groups (
        lesson_id INTEGER,
        group_name TEXT,
        PRIMARY KEY (lesson_id, group_name),
        FOREIGN KEY (lesson_id) REFERENCES theory_lessons (id)
    )
    """)

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS theory_lesson_teachers (
        lesson_id INTEGER,
        teacher_username TEXT,
        PRIMARY KEY (lesson_id, teacher_username),
        FOREIGN KEY (lesson_id) REFERENCES theory_lessons (id)
    )
    """)

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS theory_lesson_checkpoints (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        lesson_id INTEGER,
        slide_number INTEGER,
        question_text TEXT,
        option_a TEXT,
        option_b TEXT,
        option_c TEXT,
        option_d TEXT,
        correct_option TEXT,
        explanation TEXT,
        sort_order INTEGER DEFAULT 0,
        FOREIGN KEY (lesson_id) REFERENCES theory_lessons (id)
    )
    """)

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS theory_lesson_progress (
        lesson_id INTEGER,
        username TEXT,
        current_slide INTEGER DEFAULT 1,
        completed_at TEXT,
        PRIMARY KEY (lesson_id, username),
        FOREIGN KEY (lesson_id) REFERENCES theory_lessons (id)
    )
    """)

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS theory_lesson_checkpoint_answers (
        lesson_id INTEGER,
        checkpoint_id INTEGER,
        username TEXT,
        selected_option TEXT,
        is_correct INTEGER DEFAULT 0,
        answered_at TEXT,
        PRIMARY KEY (checkpoint_id, username),
        FOREIGN KEY (lesson_id) REFERENCES theory_lessons (id),
        FOREIGN KEY (checkpoint_id) REFERENCES theory_lesson_checkpoints (id)
    )
    """)

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS tasks (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        subject_id INTEGER,
        name TEXT,
        assign_date TEXT,
        marking_script TEXT,
        theory_test_id INTEGER,
        task_type TEXT DEFAULT 'practical',
        allow_multiple INTEGER DEFAULT 0,
        max_attempts INTEGER DEFAULT 1,
        is_active INTEGER DEFAULT 1,
        marking_setup_id INTEGER,
        created_by TEXT,
        created_at TEXT,
        FOREIGN KEY (subject_id) REFERENCES subjects (id)
    )
    """)

    # Migration: add columns to tasks if missing
    try:
        cursor.execute("PRAGMA table_info(tasks)")
        columns = [col[1] for col in cursor.fetchall()]
        if "marking_script" not in columns:
            cursor.execute("ALTER TABLE tasks ADD COLUMN marking_script TEXT")
            print("Migration: added marking_script column to tasks")
        if "theory_test_id" not in columns:
            cursor.execute("ALTER TABLE tasks ADD COLUMN theory_test_id INTEGER")
            print("Migration: added theory_test_id column to tasks")
        if "task_type" not in columns:
            cursor.execute("ALTER TABLE tasks ADD COLUMN task_type TEXT DEFAULT 'practical'")
            print("Migration: added task_type column to tasks")
        if "allow_multiple" not in columns:
            cursor.execute("ALTER TABLE tasks ADD COLUMN allow_multiple INTEGER DEFAULT 0")
            print("Migration: added allow_multiple column to tasks")
        if "max_attempts" not in columns:
            cursor.execute("ALTER TABLE tasks ADD COLUMN max_attempts INTEGER DEFAULT 1")
            print("Migration: added max_attempts column to tasks")
        if "is_active" not in columns:
            cursor.execute("ALTER TABLE tasks ADD COLUMN is_active INTEGER DEFAULT 1")
        if "marking_setup_id" not in columns:
            cursor.execute("ALTER TABLE tasks ADD COLUMN marking_setup_id INTEGER")
            print("Migration: added marking_setup_id column to tasks")
        
        if "question_text" not in columns:
            cursor.execute("ALTER TABLE tasks ADD COLUMN question_text TEXT")
            print("Migration: added question_text column to tasks")

        if "sample_file" not in columns:
            cursor.execute("ALTER TABLE tasks ADD COLUMN sample_file BLOB")
            print("Migration: added sample_file BLOB column to tasks")

        if "sample_file_name" not in columns:
            cursor.execute("ALTER TABLE tasks ADD COLUMN sample_file_name TEXT")
            print("Migration: added sample_file_name column to tasks")
    except Exception as e:
        print(f"Note: tasks migration check: {e}")

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS learner_notes (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        username TEXT,
        note TEXT,
        flag TEXT,
        created_by TEXT,
        created_at TEXT
    )
    """)

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS result_removals (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        username TEXT,
        task_type TEXT,
        subject TEXT,
        task_name TEXT,
        test_id INTEGER,
        removed_by TEXT,
        reason TEXT,
        removed_at TEXT
    )
    """)

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS task_groups (
        task_id INTEGER,
        group_name TEXT,
        PRIMARY KEY (task_id, group_name),
        FOREIGN KEY (task_id) REFERENCES tasks (id)
    )
    """)

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS task_teachers (
        task_id INTEGER,
        teacher_username TEXT,
        PRIMARY KEY (task_id, teacher_username),
        FOREIGN KEY (task_id) REFERENCES tasks (id)
    )
    """)

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS theory_test_teachers (
        test_id INTEGER,
        teacher_username TEXT,
        PRIMARY KEY (test_id, teacher_username),
        FOREIGN KEY (test_id) REFERENCES theory_tests (id)
    )
    """)

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS term_dates (
        term INTEGER PRIMARY KEY,  -- 1, 2, 3, or 4
        start_date TEXT,
        end_date TEXT
    )
    """)

    cursor.execute("CREATE INDEX IF NOT EXISTS idx_results_user_subject_task ON results (username, subject, task)")
    cursor.execute("CREATE INDEX IF NOT EXISTS idx_results_user_timestamp ON results (username, timestamp)")
    cursor.execute("CREATE INDEX IF NOT EXISTS idx_login_history_user_date ON login_history (username, date)")
    cursor.execute("CREATE INDEX IF NOT EXISTS idx_login_history_date_user ON login_history (date, username)")
    cursor.execute("CREATE INDEX IF NOT EXISTS idx_attendance_override_user_date ON attendance_override (username, date)")
    cursor.execute("CREATE INDEX IF NOT EXISTS idx_users_role_teacher_group ON users (role, teacher_username, group_name)")
    cursor.execute("CREATE INDEX IF NOT EXISTS idx_users_group_role ON users (group_name, role)")
    cursor.execute("CREATE INDEX IF NOT EXISTS idx_theory_submissions_user_test ON theory_submissions (username, test_id)")
    cursor.execute("CREATE INDEX IF NOT EXISTS idx_theory_submissions_test_user ON theory_submissions (test_id, username)")
    cursor.execute("CREATE INDEX IF NOT EXISTS idx_task_groups_group_task ON task_groups (group_name, task_id)")
    cursor.execute("CREATE INDEX IF NOT EXISTS idx_theory_test_groups_group_test ON theory_test_groups (group_name, test_id)")
    cursor.execute("CREATE INDEX IF NOT EXISTS idx_tasks_assign_type ON tasks (assign_date, task_type)")
    cursor.execute("CREATE INDEX IF NOT EXISTS idx_theory_tests_assign_date ON theory_tests (assign_date)")
    cursor.execute("CREATE INDEX IF NOT EXISTS idx_excluded_dates_group_date ON excluded_dates (group_name, date)")
    cursor.execute("CREATE INDEX IF NOT EXISTS idx_communication_threads_teacher_updated ON communication_threads (teacher_username, updated_at)")
    cursor.execute("CREATE INDEX IF NOT EXISTS idx_communication_messages_thread_created ON communication_messages (thread_id, created_at)")

    # Seed initial subjects if empty
    cursor.execute("SELECT COUNT(*) FROM subjects")
    if cursor.fetchone()[0] == 0:
        initial_subjects = ["Word", "Excel", "Access", "HTML"]
        for subj in initial_subjects:
            cursor.execute("INSERT INTO subjects (name, created_by, created_at) VALUES (?, ?, ?)",
                           (subj, "system", datetime.now().isoformat()))

    # Seed initial tasks if empty
    cursor.execute("SELECT COUNT(*) FROM tasks")
    if cursor.fetchone()[0] == 0:
        cursor.execute("SELECT id, name FROM subjects")
        subjects = cursor.fetchall()
        today = datetime.now().date().isoformat()
        for subj_id, subj_name in subjects:
            for i in range(1, 4):
                task_name = f"Task {i}"
                cursor.execute("INSERT INTO tasks (subject_id, name, assign_date, created_by, created_at) VALUES (?, ?, ?, ?, ?)",
                               (subj_id, task_name, today, "system", datetime.now().isoformat()))
                task_id = cursor.lastrowid
                # Assign to all groups
                cursor.execute("SELECT DISTINCT group_name FROM users WHERE group_name IS NOT NULL")
                groups = cursor.fetchall()
                for group_row in groups:
                    cursor.execute("INSERT INTO task_groups (task_id, group_name) VALUES (?, ?)", (task_id, group_row[0]))

    conn.commit()
    conn.close()

def get_user_role(username):
    conn = get_db()
    cursor = conn.cursor()

    cursor.execute("SELECT role FROM users WHERE username = ?", (username,))
    result = cursor.fetchone()

    conn.close()

    return result[0] if result else "student"  # default = student

def log_login(username):
    conn = get_db()
    cursor = conn.cursor()

    now = datetime.now()

    cursor.execute("""
    INSERT INTO login_history (username, login_time, date)
    VALUES (?, ?, ?)
    """, (username, str(now), now.strftime("%Y-%m-%d")))

    conn.commit()
    conn.close()

def should_record_attendance_login():
    return True

def log_activity(username, action):
    conn = get_db()
    cursor = conn.cursor()

    now = datetime.now().isoformat()

    cursor.execute("""
    INSERT INTO activities (username, action, timestamp)
    VALUES (?, ?, ?)
    """, (username, action, now))

    conn.commit()
    conn.close()

def create_user_if_not_exists(username):
    conn = get_db()
    cursor = conn.cursor()

    cursor.execute("""
    INSERT INTO users (username, role, last_active)
    VALUES (?, 'student', ?)
    ON CONFLICT(username) DO NOTHING
    """, (username, str(datetime.now())))

    conn.commit()
    conn.close()

def update_weakness(username, skills):
    conn = get_db()
    cursor = conn.cursor()

    for skill in skills:
        cursor.execute("""
        INSERT INTO weaknesses (username, skill, count)
        VALUES (?, ?, 1)
        ON CONFLICT(username, skill)
        DO UPDATE SET count = count + 1
        """, (username, skill))

    conn.commit()
    conn.close()

def get_weaknesses(username):
    conn = get_db()
    cursor = conn.cursor()

    cursor.execute("""
    SELECT skill, count
    FROM weaknesses
    WHERE username = ?
    ORDER BY count DESC
    LIMIT 5
    """, (username,))

    results = cursor.fetchall()
    conn.close()

    return results

def save_result(username, subject, task, score, feedback):
    conn = get_db()
    cursor = conn.cursor()

    cursor.execute("""
    INSERT INTO results (username, subject, task, score, feedback, timestamp)
    VALUES (?, ?, ?, ?, ?, ?)
    """, (username, subject, task, score, feedback, str(datetime.now())))

    conn.commit()
    conn.close()

def mark_file(filepath, marking_script, marking_setup_id=None):
    """
    Route a submitted file to the correct task marker using the
    marking_script name stored on the task in the database.
    """
    import importlib
    if not marking_script:
        return {
            "task_name": "Unknown Task",
            "score": 0,
            "total": 0,
            "percentage": 0,
            "results": [],
            "error": "No marking script assigned to this task. Please contact your teacher."
        }
    try:
        module = importlib.import_module(f"marking.tasks.{marking_script}")
        if marking_setup_id is not None and hasattr(module, "mark_with_setup"):
            return module.mark_with_setup(filepath, int(marking_setup_id))
        return module.mark(filepath)
    except ModuleNotFoundError:
        return {
            "task_name": marking_script,
            "score": 0,
            "total": 0,
            "percentage": 0,
            "results": [],
            "error": f"Marking script '{marking_script}' not found. Please contact your teacher."
        }

app = Flask(__name__)
app.secret_key = os.getenv('SECRET_KEY', 'dev-key-change-in-prod')

@app.after_request
def inject_global_mobile_css(response):
    if response.content_type and response.content_type.startswith("text/html"):
        html = response.get_data(as_text=True)
        changed = False
        viewport_meta = '<meta name="viewport" content="width=device-width, initial-scale=1.0">'
        mobile_css = '<link rel="stylesheet" href="/static/css/mobile_global.css">'
        mobile_nav_script = '''<script>
(function () {
    function closeMenu() {
        document.body.classList.remove('mobile-nav-open');
        document.querySelectorAll('.mobile-menu-toggle').forEach(function (button) {
            button.setAttribute('aria-expanded', 'false');
        });
    }

    document.addEventListener('DOMContentLoaded', function () {
        document.querySelectorAll('.topbar').forEach(function (topbar) {
            var links = topbar.querySelector('.topbar-links');
            if (!links) {
                var title = topbar.querySelector('h1, strong');
                var movable = Array.prototype.filter.call(topbar.children, function (child) {
                    return child !== title && !child.classList.contains('mobile-menu-toggle');
                });
                if (!movable.length) return;
                links = document.createElement('div');
                links.className = 'topbar-links';
                movable.forEach(function (child) {
                    links.appendChild(child);
                });
                topbar.appendChild(links);
            }
            if (!links || topbar.querySelector('.mobile-menu-toggle')) return;

            var button = document.createElement('button');
            button.type = 'button';
            button.className = 'mobile-menu-toggle';
            button.setAttribute('aria-label', 'Open menu');
            button.setAttribute('aria-expanded', 'false');
            button.innerHTML = '&#9776;';
            button.addEventListener('click', function (event) {
                event.stopPropagation();
                var isOpen = document.body.classList.toggle('mobile-nav-open');
                button.setAttribute('aria-expanded', isOpen ? 'true' : 'false');
            });
            topbar.insertBefore(button, topbar.firstChild);
        });

        document.addEventListener('click', function (event) {
            if (!document.body.classList.contains('mobile-nav-open')) return;
            if (event.target.closest('.topbar-links') || event.target.closest('.mobile-menu-toggle')) return;
            closeMenu();
        });

        document.addEventListener('keydown', function (event) {
            if (event.key === 'Escape') closeMenu();
        });

        document.querySelectorAll('.topbar-links a').forEach(function (link) {
            link.addEventListener('click', closeMenu);
        });
    });
})();
</script>'''
        if 'name="viewport"' not in html and "</head>" in html:
            html = html.replace("</head>", f"    {viewport_meta}\n</head>", 1)
            changed = True
        if mobile_css not in html and "</head>" in html:
            html = html.replace("</head>", f"    {mobile_css}\n</head>", 1)
            changed = True
        if "mobile-menu-toggle" not in html and "</body>" in html:
            html = html.replace("</body>", f"    {mobile_nav_script}\n</body>", 1)
            changed = True
        if changed:
            response.set_data(html)
            response.headers["Content-Length"] = str(len(response.get_data()))
    return response

@app.context_processor
def inject_session_user():
    from flask import session as _s
    uname = _s.get('username', '')
    role = ''
    teacher_unread_messages = 0
    student_unread_messages = 0
    if uname:
        try:
            role = get_user_role(uname)
            if role in ["teacher", "admin"]:
                teacher_unread_messages = get_teacher_unread_message_count(uname)
            elif role == "student":
                student_unread_messages = get_student_unread_message_count(uname)
        except Exception:
            pass
    return dict(
        session_username=uname,
        session_role=role,
        teacher_unread_messages=teacher_unread_messages,
        student_unread_messages=student_unread_messages,
    )


# Store active users in memory
active_users = {}
lock = threading.Lock()

def get_recent_results(username):
    conn = get_db()
    cursor = conn.cursor()

    cursor.execute("""
    SELECT subject, task, score
    FROM results
    WHERE username = ?
    ORDER BY id DESC
    LIMIT 5
    """, (username,))

    results = cursor.fetchall()
    conn.close()

    return results

def attendance_status_counts_as_present(status):
    return (status or "").strip().lower() in {"present", "late"}

def fetch_first_login_times(cursor, usernames, days):
    login_map = {}
    if not usernames or not days:
        return login_map
    user_placeholders = ",".join("?" for _ in usernames)
    day_placeholders = ",".join("?" for _ in days)
    cursor.execute(
        f"""
        SELECT username, date, MIN(login_time)
        FROM login_history
        WHERE username IN ({user_placeholders}) AND date IN ({day_placeholders})
        GROUP BY username, date
        """,
        list(usernames) + list(days),
    )
    for username, day, login_time in cursor.fetchall():
        login_map[(username, day)] = login_time
    return login_map

def fetch_attendance_override_statuses(cursor, usernames, days):
    override_map = {}
    if not usernames or not days:
        return override_map
    user_placeholders = ",".join("?" for _ in usernames)
    day_placeholders = ",".join("?" for _ in days)
    cursor.execute(
        f"""
        SELECT username, date, status
        FROM attendance_override
        WHERE username IN ({user_placeholders}) AND date IN ({day_placeholders})
        """,
        list(usernames) + list(days),
    )
    for username, day, status in cursor.fetchall():
        override_map[(username, day)] = status
    return override_map

def fetch_group_late_thresholds(cursor, group, days):
    late_cutoffs = {day: None for day in days}
    if not group or not days:
        return late_cutoffs
    day_placeholders = ",".join("?" for _ in days)
    cursor.execute(
        f"""
        SELECT lh.date, MIN(lh.login_time)
        FROM login_history lh
        JOIN users u ON u.username = lh.username
        WHERE u.group_name = ? AND lh.date IN ({day_placeholders})
        GROUP BY lh.date
        """,
        [group, *days],
    )
    for day, min_login in cursor.fetchall():
        if not min_login:
            continue
        try:
            min_dt = datetime.fromisoformat(min_login)
        except ValueError:
            min_dt = datetime.strptime(min_login, "%Y-%m-%d %H:%M:%S")
        late_cutoffs[day] = (min_dt + timedelta(minutes=4)).strftime("%H:%M")
    return late_cutoffs

def fetch_group_excluded_dates(cursor, groups):
    excluded_by_group = {group: set() for group in groups if group}
    if not excluded_by_group:
        return excluded_by_group
    placeholders = ",".join("?" for _ in excluded_by_group)
    cursor.execute(
        f"""
        SELECT date, group_name
        FROM excluded_dates
        WHERE group_name IS NULL OR group_name IN ({placeholders})
        """,
        list(excluded_by_group.keys()),
    )
    rows = cursor.fetchall()
    global_dates = {date_str for date_str, group_name in rows if group_name is None}
    for group in excluded_by_group:
        excluded_by_group[group] = set(global_dates)
    for date_str, group_name in rows:
        if group_name in excluded_by_group:
            excluded_by_group[group_name].add(date_str)
    return excluded_by_group

def fetch_present_day_map(cursor, usernames, days):
    present_map = {username: set() for username in usernames}
    if not present_map or not days:
        return present_map
    user_placeholders = ",".join("?" for _ in present_map)
    day_placeholders = ",".join("?" for _ in days)
    cursor.execute(
        f"""
        SELECT DISTINCT username, date
        FROM login_history
        WHERE username IN ({user_placeholders}) AND date IN ({day_placeholders})
        """,
        list(present_map.keys()) + list(days),
    )
    for username, day in cursor.fetchall():
        present_map.setdefault(username, set()).add(day)
    cursor.execute(
        f"""
        SELECT username, date, status
        FROM attendance_override
        WHERE username IN ({user_placeholders}) AND date IN ({day_placeholders})
        """,
        list(present_map.keys()) + list(days),
    )
    for username, day, status in cursor.fetchall():
        if attendance_status_counts_as_present(status):
            present_map.setdefault(username, set()).add(day)
        else:
            present_map.setdefault(username, set()).discard(day)
    return present_map

def build_attendance_history(cursor, username, group_name, days, login_map=None, override_map=None, late_cutoffs=None, class_checked_dates=None, excluded_dates=None):
    history = []
    login_map = login_map or {}
    override_map = override_map or {}
    late_cutoffs = late_cutoffs or {}
    class_checked_dates = class_checked_dates or set()
    excluded_dates = excluded_dates or set()

    for day in days:
        weekday = datetime.strptime(day, "%Y-%m-%d").weekday()
        login_time = login_map.get((username, day))
        if day in excluded_dates:
            history.append({"date": day, "status": "Normal", "time": "", "late": False, "note": "", "weekday": weekday})
            continue

        if login_time:
            time_str = login_time.split(" ")[1][:5]
            cutoff = late_cutoffs.get(day)
            late = cutoff is not None and time_str > cutoff
            history.append({"date": day, "status": "Present", "time": time_str, "late": late, "note": "Auto", "weekday": weekday})
            continue

        override_status = (override_map.get((username, day)) or "").strip().lower()
        if attendance_status_counts_as_present(override_status):
            manual_time = "12:00" if override_status == "present" else ""
            history.append({"date": day, "status": "Present", "time": manual_time, "late": override_status == "late", "note": "Manual", "weekday": weekday})
        elif override_status == "absent" or day in class_checked_dates:
            note = "Manual" if override_status == "absent" else "Class checked in"
            history.append({"date": day, "status": "Absent", "time": "", "late": False, "note": note, "weekday": weekday})
        else:
            history.append({"date": day, "status": "Normal", "time": "", "late": False, "note": "", "weekday": weekday})

    return history

def summarize_attendance_history(history):
    counted_history = [item for item in history if item["status"] in ("Present", "Absent")]
    total_days = len(counted_history)
    present_days = sum(1 for item in history if item["status"] == "Present")
    absent_days = sum(1 for item in history if item["status"] == "Absent")
    late_days = sum(1 for item in history if item["late"])
    attendance_pct = round((present_days / total_days) * 100) if total_days else 0
    return {
        "total_days": total_days,
        "present_days": present_days,
        "absent_days": absent_days,
        "late_days": late_days,
        "attendance_pct": attendance_pct,
    }

def build_attendance_group_summary(cursor, groups, teacher_username=None, days=None):
    summaries = []
    days = days or get_last_21_days()
    for group_name in groups:
        if not group_name:
            continue
        if teacher_username:
            cursor.execute("""
                SELECT username, full_name
                FROM users
                WHERE group_name = ? AND role = 'student' AND teacher_username = ?
                ORDER BY full_name
            """, (group_name, teacher_username))
        else:
            cursor.execute("""
                SELECT username, full_name
                FROM users
                WHERE group_name = ? AND role = 'student'
                ORDER BY full_name
            """, (group_name,))
        learners = cursor.fetchall()
        usernames = [row[0] for row in learners]
        auto_exclude_empty_attendance_days(cursor, [group_name], created_by=teacher_username or "system", days=get_current_year_attendance_days())
        excluded_dates = fetch_group_excluded_dates(cursor, [group_name]).get(group_name, set())
        filtered_days = [day for day in days if day not in excluded_dates]
        login_map = fetch_first_login_times(cursor, usernames, filtered_days)
        override_map = fetch_attendance_override_statuses(cursor, usernames, filtered_days)
        late_cutoffs = fetch_group_late_thresholds(cursor, group_name, filtered_days)
        cursor.execute("""
            SELECT DISTINCT lh.date
            FROM login_history lh
            JOIN users u ON u.username = lh.username
            WHERE u.group_name = ? AND u.role = 'student'
        """, (group_name,))
        class_checked_dates = {row[0] for row in cursor.fetchall()}

        present_total = 0
        absent_total = 0
        late_total = 0
        counted_total = 0
        for uname in usernames:
            history = build_attendance_history(
                cursor,
                uname,
                group_name,
                filtered_days,
                login_map=login_map,
                override_map=override_map,
                late_cutoffs=late_cutoffs,
                class_checked_dates=class_checked_dates,
                excluded_dates=excluded_dates,
            )
            summary = summarize_attendance_history(history)
            present_total += summary["present_days"]
            absent_total += summary["absent_days"]
            late_total += summary["late_days"]
            counted_total += summary["total_days"]

        att_pct = round((present_total / counted_total) * 100) if counted_total else 0
        summaries.append({
            "group": group_name,
            "students": len(learners),
            "attendance_pct": att_pct,
            "present_days": present_total,
            "absent_days": absent_total,
            "late_days": late_total,
        })
    return summaries

def build_attendance_months(history):
    attendance_months = []
    for item in history:
        month_key = item["date"][:7]
        if not attendance_months or attendance_months[-1]["key"] != month_key:
            month_date = datetime.strptime(item["date"], "%Y-%m-%d")
            attendance_months.append({
                "key": month_key,
                "name": month_date.strftime("%B"),
                "days": []
            })
        attendance_months[-1]["days"].append(item)
    return attendance_months

def auto_exclude_empty_attendance_days(cursor, groups, created_by="system", days=None):
    groups = [group for group in groups if group]
    if not groups:
        return 0

    if days is None:
        days = get_current_year_attendance_days()
    today = datetime.now().date()
    filtered_days = []
    for day in days:
        try:
            day_obj = datetime.strptime(day, "%Y-%m-%d").date()
        except ValueError:
            continue
        if day_obj >= today:
            continue
        if day_obj.weekday() >= 5:
            continue
        filtered_days.append(day)

    days = filtered_days
    if not days:
        return 0

    group_placeholders = ",".join("?" for _ in groups)
    day_placeholders = ",".join("?" for _ in days)

    cursor.execute(
        f"""
        SELECT u.group_name, lh.date, COUNT(DISTINCT lh.username)
        FROM login_history lh
        JOIN users u ON u.username = lh.username
        WHERE u.group_name IN ({group_placeholders})
          AND u.role = 'student'
          AND lh.date IN ({day_placeholders})
        GROUP BY u.group_name, lh.date
        """,
        [*groups, *days],
    )
    login_counts = {(group_name, day): count for group_name, day, count in cursor.fetchall()}

    cursor.execute(
        f"""
        SELECT u.group_name, ao.date, COUNT(DISTINCT ao.username)
        FROM attendance_override ao
        JOIN users u ON u.username = ao.username
        WHERE u.group_name IN ({group_placeholders})
          AND u.role = 'student'
          AND ao.date IN ({day_placeholders})
          AND LOWER(COALESCE(ao.status, '')) IN ('present', 'late')
        GROUP BY u.group_name, ao.date
        """,
        [*groups, *days],
    )
    override_counts = {(group_name, day): count for group_name, day, count in cursor.fetchall()}

    cursor.execute(
        f"""
        SELECT date, group_name
        FROM excluded_dates
        WHERE (group_name IN ({group_placeholders}) OR group_name IS NULL)
          AND date IN ({day_placeholders})
        """,
        [*groups, *days],
    )
    existing = {(date_str, group_name) for date_str, group_name in cursor.fetchall()}

    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    inserted = 0
    for group_name in groups:
        for day in days:
            if (day, group_name) in existing:
                continue
            if (day, None) in existing:
                continue
            if login_counts.get((group_name, day), 0) > 0:
                continue
            if override_counts.get((group_name, day), 0) > 0:
                continue
            cursor.execute(
                """
                INSERT OR IGNORE INTO excluded_dates (date, group_name, reason, created_by, created_at)
                VALUES (?, ?, ?, ?, ?)
                """,
                (day, group_name, "Auto excluded: no class attendance recorded", created_by, timestamp),
            )
            if cursor.rowcount:
                inserted += 1
    return inserted

def fetch_student_practical_averages(cursor, teacher_username=None):
    params = []
    teacher_filter = ""
    if teacher_username:
        teacher_filter = " AND u.teacher_username = ?"
        params.append(teacher_username)
    cursor.execute(
        f"""
        SELECT best_scores.username, ROUND(AVG(best_scores.best_score), 1)
        FROM (
            SELECT username, subject, task, MAX(score) AS best_score
            FROM results
            GROUP BY username, subject, task
        ) best_scores
        JOIN users u ON u.username = best_scores.username
        WHERE u.role = 'student'{teacher_filter}
        GROUP BY best_scores.username
        """,
        params,
    )
    return {username: average for username, average in cursor.fetchall()}

def fetch_student_theory_averages(cursor, teacher_username=None):
    params = []
    teacher_filter = ""
    if teacher_username:
        teacher_filter = " AND u.teacher_username = ?"
        params.append(teacher_username)
    cursor.execute(
        f"""
        SELECT best_scores.username, ROUND(AVG(best_scores.best_pct), 1)
        FROM (
            SELECT username, test_id, MAX(percentage) AS best_pct
            FROM theory_submissions
            GROUP BY username, test_id
        ) best_scores
        JOIN users u ON u.username = best_scores.username
        WHERE u.role = 'student'{teacher_filter}
        GROUP BY best_scores.username
        """,
        params,
    )
    return {username: average for username, average in cursor.fetchall()}

def fetch_group_practical_averages(cursor, groups):
    if not groups:
        return {}
    placeholders = ",".join("?" for _ in groups)
    cursor.execute(
        f"""
        SELECT u.group_name, ROUND(AVG(best_scores.best_score), 1)
        FROM (
            SELECT username, subject, task, MAX(score) AS best_score
            FROM results
            GROUP BY username, subject, task
        ) best_scores
        JOIN users u ON u.username = best_scores.username
        WHERE u.role = 'student' AND u.group_name IN ({placeholders})
        GROUP BY u.group_name
        """,
        groups,
    )
    return {group: average for group, average in cursor.fetchall()}

def fetch_group_theory_averages(cursor, groups):
    if not groups:
        return {}
    placeholders = ",".join("?" for _ in groups)
    cursor.execute(
        f"""
        SELECT u.group_name, ROUND(AVG(best_scores.best_pct), 1)
        FROM (
            SELECT username, test_id, MAX(percentage) AS best_pct
            FROM theory_submissions
            GROUP BY username, test_id
        ) best_scores
        JOIN users u ON u.username = best_scores.username
        WHERE u.role = 'student' AND u.group_name IN ({placeholders})
        GROUP BY u.group_name
        """,
        groups,
    )
    return {group: average for group, average in cursor.fetchall()}

def get_teacher_quick_action_catalog():
    return [
        {"key": "manage_subjects", "label": "Practical", "icon": "📁", "href": "/manage_subjects", "kind": "link"},
        {"key": "manage_lessons", "label": "Lesson Setup", "icon": "📘", "href": "/manage_lessons", "kind": "link"},
        {"key": "manage_tests", "label": "Theory Tests", "icon": "📝", "href": "/manage_tests", "kind": "link"},
        {"key": "response_review", "label": "Review Responses", "icon": "🧾", "href": "/response_review", "kind": "link"},
        {"key": "marking_setup", "label": "Marking Setup", "icon": "🛠️", "href": "/marking_setup", "kind": "link"},
        {"key": "attendance", "label": "Attendance", "icon": "📅", "href": "/attendance", "kind": "link"},
        {"key": "group_results", "label": "Results", "icon": "📊", "href": "/group_results", "kind": "link"},
        {"key": "communications", "label": "Messages", "icon": "💬", "href": "/communications", "kind": "link"},
        {"key": "export", "label": "Export", "icon": "⬇️", "kind": "export"},
        {"key": "view_group", "label": "View Group", "icon": "👁", "kind": "view_group"},
    ]

def get_teacher_selected_quick_actions(username):
    catalog = get_teacher_quick_action_catalog()
    catalog_map = {item["key"]: item for item in catalog}
    default_keys = ["manage_subjects", "manage_tests", "attendance", "communications", "response_review", "export", "view_group"]
    conn = get_db()
    cursor = conn.cursor()
    cursor.execute(
        """
        SELECT action_key
        FROM teacher_quick_actions
        WHERE username = ?
        ORDER BY rowid
        """,
        (username,),
    )
    stored_keys = [row[0] for row in cursor.fetchall()]
    conn.close()
    selected = [catalog_map[key] for key in stored_keys if key in catalog_map]
    if not selected:
        selected = [catalog_map[key] for key in default_keys if key in catalog_map]
    selected_keys = {item["key"] for item in selected}
    available = [dict(item, selected=item["key"] in selected_keys) for item in catalog]
    return selected, available

def get_teacher_unread_message_count(username):
    conn = get_db()
    cursor = conn.cursor()
    cursor.execute(
        """
        SELECT COUNT(*)
        FROM communication_threads t
        WHERE t.teacher_username = ?
          AND COALESCE(t.last_sender_role, '') = 'student'
        """,
        (username,),
    ) if False else None
    cursor.execute(
        """
        SELECT COUNT(*)
        FROM communication_threads t
        WHERE t.teacher_username = ?
          AND EXISTS (
              SELECT 1
              FROM communication_messages m
              WHERE m.thread_id = t.id
                AND COALESCE(m.sender_role, '') = 'student'
                AND (t.teacher_read_at IS NULL OR COALESCE(m.created_at, '') > t.teacher_read_at)
          )
        """,
        (username,),
    )
    count = cursor.fetchone()[0] or 0
    conn.close()
    return count

def get_student_unread_message_count(username):
    conn = get_db()
    cursor = conn.cursor()
    cursor.execute(
        """
        SELECT COUNT(*)
        FROM communication_threads t
        WHERE t.student_username = ?
          AND EXISTS (
              SELECT 1
              FROM communication_messages m
              WHERE m.thread_id = t.id
                AND COALESCE(m.sender_role, '') IN ('teacher', 'admin')
                AND (t.student_read_at IS NULL OR COALESCE(m.created_at, '') > t.student_read_at)
          )
        """,
        (username,),
    )
    count = cursor.fetchone()[0] or 0
    conn.close()
    return count

def student_has_fresh_teacher_reply(username):
    return get_student_unread_message_count(username) > 0

def get_student_message_threads(username):
    conn = get_db()
    conn.row_factory = sqlite3.Row
    cursor = conn.cursor()
    cursor.execute("""
        SELECT *
        FROM communication_threads
        WHERE student_username = ?
        ORDER BY COALESCE(updated_at, created_at) DESC, id DESC
    """, (username,))
    threads = []
    for row in cursor.fetchall():
        thread = dict(row)
        cursor.execute("""
            SELECT *
            FROM communication_messages
            WHERE thread_id = ?
            ORDER BY COALESCE(created_at, '') ASC, id ASC
        """, (thread["id"],))
        thread["messages"] = [dict(message_row) for message_row in cursor.fetchall()]
        threads.append(thread)
    conn.close()
    return threads

def mark_student_threads_read(username):
    now_text = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("""
        UPDATE communication_threads
        SET student_read_at = ?
        WHERE student_username = ?
    """, (now_text, username))
    conn.commit()
    conn.close()

def create_communication_thread(cursor, student_username, topic, subject_line="", attendance_date="", initial_message="", chat_session_id=""):
    created_at = datetime.now().isoformat()
    cursor.execute("SELECT full_name, group_name, teacher_username FROM users WHERE username = ?", (student_username,))
    user_row = cursor.fetchone()
    group_name = user_row[1] if user_row else ""
    teacher_username = user_row[2] if user_row else ""
    cursor.execute(
        """
        INSERT INTO communication_threads (
            student_username, teacher_username, group_name, topic, subject_line,
            attendance_date, status, chat_session_id, created_at, updated_at
        )
        VALUES (?, ?, ?, ?, ?, ?, 'open', ?, ?, ?)
        """,
        (student_username, teacher_username, group_name, topic, subject_line, attendance_date, chat_session_id, created_at, created_at),
    )
    thread_id = cursor.lastrowid
    if initial_message:
        add_communication_message(cursor, thread_id, student_username, "student", initial_message)
    return thread_id

def add_communication_message(cursor, thread_id, sender_username, sender_role, message):
    created_at = datetime.now().isoformat()
    cursor.execute(
        """
        INSERT INTO communication_messages (thread_id, sender_username, sender_role, message, created_at)
        VALUES (?, ?, ?, ?, ?)
        """,
        (thread_id, sender_username, sender_role, message, created_at),
    )
    cursor.execute(
        """
        UPDATE communication_threads
        SET updated_at = ?,
            teacher_read_at = CASE WHEN ? IN ('teacher', 'admin') THEN ? ELSE teacher_read_at END,
            student_read_at = CASE WHEN ? = 'student' THEN ? ELSE student_read_at END
        WHERE id = ?
        """,
        (created_at, sender_role, created_at, sender_role, created_at, thread_id),
    )

def normalize_review_text(value):
    return " ".join((value or "").strip().lower().split())

def get_true_false_option_data(options):
    correct_choice = ""
    accepted_corrections = []
    for option in options:
        option_text = option[1] if len(option) > 1 else ""
        is_correct = option[2] if len(option) > 2 else 0
        match_pair = option[3] if len(option) > 3 else None
        if is_correct == 1 and match_pair != "correction":
            correct_choice = option_text
        if match_pair == "correction" and option_text:
            accepted_corrections.append(option_text)
    return correct_choice, accepted_corrections

def score_true_false_answer(selected, correction_submitted, options, effective_marks):
    correct_choice, accepted_corrections = get_true_false_option_data(options)
    if selected != correct_choice:
        return 0
    if selected == "False" and accepted_corrections:
        submitted_key = normalize_review_text(correction_submitted)
        accepted_keys = {normalize_review_text(item) for item in accepted_corrections if item}
        return effective_marks if submitted_key and submitted_key in accepted_keys else 0
    return effective_marks

def get_fill_in_accepted_answers(options):
    return [option[1] for option in options if len(option) > 2 and option[2] == 1 and option[1]]

def score_fill_in_answer(answer_text, options, marks):
    submitted_key = normalize_review_text(answer_text)
    accepted_keys = {normalize_review_text(item) for item in get_fill_in_accepted_answers(options)}
    return marks if submitted_key and submitted_key in accepted_keys else 0

def parse_true_false_answer_text(answer_text):
    raw = (answer_text or "").strip()
    if "(correction:" in raw:
        selected, correction = raw.split("(correction:", 1)
        return selected.strip(), correction.rstrip(") ").strip()
    return raw, ""

def compute_theory_answer_award(question_type, marks, options, answer_text):
    if question_type == "mcq_single":
        correct_values = {option[1] for option in options if len(option) > 2 and option[2] == 1}
        return marks if answer_text in correct_values else 0
    if question_type == "mcq_multi":
        correct_values = sorted(option[1] for option in options if len(option) > 2 and option[2] == 1)
        selected_values = sorted([item.strip() for item in (answer_text or "").split(",") if item.strip()])
        return marks if selected_values == correct_values else 0
    if question_type == "true_false":
        selected, correction = parse_true_false_answer_text(answer_text)
        return score_true_false_answer(selected, correction, options, marks)
    if question_type == "fill_in":
        return score_fill_in_answer(answer_text, options, marks)
    if question_type == "match":
        expected_map = {option[1]: option[3] for option in options if len(option) > 3 and option[3] and option[3] != "correction"}
        awarded = 0
        for pair in (answer_text or "").split(";"):
            if "=" not in pair:
                continue
            left, chosen = pair.split("=", 1)
            if expected_map.get(left.strip()) == chosen.strip():
                awarded += 1
        return awarded
    return 0

QUESTION_BANK_SUPPORTED_TYPES = ["mcq_single", "fill_in", "true_false", "match"]

def parse_module_names(raw_value):
    raw_value = (raw_value or "").replace("\r", "\n")
    parts = []
    for chunk in raw_value.replace(";", ",").split(","):
        for line in chunk.split("\n"):
            cleaned = line.strip()
            if cleaned:
                parts.append(cleaned)
    seen = set()
    ordered = []
    for item in parts:
        key = item.lower()
        if key in seen:
            continue
        seen.add(key)
        ordered.append(item)
    return ordered

def build_bank_option_signature(question_type, options):
    normalized = []
    for option_text, is_correct, match_pair in options:
        normalized.append((
            (option_text or "").strip().lower(),
            safe_int(is_correct, 0),
            (match_pair or "").strip().lower(),
        ))
    if question_type == "match":
        normalized.sort()
    return tuple(normalized)

def bank_question_exists(cursor, question_text, question_type, subject, modules, options):
    normalized_text = (question_text or "").strip().lower()
    normalized_subject = (subject or "").strip().lower()
    normalized_modules = (modules or "").strip().lower()
    target_signature = build_bank_option_signature(question_type, options)

    cursor.execute("""
        SELECT id
        FROM question_bank_questions
        WHERE LOWER(TRIM(question_text)) = ?
          AND question_type = ?
          AND LOWER(TRIM(COALESCE(subject, ''))) = ?
          AND LOWER(TRIM(COALESCE(modules, ''))) = ?
    """, (normalized_text, question_type, normalized_subject, normalized_modules))
    candidate_ids = [row[0] for row in cursor.fetchall()]
    for candidate_id in candidate_ids:
        cursor.execute("""
            SELECT option_text, is_correct, match_pair
            FROM question_bank_options
            WHERE bank_question_id = ?
            ORDER BY id
        """, (candidate_id,))
        candidate_signature = build_bank_option_signature(question_type, cursor.fetchall())
        if candidate_signature == target_signature:
            return True
    return False

def get_question_bank_counts(cursor, modules=None, subjects=None):
    modules = parse_module_names(",".join(modules or []))
    subjects = [item.strip() for item in (subjects or []) if item and item.strip()]
    counts = {}
    for q_type in QUESTION_BANK_SUPPORTED_TYPES:
        where_parts = ["question_type = ?"]
        params = [q_type]
        if modules:
            module_filters = " OR ".join(["LOWER(COALESCE(modules, '')) LIKE ?" for _ in modules])
            where_parts.append(f"({module_filters})")
            params.extend(f"%{module.lower()}%" for module in modules)
        if subjects:
            subject_filters = " OR ".join(["LOWER(COALESCE(subject, '')) = ?" for _ in subjects])
            where_parts.append(f"({subject_filters})")
            params.extend(subject.lower() for subject in subjects)
        cursor.execute(f"""
            SELECT COUNT(*)
            FROM question_bank_questions
            WHERE {' AND '.join(where_parts)}
        """, params)
        counts[q_type] = cursor.fetchone()[0] or 0
    return counts

def clone_bank_question_to_test(cursor, bank_question_id, test_id, order_index):
    cursor.execute("""
        SELECT question_text, question_type, marks
        FROM question_bank_questions
        WHERE id = ?
    """, (bank_question_id,))
    row = cursor.fetchone()
    if not row:
        return None
    question_text, question_type, marks = row
    cursor.execute("""
        INSERT INTO theory_questions (test_id, question_text, question_type, marks, order_index)
        VALUES (?, ?, ?, ?, ?)
    """, (test_id, question_text, question_type, marks, order_index))
    new_question_id = cursor.lastrowid
    cursor.execute("""
        SELECT option_text, is_correct, match_pair
        FROM question_bank_options
        WHERE bank_question_id = ?
        ORDER BY id
    """, (bank_question_id,))
    for option_text, is_correct, match_pair in cursor.fetchall():
        cursor.execute("""
            INSERT INTO theory_options (question_id, option_text, is_correct, match_pair)
            VALUES (?, ?, ?, ?)
        """, (new_question_id, option_text, is_correct, match_pair))
    return new_question_id

def create_generated_match_question(cursor, bank_question_ids, test_id, order_index):
    if not bank_question_ids:
        return None

    placeholders = ",".join("?" for _ in bank_question_ids)
    cursor.execute(f"""
        SELECT id, question_text, marks
        FROM question_bank_questions
        WHERE id IN ({placeholders})
    """, bank_question_ids)
    question_rows = cursor.fetchall()
    if not question_rows:
        return None

    question_map = {row[0]: row for row in question_rows}
    ordered_rows = [question_map[q_id] for q_id in bank_question_ids if q_id in question_map]
    if not ordered_rows:
        return None

    question_text = ordered_rows[0][1] or "Match Column A to B"
    total_marks = sum(max(1, safe_int(row[2], 1)) for row in ordered_rows)
    cursor.execute("""
        INSERT INTO theory_questions (test_id, question_text, question_type, marks, order_index)
        VALUES (?, ?, 'match', ?, ?)
    """, (test_id, question_text, total_marks, order_index))
    new_question_id = cursor.lastrowid

    for bank_question_id in bank_question_ids:
        cursor.execute("""
            SELECT option_text, is_correct, match_pair
            FROM question_bank_options
            WHERE bank_question_id = ? AND COALESCE(match_pair, '') != ''
            ORDER BY id
            LIMIT 1
        """, (bank_question_id,))
        option_row = cursor.fetchone()
        if not option_row:
            continue
        option_text, is_correct, match_pair = option_row
        cursor.execute("""
            INSERT INTO theory_options (question_id, option_text, is_correct, match_pair)
            VALUES (?, ?, ?, ?)
        """, (new_question_id, option_text, is_correct, match_pair))
    return new_question_id

def regrade_theory_question_answers(test_id, question_id, selected_group=None, learner_username=None):
    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("""
        SELECT question_type, marks
        FROM theory_questions
        WHERE id = ? AND test_id = ?
    """, (question_id, test_id))
    question_row = cursor.fetchone()
    if not question_row:
        conn.close()
        return 0

    q_type, marks = question_row
    cursor.execute("""
        SELECT id, option_text, is_correct, match_pair
        FROM theory_options
        WHERE question_id = ?
        ORDER BY id
    """, (question_id,))
    options = cursor.fetchall()

    query = """
        SELECT a.id, a.answer_text, a.submission_id
        FROM theory_answers a
        JOIN theory_submissions s ON s.id = a.submission_id
        JOIN users u ON u.username = s.username
        WHERE s.test_id = ? AND a.question_id = ?
    """
    params = [test_id, question_id]
    if selected_group:
        query += " AND u.group_name = ?"
        params.append(selected_group)
    if learner_username:
        query += " AND s.username = ?"
        params.append(learner_username)
    cursor.execute(query, params)
    answer_rows = cursor.fetchall()

    updated = 0
    submission_ids = set()
    for answer_id, answer_text, submission_id in answer_rows:
        awarded = compute_theory_answer_award(q_type, marks, options, answer_text)
        is_correct = 1 if awarded == marks else 0
        if q_type == "match":
            expected_total = len([option for option in options if len(option) > 3 and option[3] and option[3] != "correction"])
            is_correct = 1 if awarded == expected_total else 0
        cursor.execute("""
            UPDATE theory_answers
            SET marks_awarded = ?, is_correct = ?
            WHERE id = ?
        """, (awarded, is_correct, answer_id))
        updated += cursor.rowcount
        submission_ids.add(submission_id)

    for submission_id in submission_ids:
        cursor.execute("""
            SELECT COALESCE(SUM(marks_awarded), 0)
            FROM theory_answers
            WHERE submission_id = ?
        """, (submission_id,))
        score = cursor.fetchone()[0] or 0
        cursor.execute("SELECT total FROM theory_submissions WHERE id = ?", (submission_id,))
        total_row = cursor.fetchone()
        total = total_row[0] if total_row and total_row[0] else 0
        percentage = round((score / total) * 100) if total else 0
        cursor.execute("""
            UPDATE theory_submissions
            SET score = ?, percentage = ?
            WHERE id = ?
        """, (score, percentage, submission_id))

    conn.commit()
    conn.close()
    return updated

def get_student_dashboard_data(username):
    conn = get_db()
    cursor = conn.cursor()

    cursor.execute("""
    SELECT subject, ROUND(AVG(score), 1)
    FROM results
    WHERE username = ?
    GROUP BY subject
    """, (username,))
    subject_averages = cursor.fetchall()

    cursor.execute("""
    SELECT ROUND(AVG(score), 1)
    FROM results
    WHERE username = ?
    """, (username,))

    result = cursor.fetchone()
    overall_avg = result[0] if result and result[0] is not None else 0

    cursor.execute("""
    SELECT subject, task, score
    FROM results
    WHERE username = ?
    ORDER BY id DESC
    LIMIT 5
    """, (username,))
    recent_results = cursor.fetchall()

    conn.close()

    return (
        subject_averages or [],
        overall_avg,
        recent_results or []
    )

def get_overall_average(username):
    conn = get_db()
    cursor = conn.cursor()

    
    cursor.execute("""
    SELECT AVG(score)
    FROM results
    WHERE username = ?
    """, (username,))

    result = cursor.fetchone()
    conn.close()

    return round(result[0], 1) if result[0] else 0

def get_low_attendance_learners(limit=10, groups=None, teacher_username=None):
    days = get_last_21_days()
    conn = get_db()
    cursor = conn.cursor()

    # Get excluded dates (global only for cross-group summary)
    cursor.execute("SELECT date FROM excluded_dates WHERE group_name IS NULL")
    excluded = {row[0] for row in cursor.fetchall()}
    days = [d for d in days if d not in excluded]

    if teacher_username:
        cursor.execute(
            "SELECT username, full_name FROM users WHERE role = 'student' AND teacher_username = ?",
            (teacher_username,))
    elif groups:
        placeholders = ",".join("?" for _ in groups)
        cursor.execute(f"SELECT username, full_name FROM users WHERE role = 'student' AND group_name IN ({placeholders})", groups)
    else:
        cursor.execute("SELECT username, full_name FROM users WHERE role = 'student'")
    learners = cursor.fetchall()

    results = []
    for username, full_name in learners:
        present = 0
        for day in days:
            cursor.execute("SELECT 1 FROM login_history WHERE username = ? AND date = ?", (username, day))
            if cursor.fetchone():
                present += 1
                continue
            cursor.execute("SELECT status FROM attendance_override WHERE username = ? AND date = ?", (username, day))
            override = cursor.fetchone()
            if override and override[0] == "present":
                present += 1

        absent = len(days) - present
        results.append((full_name or username, username, absent))

    conn.close()
    results.sort(key=lambda x: x[2], reverse=True)
    return results[:limit]


def get_teachers():
    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("SELECT username, full_name FROM users WHERE role = 'teacher' ORDER BY full_name")
    teachers = cursor.fetchall()
    conn.close()
    return teachers


@app.route("/marking_setup", methods=["GET", "POST"])
def marking_setup():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    error = None
    message = None
    generated_json_preview = ""
    llm_output_preview = ""
    llm_normalized_output_preview = ""
    generation_warnings_preview = []
    llm_status = ""

    conn = get_marking_db()
    cursor = conn.cursor()

    if request.method == "POST":
        action = request.form.get("action")
        if action == "delete_one":
            setup_id = request.form.get("setup_id")
            cursor.execute("DELETE FROM marking_setups WHERE id = ?", (setup_id,))
            conn.commit()
            message = "Marking setup deleted."
        elif action == "delete_selected":
            setup_ids = [sid for sid in request.form.getlist("setup_ids") if sid.strip()]
            if setup_ids:
                placeholders = ",".join("?" for _ in setup_ids)
                cursor.execute(f"DELETE FROM marking_setups WHERE id IN ({placeholders})", setup_ids)
                conn.commit()
                message = f"Deleted {len(setup_ids)} marking setup(s)."
            else:
                error = "Select at least one marking setup to delete."
        else:
            title = (request.form.get("title") or "").strip()
            notes = (request.form.get("notes") or "").strip()
            llm_model = (request.form.get("llm_model") or "").strip() or None
            question_paper_file = request.files.get("question_paper_file")
            memo_file = request.files.get("memo_file")

            if not title or not question_paper_file or not question_paper_file.filename or not memo_file or not memo_file.filename:
                error = "Title, question paper, and memo are required."
            else:
                from Marking_Experiment.generation_service import generate_marking_task_json

                qp_suffix = Path(question_paper_file.filename).suffix or ".docx"
                memo_suffix = Path(memo_file.filename).suffix or ".docx"
                qp_temp = None
                memo_temp = None
                try:
                    with tempfile.NamedTemporaryFile(delete=False, suffix=qp_suffix) as qp_handle:
                        question_paper_blob = question_paper_file.read()
                        qp_handle.write(question_paper_blob)
                        qp_temp = Path(qp_handle.name)
                    with tempfile.NamedTemporaryFile(delete=False, suffix=memo_suffix) as memo_handle:
                        memo_blob = memo_file.read()
                        memo_handle.write(memo_blob)
                        memo_temp = Path(memo_handle.name)

                    generated = generate_marking_task_json(qp_temp, memo_temp, title, model=llm_model)
                    task_definition = generated["task_definition"]
                    generated_json_text = json.dumps(task_definition, indent=2, ensure_ascii=False)
                    generated_json_preview = generated_json_text
                    llm_output_preview = generated.get("llm_output") or ""
                    llm_normalized_output_preview = generated.get("llm_normalized_output") or ""
                    generation_warnings_preview = generated.get("warnings") or []
                    llm_status = generated.get("llm_status") or ""

                    now = datetime.now().isoformat()
                    cursor.execute(
                        """
                        INSERT INTO marking_setups (
                            title, created_by, created_at, updated_at,
                            question_paper_filename, question_paper_blob,
                            memo_filename, memo_blob,
                            json_script_filename, json_script_blob,
                            llm_model, generation_warnings, generation_error, notes
                        ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                        """,
                        (
                            title, username, now, now,
                            secure_filename(question_paper_file.filename), question_paper_blob,
                            secure_filename(memo_file.filename), memo_blob,
                            f"{secure_filename(title) or 'marking_setup'}.json", generated_json_text.encode("utf-8"),
                            llm_model, json.dumps(generation_warnings_preview, ensure_ascii=False), None, notes,
                        ),
                    )
                    conn.commit()
                    message = "Marking setup generated and saved."
                except Exception as exc:
                    error = f"Failed to generate marking setup: {exc}"
                finally:
                    for temp_path in (qp_temp, memo_temp):
                        if temp_path and temp_path.exists():
                            temp_path.unlink()

    cursor.execute("""
        SELECT id, title, created_by, created_at,
               question_paper_filename, memo_filename, json_script_filename, notes
        FROM marking_setups
        ORDER BY created_at DESC, id DESC
    """)
    setups = cursor.fetchall()
    conn.close()

    return render_template(
        "marking_setup.html",
        setups=setups,
        error=error,
        message=message,
        generated_json_preview=generated_json_preview,
        llm_output_preview=llm_output_preview,
        llm_normalized_output_preview=llm_normalized_output_preview,
        generation_warnings_preview=generation_warnings_preview,
        llm_status=llm_status,
    )


@app.route("/marking_setup/<int:setup_id>/download/<field>")
def download_marking_blob(setup_id, field):
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    field_map = {
        "question_paper": ("question_paper_filename", "question_paper_blob"),
        "memo": ("memo_filename", "memo_blob"),
        "json_script": ("json_script_filename", "json_script_blob"),
    }
    if field not in field_map:
        return "File type not found", 404

    filename_field, blob_field = field_map[field]
    conn = get_marking_db()
    cursor = conn.cursor()
    cursor.execute(f"SELECT {filename_field}, {blob_field} FROM marking_setups WHERE id = ?", (setup_id,))
    row = cursor.fetchone()
    conn.close()
    if not row or not row[1]:
        return "File not found", 404

    return send_file(
        io.BytesIO(row[1]),
        as_attachment=True,
        download_name=row[0] or f"{field}_{setup_id}",
        mimetype="application/octet-stream",
    )


@app.route("/marking_setup/<int:setup_id>/edit", methods=["GET", "POST"])
def edit_marking_setup(setup_id):
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    conn = get_marking_db()
    cursor = conn.cursor()
    cursor.execute("SELECT title, notes, json_script_blob FROM marking_setups WHERE id = ?", (setup_id,))
    row = cursor.fetchone()
    if not row:
        conn.close()
        return "Marking setup not found", 404

    title, notes, json_blob = row
    error = None
    message = None
    try:
        task_definition = json.loads((json_blob or b"{}").decode("utf-8") if isinstance(json_blob, bytes) else (json_blob or "{}"))
    except Exception:
        task_definition = {}
    questions = task_definition.get("questions", []) if isinstance(task_definition, dict) else []

    if request.method == "POST":
        try:
            row_count = safe_int(request.form.get("row_count"), len(questions))
            updated_questions = []
            for idx in range(row_count):
                if not request.form.get(f"enabled_{idx}"):
                    continue
                target_json = request.form.get(f"target_{idx}", "{}").strip() or "{}"
                expected_json = request.form.get(f"expected_{idx}", "true").strip() or "true"
                updated_questions.append({
                    "question_number": request.form.get(f"question_number_{idx}", "").strip(),
                    "description": request.form.get(f"description_{idx}", "").strip(),
                    "domain": request.form.get(f"domain_{idx}", "").strip(),
                    "type": request.form.get(f"type_{idx}", "").strip(),
                    "target": json.loads(target_json),
                    "expected": json.loads(expected_json),
                    "marks": safe_int(request.form.get(f"marks_{idx}"), 1),
                })

            task_definition["task_name"] = (request.form.get("task_name") or title).strip()
            task_definition["program"] = (request.form.get("program") or task_definition.get("program") or "word").strip()
            task_definition["file"] = (request.form.get("file") or task_definition.get("file") or "student_file.docx").strip()
            task_definition["questions"] = updated_questions
            task_definition["total_marks"] = sum(int(question.get("marks", 1)) for question in updated_questions)

            encoded = json.dumps(task_definition, indent=2, ensure_ascii=False).encode("utf-8")
            notes = request.form.get("notes", "")
            cursor.execute(
                "UPDATE marking_setups SET title = ?, notes = ?, json_script_blob = ?, updated_at = ? WHERE id = ?",
                (task_definition["task_name"], notes, encoded, datetime.now().isoformat(), setup_id),
            )
            conn.commit()
            title = task_definition["task_name"]
            questions = updated_questions
            message = "Marking setup saved."
        except Exception as exc:
            error = f"Could not save marking setup: {exc}"

    conn.close()
    return render_template(
        "edit_marking_setup.html",
        title=title,
        notes=notes,
        task_definition=task_definition,
        questions=questions,
        error=error,
        message=message,
    )


def get_marking_scripts():
    """Return a list of available marking script names from marking/tasks/."""
    tasks_dir = os.path.join(os.path.dirname(__file__), "marking", "tasks")
    scripts = []
    if os.path.exists(tasks_dir):
        for f in sorted(os.listdir(tasks_dir)):
            if f.endswith(".py") and f != "__init__.py":
                scripts.append(f[:-3])  # strip .py
    return scripts

def get_groups(username=None):
    conn = get_db()
    cursor = conn.cursor()

    if username:
        role = get_user_role(username)
        if role == 'teacher':
            group_set = set()
            cursor.execute("SELECT group_name FROM group_teachers WHERE teacher_username = ?", (username,))
            group_set.update(g[0] for g in cursor.fetchall() if g[0])
            cursor.execute("SELECT DISTINCT group_name FROM users WHERE role = 'student' AND teacher_username = ? AND group_name IS NOT NULL", (username,))
            group_set.update(g[0] for g in cursor.fetchall() if g[0])
            conn.close()
            return sorted(group_set)

    cursor.execute("SELECT DISTINCT group_name FROM users WHERE group_name IS NOT NULL")
    groups = [g[0] for g in cursor.fetchall()]

    conn.close()
    return groups

def add_learner_note_entry(cursor, username, note, created_by, flag=""):
    cursor.execute("""
        INSERT INTO learner_notes (username, note, flag, created_by, created_at)
        VALUES (?, ?, ?, ?, ?)
    """, (username, note, flag, created_by, datetime.now().isoformat()))


def get_group_late_threshold(group, date):
    conn = get_db()
    cursor = conn.cursor()

    cursor.execute("""
    SELECT MIN(login_time)
    FROM login_history
    WHERE date = ?
      AND username IN (
          SELECT username FROM users WHERE group_name = ?
      )
    """, (date, group))
    min_login = cursor.fetchone()[0]
    conn.close()

    if not min_login:
        return None

    try:
        min_dt = datetime.fromisoformat(min_login)
    except ValueError:
        min_dt = datetime.strptime(min_login, "%Y-%m-%d %H:%M:%S")

    late_cutoff = min_dt + timedelta(minutes=4)
    return late_cutoff.strftime("%H:%M")


def is_attendance_editable(date_str, role='teacher'):
    """Check if attendance for a date can be edited (not locked after 10 days)."""
    if role == 'admin':  # Admins can always edit
        return True
    try:
        date_obj = datetime.strptime(date_str, '%Y-%m-%d').date()
        days_diff = (datetime.now().date() - date_obj).days
        return days_diff <= 10
    except ValueError:
        return False  # Invalid date, treat as locked


def get_attendance_data(group, start_date=None, end_date=None, teacher_username=None):
    conn = get_db()
    cursor = conn.cursor()

    # If no dates provided, use active term range or fall back to last 7 days
    if not start_date or not end_date:
        term_range = get_active_term_range()
        if term_range:
            start_date, end_date = term_range
            current = datetime.strptime(start_date, "%Y-%m-%d").date()
            end = datetime.strptime(end_date, "%Y-%m-%d").date()
            days = []
            while current <= end:
                if current.weekday() < 5:
                    days.append(current.strftime("%Y-%m-%d"))
                current += timedelta(days=1)
        else:
            days = get_last_7_days()
    else:
        # Generate business days (Mon-Fri) between start_date and end_date
        days = []
        current = datetime.strptime(start_date, "%Y-%m-%d").date()
        end = datetime.strptime(end_date, "%Y-%m-%d").date()
        while current <= end:
            if current.weekday() < 5:
                days.append(current.strftime("%Y-%m-%d"))
            current += timedelta(days=1)

    # Remove future dates
    today_str = datetime.now().date().isoformat()
    days = [day for day in days if day <= today_str]

    auto_exclude_empty_attendance_days(cursor, [group], created_by=teacher_username or "system", days=days)
    cursor.execute("""
    SELECT date FROM excluded_dates 
    WHERE group_name IS NULL OR group_name = ?
    """, (group,))
    excluded_dates = {row[0] for row in cursor.fetchall()}
    days = [day for day in days if day not in excluded_dates]

    year_days = [day for day in get_current_year_attendance_days() if day <= today_str and day not in excluded_dates]
    auto_exclude_empty_attendance_days(cursor, [group], created_by=teacher_username or "system", days=year_days)
    cursor.execute("""
    SELECT date FROM excluded_dates 
    WHERE group_name IS NULL OR group_name = ?
    """, (group,))
    excluded_dates = {row[0] for row in cursor.fetchall()}
    days = [day for day in days if day not in excluded_dates]
    year_days = [day for day in get_current_year_attendance_days() if day <= today_str and day not in excluded_dates]

    query = """
    SELECT username, full_name, group_name
    FROM users
    WHERE group_name = ? AND role = 'student'
    """
    params = [group]
    if teacher_username:
        query += " AND teacher_username = ?"
        params.append(teacher_username)
    
    cursor.execute(query, params)
    learners = cursor.fetchall()

    attendance = []
    late_cutoffs = fetch_group_late_thresholds(cursor, group, days)
    year_late_cutoffs = fetch_group_late_thresholds(cursor, group, year_days)
    usernames = [user for user, _, _ in learners]
    login_map = fetch_first_login_times(cursor, usernames, days)
    override_map = fetch_attendance_override_statuses(cursor, usernames, days)
    year_login_map = fetch_first_login_times(cursor, usernames, year_days)
    year_override_map = fetch_attendance_override_statuses(cursor, usernames, year_days)
    cursor.execute("""
        SELECT DISTINCT lh.date
        FROM login_history lh
        JOIN users u ON u.username = lh.username
        WHERE u.group_name = ? AND u.role = 'student'
    """, (group,))
    class_checked_dates = {row[0] for row in cursor.fetchall()}

    for user, name, user_group_name in learners:
        history = build_attendance_history(
            cursor,
            user,
            group,
            days,
            login_map=login_map,
            override_map=override_map,
            late_cutoffs=late_cutoffs,
            class_checked_dates=class_checked_dates,
        )
        year_history = build_attendance_history(
            cursor,
            user,
            group,
            year_days,
            login_map=year_login_map,
            override_map=year_override_map,
            late_cutoffs=year_late_cutoffs,
            class_checked_dates=class_checked_dates,
        )
        summary = summarize_attendance_history(history)
        year_summary = summarize_attendance_history(year_history)
        row = {
            "username": user,
            "name": name,
            "group": user_group_name,
            "days": {},
            "attendance_pct": summary["attendance_pct"],
            "present_days": summary["present_days"],
            "absent_days": year_summary["absent_days"],
            "year_present_days": year_summary["present_days"],
            "year_attendance_pct": year_summary["attendance_pct"],
            "year_total_days": year_summary["total_days"],
            "late_days": summary["late_days"],
        }
        for item in history:
            if item["status"] == "Present":
                row["days"][item["date"]] = {
                    "time": item["time"],
                    "late": item["late"],
                    "manual": item["note"] == "Manual"
                }
            elif item["status"] == "Absent":
                row["days"][item["date"]] = None
            else:
                row["days"][item["date"]] = None

        attendance.append(row)

    conn.commit()
    conn.close()
    return days, attendance

def import_users_from_excel():
    """Import users from Excel file Users/grade12.xlsx"""
    try:
        df = pd.read_excel("Users/grade12.xlsx")
        
        conn = get_db()
        cursor = conn.cursor()
        
        imported_count = 0
        updated_count = 0
        
        for _, row in df.iterrows():
            username = str(row["username"]).strip().upper()
            full_name = str(row["full_name"]).strip()
            group_name = str(row["group"]).strip()
            
            # Check if user already exists
            cursor.execute("SELECT username FROM users WHERE username = ?", (username,))
            existing = cursor.fetchone()
            
            cursor.execute("""
            INSERT INTO users (username, full_name, group_name, role)
            VALUES (?, ?, ?, 'student')
            ON CONFLICT(username) DO UPDATE SET
                full_name = excluded.full_name,
                group_name = excluded.group_name
            """, (username, full_name, group_name))
            
            if existing:
                updated_count += 1
            else:
                imported_count += 1
        
        conn.commit()
        conn.close()
        
        return f"Successfully imported {imported_count} new users and updated {updated_count} existing users."
    
    except FileNotFoundError:
        return "Error: Users/grade12.xlsx file not found."
    except Exception as e:
        return f"Error importing users: {str(e)}"

# 🔹 Login Page
@app.route("/", methods=["GET", "POST"])
def login():
    if request.method == "POST":
        username = request.form.get("username", "").strip().upper()

        if username and len(username) <= 20 and username.isalnum():
            session["username"] = username

            create_user_if_not_exists(username)
            update_last_active(username)
            if get_user_role(username) == "student" and should_record_attendance_login():
                log_login(username)
            log_activity(username, "logged in")

            with lock:
                active_users[username] = datetime.now()

            role = get_user_role(username)
            if role in ["teacher", "admin"]:
                return redirect(url_for("teacher_dashboard"))
            return redirect(url_for("student_dashboard"))
        else:
            return "Invalid username", 400

    return render_template("login.html")


# 🔹 Heartbeat (updates activity)
@app.route("/heartbeat/<username>")
def heartbeat(username):
    with lock:
        if username in active_users:
            active_users[username] = datetime.now()
            update_last_active(username)
        else:
            return "Invalid", 401
    return "OK"


def cleanup_thread():
    while True:
        now = datetime.now()
        with lock:
            to_remove = []
            for user, last_seen in list(active_users.items()):
                if now - last_seen > timedelta(seconds=TIMEOUT):
                    to_remove.append(user)
            for user in to_remove:
                del active_users[user]
        time.sleep(30)


@app.route("/student_dashboard")
def student_dashboard():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    update_active_user(username)
    role = get_user_role(username)

    # Teachers go to teacher dashboard
    if role in ["teacher", "admin"]:
        return redirect(url_for("teacher_dashboard"))

    conn = get_db()
    cursor = conn.cursor()

    cursor.execute("SELECT full_name, group_name FROM users WHERE username = ?", (username,))
    user_row = cursor.fetchone()
    display_name = user_row[0] if user_row and user_row[0] else username
    user_group = user_row[1] if user_row else None

    # ─ Academic (best score per task)
    cursor.execute("""
        SELECT subject, ROUND(AVG(best_score),1)
        FROM (SELECT subject, task, MAX(score) as best_score FROM results WHERE username = ? GROUP BY subject, task)
        GROUP BY subject
    """, (username,))
    subject_avgs = cursor.fetchall()

    cursor.execute("""
        SELECT ROUND(AVG(best_score),1)
        FROM (SELECT subject, task, MAX(score) as best_score FROM results WHERE username = ? GROUP BY subject, task)
    """, (username,))
    overall_row = cursor.fetchone()
    practical_avg = overall_row[0] if overall_row and overall_row[0] else 0

    cursor.execute("""
        SELECT ROUND(AVG(best_pct),1)
        FROM (SELECT test_id, MAX(percentage) as best_pct FROM theory_submissions WHERE username = ? GROUP BY test_id)
    """, (username,))
    theory_avg_row = cursor.fetchone()
    theory_avg = theory_avg_row[0] if theory_avg_row and theory_avg_row[0] else None

    if practical_avg and theory_avg:
        overall_avg = round((practical_avg + theory_avg) / 2, 1)
    else:
        overall_avg = practical_avg or theory_avg or 0

    cursor.execute("""
        SELECT subject, task, MAX(score) as score, feedback, MAX(timestamp) as timestamp
        FROM results WHERE username = ? GROUP BY subject, task ORDER BY timestamp DESC LIMIT 5
    """, (username,))
    recent_results = cursor.fetchall()

    cursor.execute("""
        SELECT tt.subject, tt.title, MAX(ts.percentage) as best_pct, MAX(ts.submitted_at) as latest
        FROM theory_submissions ts
        JOIN theory_tests tt ON ts.test_id = tt.id
        WHERE ts.username = ? GROUP BY ts.test_id ORDER BY latest DESC LIMIT 5
    """, (username,))
    recent_theory_results = cursor.fetchall()

    cursor.execute("""
        SELECT tt.subject, ROUND(AVG(best_pct),1)
        FROM (SELECT test_id, MAX(percentage) as best_pct FROM theory_submissions WHERE username = ? GROUP BY test_id) b
        JOIN theory_tests tt ON b.test_id = tt.id
        WHERE tt.subject IS NOT NULL AND tt.subject != ''
        GROUP BY tt.subject
    """, (username,))
    theory_subject_avgs = cursor.fetchall()

    attendance_year = datetime.now().year
    attendance_days = get_current_year_attendance_days()
    auto_exclude_empty_attendance_days(cursor, [user_group], created_by=username, days=attendance_days)
    cursor.execute("SELECT date FROM excluded_dates WHERE group_name IS NULL OR group_name = ?", (user_group,))
    excluded_dates = {r[0] for r in cursor.fetchall()}
    login_map = fetch_first_login_times(cursor, [username], attendance_days)
    override_map = fetch_attendance_override_statuses(cursor, [username], attendance_days)
    late_cutoffs = fetch_group_late_thresholds(cursor, user_group, attendance_days) if user_group else {}
    cursor.execute("""
        SELECT DISTINCT lh.date
        FROM login_history lh
        JOIN users u ON u.username = lh.username
        WHERE u.group_name = ? AND u.role = 'student' AND lh.username != ?
    """, (user_group, username))
    class_checked_dates = {row[0] for row in cursor.fetchall()}
    att_history = build_attendance_history(
        cursor,
        username,
        user_group,
        attendance_days,
        login_map=login_map,
        override_map=override_map,
        late_cutoffs=late_cutoffs,
        class_checked_dates=class_checked_dates,
        excluded_dates=excluded_dates,
    )
    attendance_months = build_attendance_months(att_history)
    attendance_summary = summarize_attendance_history(att_history)
    att_pct = attendance_summary["attendance_pct"]

    # ─ Practical subjects and tasks for inline panel
    today = datetime.now().date().isoformat()
    cursor.execute("SELECT id, name FROM subjects ORDER BY name")
    all_subjects = cursor.fetchall()
    subjects_with_tasks = []
    for subj_id, subj_name in all_subjects:
        cursor.execute("""
            SELECT t.id, t.name
            FROM tasks t
            JOIN task_groups tg ON t.id = tg.task_id
            WHERE t.subject_id = ? AND tg.group_name = ? AND t.assign_date <= ?
            AND t.task_type = 'practical' AND t.is_active = 1
            ORDER BY t.name
        """, (subj_id, user_group, today))
        tasks_for_subj = cursor.fetchall()
        if tasks_for_subj:
            subjects_with_tasks.append({'id': subj_id, 'name': subj_name, 'tasks': tasks_for_subj})

    # ─ Missing tasks (practical + theory) assigned to the student's group
    today = datetime.now().date().isoformat()

    # Practical missing tasks
    cursor.execute("""
        SELECT t.id,
               s.name AS subject_name,
               t.name AS task_name,
               t.assign_date,
               t.marking_script
        FROM tasks t
        JOIN subjects s ON t.subject_id = s.id
        JOIN task_groups tg ON t.id = tg.task_id
        WHERE tg.group_name = ?
          AND t.assign_date <= ?
          AND t.task_type = 'practical'
          AND t.is_active = 1
          AND NOT EXISTS (
              SELECT 1 FROM results r
              WHERE r.username = ? AND r.subject = s.name AND r.task = t.name
          )
        ORDER BY t.assign_date, t.name
        LIMIT 10
    """, (user_group, today, username))
    practical_missing = cursor.fetchall()

    practical_missing_assignments = [
        {
            "type": "practical",
            "subject": row[1],
            "activity": row[2],
            "due": row[3],
            "task_id": row[0],
            "subject_id": None,
            "start_url": None,
        }
        for row in practical_missing
    ]

    # Need subject_id for upload link
    if practical_missing_assignments:
        task_ids = [a["task_id"] for a in practical_missing_assignments]
        placeholders = ",".join(["?"] * len(task_ids))
        cursor.execute(f"""
            SELECT id, subject_id
            FROM tasks
            WHERE id IN ({placeholders})
        """, task_ids)
        id_to_subject = {row[0]: row[1] for row in cursor.fetchall()}

        for a in practical_missing_assignments:
            a["subject_id"] = id_to_subject.get(a["task_id"])
            if a["subject_id"] is not None:
                a["start_url"] = f"/upload/{username}/{a['subject_id']}/{a['task_id']}"

    # Theory missing tasks for the group (no submission by this student)
    cursor.execute("""
        SELECT tt.id,
               tt.subject,
               tt.title,
               tt.time_limit,
               tt.allow_multiple,
               tt.max_attempts
        FROM theory_tests tt
        JOIN theory_test_groups ttg ON tt.id = ttg.test_id
        WHERE tt.is_active = 1
          AND ttg.group_name = ?
          AND NOT EXISTS (
              SELECT 1 FROM theory_submissions ts
              WHERE ts.username = ? AND ts.test_id = tt.id
          )
        ORDER BY tt.subject, tt.title
        LIMIT 10
    """, (user_group, username))
    theory_missing = cursor.fetchall()

    theory_missing_assignments = [
        {
            "type": "theory",
            "subject": row[1] or "Theory",
            "activity": row[2],
            "due": None,
            "test_id": row[0],
            "start_url": f"/take_test/{row[0]}",
        }
        for row in theory_missing
    ]

    missing_assignments = practical_missing_assignments + theory_missing_assignments
    # Backward-compatible variable for the existing template placeholder
    missing_tasks = [(None, a['subject'], a['activity'], a['due'] or '') for a in missing_assignments if a['type']=='practical']

    # If template expects best-effort ordering, put practical first then theory.
    missing_assignments = practical_missing_assignments + theory_missing_assignments


    # ─ Weaknesses
    cursor.execute("SELECT skill, count FROM weaknesses WHERE username = ? ORDER BY count DESC LIMIT 5", (username,))
    weaknesses = cursor.fetchall()

    # ─ Recent feedback (non-empty)
    cursor.execute("""
        SELECT subject, task, feedback, timestamp
        FROM results WHERE username = ? AND feedback IS NOT NULL AND feedback != ''
        ORDER BY timestamp DESC LIMIT 3
    """, (username,))
    recent_feedback = cursor.fetchall()

    message_threads = get_student_message_threads(username)
    unread_teacher_messages = get_student_unread_message_count(username)
    auto_open_chat = student_has_fresh_teacher_reply(username)

    conn.close()
    return render_template(
        "student_dashboard.html",
        username=username,
        display_name=display_name,
        overall_avg=overall_avg,
        practical_avg=practical_avg,
        theory_avg=theory_avg,
        subject_avgs=subject_avgs,
        theory_subject_avgs=theory_subject_avgs,
        recent_results=recent_results,
        recent_theory_results=recent_theory_results,
        att_history=att_history,
        attendance_months=attendance_months,
        attendance_year=attendance_year,
        att_pct=att_pct,
        present_days=attendance_summary["present_days"],
        absent_days=attendance_summary["absent_days"],
        late_days=attendance_summary["late_days"],
        total_attendance_days=attendance_summary["total_days"],
        subjects_with_tasks=subjects_with_tasks,
        missing_assignments=missing_assignments,
        missing_tasks=missing_tasks,
        weaknesses=weaknesses,
        recent_feedback=recent_feedback,
        message_threads=message_threads,
        unread_teacher_messages=unread_teacher_messages,
        auto_open_chat=auto_open_chat
    )

@app.route("/auto_login")

def auto_login():
    username = request.args.get("username", "").strip().upper()

    if username and len(username) <= 20 and username.isalnum():
        session["username"] = username   # ✅ IMPORTANT

        create_user_if_not_exists(username)
        update_last_active(username)
        if get_user_role(username) == "student" and should_record_attendance_login():
            log_login(username)
        log_activity(username, "logged in")

        with lock:
            active_users[username] = datetime.now()

        role = get_user_role(username)
        if role in ["teacher", "admin"]:
            return redirect(url_for("teacher_dashboard"))
        return redirect(url_for("student_dashboard"))

    return "Invalid auto-login", 400

@app.route("/logout")
def logout():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    if username:
        with lock:
            active_users.pop(username, None)

    session.clear()   # ✅ IMPORTANT

    return redirect(url_for("login"))



# 🔹 Teacher View (see active users)
@app.route("/active")
def active():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    if not username:
        return redirect(url_for("login"))

    role = get_user_role(username)

    if role not in ["teacher", "admin"]:
        return "Access denied", 403

    now = datetime.now()

    # Remove inactive users
    with lock:
        active_list = []
        to_remove = []
        for user, last_seen in list(active_users.items()):
            if now - last_seen < timedelta(seconds=TIMEOUT):
                active_list.append(user)
            else:
                to_remove.append(user)
        for user in to_remove:
            del active_users[user]

    return render_template("active.html", users=active_list)

def update_last_active(username):
    conn = get_db()
    cursor = conn.cursor()

    cursor.execute("""
    UPDATE users
    SET last_active = ?
    WHERE username = ?
    """, (str(datetime.now()), username))

    conn.commit()
    conn.close()

@app.route("/tasks/<int:task_id>/edit", methods=["GET", "POST"])
def edit_task(task_id):
    username = session.get("username")
    if not username or get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    conn = get_db()
    cursor = conn.cursor()

    cursor.execute("""
        SELECT t.id, t.name, t.assign_date, t.marking_script, t.question_text, t.subject_id,
               t.sample_file, t.sample_file_name, t.allow_multiple, t.max_attempts
        FROM tasks t
        WHERE t.id = ? AND t.task_type = 'practical'
    """, (task_id,))
    task = cursor.fetchone()
    if not task:
        conn.close()
        return "Task not found", 404

    subject_id = task[5]

    if request.method == "POST":
        assign_date = request.form.get("assign_date")
        marking_script = request.form.get("marking_script")
        allow_multiple = 1 if request.form.get("allow_multiple") else 0
        max_attempts = int(request.form.get("max_attempts", 1)) if allow_multiple else 1
        groups = request.form.getlist("groups")

        question_text = request.form.get("question_text", "").strip()
        cursor.execute("""
            UPDATE tasks SET assign_date = ?, marking_script = ?, question_text = ?, allow_multiple = ?, max_attempts = ? WHERE id = ?
        """, (assign_date, marking_script, question_text, allow_multiple, max_attempts, task_id))

        cursor.execute("DELETE FROM task_groups WHERE task_id = ?", (task_id,))
        for g in groups:
            if g.strip():
                cursor.execute("INSERT INTO task_groups (task_id, group_name) VALUES (?, ?)", (task_id, g))

        cursor.execute("DELETE FROM task_teachers WHERE task_id = ?", (task_id,))
        teachers = request.form.getlist("teachers")
        for t in teachers:
            if t.strip():
                cursor.execute("INSERT INTO task_teachers (task_id, teacher_username) VALUES (?, ?)", (task_id, t))

        # Optional: replace sample file
        sample_file = request.files.get("sample_file")
        if sample_file and sample_file.filename:
            sample_bytes = sample_file.read()
            # keep original filename for downloading
            sample_filename = sample_file.filename
            cursor.execute(
                """UPDATE tasks
                   SET sample_file = ?, sample_file_name = ?
                   WHERE id = ?""",
                (sample_bytes, sample_filename, task_id)
            )

        conn.commit()
        conn.close()
        log_activity(username, f"edited task {task[1]}")
        return redirect(url_for("manage_tasks", subject_id=subject_id))

    # GET — load current groups and scripts
    cursor.execute("SELECT group_name FROM task_groups WHERE task_id = ?", (task_id,))
    current_groups = {row[0] for row in cursor.fetchall()}
    cursor.execute("SELECT teacher_username FROM task_teachers WHERE task_id = ?", (task_id,))
    current_teachers = {row[0] for row in cursor.fetchall()}
    all_groups = get_groups()
    available_scripts = get_marking_scripts()
    teachers = get_teachers()
    conn.close()

    script_options = '<option value="">-- No marking script --</option>'
    for s in available_scripts:
        selected = 'selected' if s == task[3] else ''
        script_options += f'<option value="{escape(s)}" {selected}>{escape(s)}</option>'

    teacher_checkboxes = ''
    for t in teachers:
        checked = 'checked' if t[0] in current_teachers else ''
        teacher_checkboxes += f'<label style="display:inline-flex;align-items:center;gap:5px;"><input type="checkbox" name="teachers" value="{escape(t[0])}" {checked}> {escape(t[1] or t[0])}</label>'

    return render_template(
        "edit_task.html",
        task=task,
        current_groups=current_groups,
        all_groups=all_groups,
        script_options=script_options,
        teacher_checkboxes=teacher_checkboxes
    )


@app.route("/tasks/<int:task_id>/toggle", methods=["POST"])
def toggle_task(task_id):
    username = session.get("username")
    if not username or get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403
    subject_id = request.form.get("subject_id")
    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("SELECT is_active FROM tasks WHERE id = ?", (task_id,))
    row = cursor.fetchone()
    if row:
        new_state = 0 if row[0] else 1
        cursor.execute("UPDATE tasks SET is_active = ? WHERE id = ?", (new_state, task_id))
        conn.commit()
        state_label = "activated" if new_state else "deactivated"
        log_activity(username, f"{state_label} task {task_id}")
    conn.close()
    return redirect(url_for("manage_tasks", subject_id=subject_id))


@app.route("/tasks/<int:task_id>/clear_uploads", methods=["GET", "POST"])
def clear_task_uploads(task_id):

    """Serve a task's sample file stored in DB (BLOB)."""
    username = session.get("username")
    if not username:
        return "Unauthorized", 401

    conn = get_db()
    cursor = conn.cursor()
    cursor.execute(
        """
        SELECT sample_file, sample_file_name
        FROM tasks
        WHERE id = ? AND task_type = 'practical'
        """,
        (task_id,)
    )
    row = cursor.fetchone()
    conn.close()

    if not row:
        return "Sample not found", 404

    sample_bytes, sample_name = row
    if not sample_bytes:
        return "No sample file uploaded", 404

    if not sample_name:
        sample_name = f"task_{task_id}_sample"

    ext = os.path.splitext(sample_name)[1].lower()

    # Import here to avoid issues if file is missing in some environments
    from io import BytesIO

    bio = BytesIO(sample_bytes)

    # Pick a reasonable mimetype
    mimetype = None
    if ext == ".docx":
        mimetype = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    elif ext == ".xlsx":
        mimetype = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    elif ext == ".html":
        mimetype = "text/html; charset=utf-8"

    from flask import send_file
    return send_file(
        bio,
        mimetype=mimetype,
        as_attachment=True,
        download_name=sample_name
    )


def clear_task_uploads_delete(task_id):
    username = session.get("username")
    if not username or get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403
    subject_id = request.form.get("subject_id")
    conn = get_db()
    cursor = conn.cursor()
    # Get subject and task name to match results table
    cursor.execute("""
        SELECT t.name, s.name FROM tasks t
        JOIN subjects s ON s.id = t.subject_id
        WHERE t.id = ?
    """, (task_id,))
    row = cursor.fetchone()
    if row:
        task_name, subject_name = row
        cursor.execute("""
            DELETE FROM results WHERE subject = ? AND task = ?
        """, (subject_name, task_name))
        conn.commit()
        log_activity(username, f"cleared uploads for {subject_name} {task_name}")
    conn.close()
    return redirect(url_for("manage_tasks", subject_id=subject_id))



@app.route("/upload/<username>/<subject_id>/<task_id>", methods=["GET", "POST"])
def upload(username, subject_id, task_id):
    # Verify the logged-in user is the one accessing
    session_user = session.get("username")
    if session_user != username:
        return "Unauthorized", 401

    # Update active user timestamp
    update_active_user(username)

    conn = get_db()
    cursor = conn.cursor()

    # Get subject name
    cursor.execute("SELECT name FROM subjects WHERE id = ?", (subject_id,))
    subject_row = cursor.fetchone()
    if not subject_row:
        conn.close()
        return "Subject not found", 404
    subject_name = subject_row[0]

    # Get task with assign_date, marking script/setup, submission rules, and is_active
    cursor.execute(
        "SELECT name, assign_date, marking_script, question_text, allow_multiple, max_attempts, "
        "is_active, marking_setup_id FROM tasks WHERE id = ?",
        (task_id,),
    )
    task_row = cursor.fetchone()
    if not task_row:
        conn.close()
        return "Task not found", 404
    task_name, assign_date, marking_script, question_text, allow_multiple, max_attempts, task_is_active, marking_setup_id = task_row

    # Block upload if task is deactivated (students only)
    user_role = get_user_role(username)
    if user_role not in ["teacher", "admin"] and not task_is_active:
        conn.close()
        return """
        <p><a href="/student_dashboard">← Back to Dashboard</a></p>
        <h2>Upload Closed</h2>
        <p style="color:#A4262C;">This task is currently not accepting uploads. Please contact your teacher.</p>
        """, 403
    cursor.execute("SELECT group_name FROM users WHERE username = ?", (username,))
    user_group_row = cursor.fetchone()
    user_group = user_group_row[0] if user_group_row else None

    # Check authorization
    if user_role not in ["teacher", "admin"]:
        # Students can only access if task is assigned to their group AND assign date has passed
        today = datetime.now().date().isoformat()
        cursor.execute("""
        SELECT COUNT(*) FROM task_groups
        WHERE task_id = ? AND group_name = ?
        """, (task_id, user_group))
        
        if cursor.fetchone()[0] == 0:
            conn.close()
            return "Access denied: Task not assigned to your group", 403
        
        if assign_date > today:
            conn.close()
            return "Access denied: Task is not yet available", 403

    if request.method == "POST":
        cursor.execute("SELECT COUNT(*) FROM results WHERE username = ? AND subject = ? AND task = ?", (username, subject_name, task_name))
        submission_count = cursor.fetchone()[0]

        if not allow_multiple and submission_count >= 1:
            conn.close()
            return "<p><a href=\"/student_dashboard\">← Back to Dashboard</a></p><h2>Upload Closed</h2><p style=\"color:#A4262C;\">This task allows only a single submission, and you have already submitted once.</p>", 403

        if allow_multiple and submission_count >= max_attempts:
            conn.close()
            return "<p><a href=\"/student_dashboard\">← Back to Dashboard</a></p><h2>Upload Closed</h2><p style=\"color:#A4262C;\">You have reached the maximum number of submissions for this task.</p>", 403

        file = request.files.get("file")

        if not file:
            conn.close()
            return "No file uploaded", 400

        original_ext = os.path.splitext(file.filename or "")[1] or ".tmp"
        temp_path = f"temp_{username}{original_ext}"
        file.save(temp_path)

        try:
            result = mark_file(temp_path, marking_script, marking_setup_id)
        finally:
            if os.path.exists(temp_path):
                os.remove(temp_path)

        if result["error"]:
            conn.close()
            return f"""
            <p><a href="/student_dashboard">← Back to Dashboard</a></p>
            <h2>Submission Error</h2>
            <p style="color:red;">{escape(result['error'])}</p>
            """

        # Build weak skills from wrong answers
        weak_skills = [r["question"] for r in result["results"] if not r["passed"]]
        update_weakness(username, weak_skills)
        save_result(username, subject_name, task_name, result["percentage"], ", ".join(weak_skills[:3]) or "Well done!")
        log_activity(username, f"submitted {subject_name} {task_name}")

        correct_items = [r for r in result["results"] if r["passed"]]
        wrong_items = [r for r in result["results"] if not r["passed"]]

        conn.close()
        return render_template(
            "upload_result.html",
            subject_name=subject_name,
            task_name=task_name,
            score=result["score"],
            total=result["total"],
            percentage=result["percentage"],
            correct_items=correct_items,
            wrong_items=wrong_items,
        )

    conn.close()

    # Include a student-safe sample download link if a sample file exists
    sample_link_html = ""
    cursor = None
    try:
        conn = get_db()
        cursor = conn.cursor()
        cursor.execute(
            "SELECT sample_file, sample_file_name FROM tasks WHERE id = ? AND task_type = 'practical'",
            (task_id,),
        )
        row = cursor.fetchone()
        if row and row[0]:
            sample_name = row[1] or f"task_{task_id}_sample"
            sample_link_html = (
                f'<p><a href="/tasks/{task_id}/sample_file" target="_blank">📎 Download task sample</a> '
                f'({escape(sample_name)})</p>'
            )
    finally:
        try:
            if cursor:
                conn.close()
        except Exception:
            pass

    return render_template(
        "upload_task.html",
        subject_name=subject_name,
        task_name=task_name,
        question_text=question_text,
        sample_link_html=sample_link_html,
    )




@app.route("/subjects/<username>")
def subjects(username):
    if not session.get("username"):
        return redirect(url_for("login"))

    session_user = session.get("username")
    if session_user != username:
        return "Unauthorized", 401
    return redirect(url_for("learner_tasks"))

@app.route("/tasks/<username>/<subject_id>")
def tasks(username, subject_id):
    if not session.get("username"):
        return redirect(url_for("login"))

    # Verify the logged-in user is the one accessing
    session_user = session.get("username")
    if session_user != username:
        return "Unauthorized", 401
    return redirect(url_for("learner_tasks"))

@app.route("/view_as_student/<group_name>")
def view_as_student(group_name):
    admin_user = session.get("username")
    if not admin_user:
        return redirect(url_for("login"))
    if get_user_role(admin_user) not in ["teacher", "admin"]:
        return "Access denied", 403

    conn = get_db()
    cursor = conn.cursor()

    # Pick first student in the group as representative
    cursor.execute("""
        SELECT username, full_name, group_name FROM users
        WHERE group_name = ? AND role = 'student' LIMIT 1
    """, (group_name,))
    rep = cursor.fetchone()

    # Group student count
    cursor.execute("SELECT COUNT(*) FROM users WHERE group_name = ? AND role = 'student'", (group_name,))
    student_count = cursor.fetchone()[0]

    # All students in group
    cursor.execute("""
        SELECT username, full_name FROM users
        WHERE group_name = ? AND role = 'student' ORDER BY full_name
    """, (group_name,))
    students = cursor.fetchall()

    # Subject averages for the group
    cursor.execute("""
        SELECT r.subject, ROUND(AVG(r.score), 1)
        FROM results r
        JOIN users u ON r.username = u.username
        WHERE u.group_name = ?
        GROUP BY r.subject
    """, (group_name,))
    subject_avgs = cursor.fetchall()

    # Overall group average
    practical_group_avg = fetch_group_practical_averages(cursor, [group_name]).get(group_name)
    theory_group_avg = fetch_group_theory_averages(cursor, [group_name]).get(group_name)
    if practical_group_avg is not None and theory_group_avg is not None:
        overall_avg = round((practical_group_avg + theory_group_avg) / 2, 1)
    else:
        overall_avg = practical_group_avg if practical_group_avg is not None else (theory_group_avg or 0)

    # Recent submissions for the group
    cursor.execute("""
        SELECT u.full_name, r.subject, r.task, r.score, r.timestamp
        FROM results r
        JOIN users u ON r.username = u.username
        WHERE u.group_name = ?
        ORDER BY r.timestamp DESC LIMIT 10
    """, (group_name,))
    recent_results = cursor.fetchall()

    # Attendance last 7 days for the group
    days = get_last_7_days()
    excluded = fetch_group_excluded_dates(cursor, [group_name]).get(group_name, set())
    days = [d for d in days if d not in excluded]
    student_usernames = [student_username for student_username, _full_name in students]
    login_times = fetch_first_login_times(cursor, student_usernames, days)
    overrides = fetch_attendance_override_statuses(cursor, student_usernames, days)

    att_summary = []
    for day in days:
        present = 0
        absent = 0
        for student_username in student_usernames:
            history = build_attendance_history(
                cursor,
                student_username,
                group_name,
                [day],
                login_map=login_times,
                override_map=overrides,
                excluded_dates=excluded,
            )
            summary = summarize_attendance_history(history)
            present += summary["present_days"]
            absent += summary["absent_days"]
        total_counted = present + absent
        pct = round((present / total_counted) * 100) if total_counted else 0
        att_summary.append({"date": day, "present": present, "total": total_counted, "pct": pct})

    # Missing tasks for the group
    today = datetime.now().date().isoformat()
    cursor.execute("""
        SELECT s.name, t.name, t.assign_date,
               COUNT(DISTINCT u.username) as missing_count
        FROM tasks t
        JOIN subjects s ON s.id = t.subject_id
        JOIN task_groups tg ON tg.task_id = t.id
        JOIN users u ON u.group_name = tg.group_name AND u.role = 'student'
        WHERE tg.group_name = ? AND t.assign_date <= ? AND t.task_type = 'practical'
          AND NOT EXISTS (
              SELECT 1 FROM results r
              WHERE r.username = u.username AND r.subject = s.name AND r.task = t.name
          )
        GROUP BY t.id
        ORDER BY t.assign_date
    """, (group_name, today))
    missing_tasks = cursor.fetchall()

    # Top weaknesses for the group
    cursor.execute("""
        SELECT w.skill, SUM(w.count) as total
        FROM weaknesses w
        JOIN users u ON w.username = u.username
        WHERE u.group_name = ?
        GROUP BY w.skill ORDER BY total DESC LIMIT 5
    """, (group_name,))
    weaknesses = cursor.fetchall()

    # Theory test results for the group
    cursor.execute("""
        SELECT t.title, ROUND(AVG(s.percentage), 1), COUNT(DISTINCT s.username)
        FROM theory_submissions s
        JOIN theory_tests t ON s.test_id = t.id
        JOIN users u ON s.username = u.username
        WHERE u.group_name = ?
        GROUP BY t.id ORDER BY s.submitted_at DESC LIMIT 5
    """, (group_name,))
    theory_avgs = cursor.fetchall()

    conn.close()
    return render_template(
        "view_as_student.html",
        group_name=group_name,
        student_count=student_count,
        students=students,
        subject_avgs=subject_avgs,
        overall_avg=overall_avg,
        recent_results=recent_results,
        att_summary=att_summary,
        missing_tasks=missing_tasks,
        weaknesses=weaknesses,
        theory_avgs=theory_avgs
    )


@app.route("/teacher_dashboard")
def teacher_dashboard():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    role = get_user_role(username)
    if role not in ["teacher", "admin"]:
        return "Access denied", 403

    conn = get_db()
    cursor = conn.cursor()
    today = datetime.now().strftime("%Y-%m-%d")
    groups = get_groups(username)
    quick_actions, available_quick_actions = get_teacher_selected_quick_actions(username)

    if role == "teacher":
        cursor.execute("""
            SELECT username, full_name, group_name
            FROM users
            WHERE role = 'student' AND teacher_username = ?
            ORDER BY group_name, full_name, username
        """, (username,))
    else:
        cursor.execute("""
            SELECT username, full_name, group_name
            FROM users
            WHERE role = 'student'
            ORDER BY group_name, full_name, username
        """)
    student_rows = cursor.fetchall()
    total_students = len(student_rows)
    student_usernames = [row[0] for row in student_rows]
    student_groups = sorted({row[2] for row in student_rows if row[2]})
    groups = student_groups or groups
    group_filter_clause = "AND u.teacher_username = ?" if role == "teacher" else ""
    group_filter_params = [username] if role == "teacher" else []

    cursor.execute("""
        SELECT COUNT(DISTINCT lh.username)
        FROM login_history lh
        JOIN users u ON u.username = lh.username
        WHERE lh.date = ? AND u.role = 'student'
        """ + (" AND u.teacher_username = ?" if role == "teacher" else ""),
        ((today, username) if role == "teacher" else (today,))
    )
    active_today = cursor.fetchone()[0] or 0

    days_21 = get_last_21_days()
    excluded_days = fetch_group_excluded_dates(cursor, groups)
    filtered_days = [day for day in days_21 if day not in excluded_days]
    login_times = fetch_first_login_times(cursor, student_usernames, filtered_days)
    overrides = fetch_attendance_override_statuses(cursor, student_usernames, filtered_days)

    group_members = defaultdict(list)
    for uname, full_name, group_name in student_rows:
        if group_name:
            group_members[group_name].append((uname, full_name or uname))

    group_att = []
    total_present_all = 0
    total_slots_all = 0
    for group_name in groups:
        members = group_members.get(group_name, [])
        group_days = [day for day in filtered_days if day not in excluded_days.get(group_name, set())]
        total_present = 0
        total_absent = 0
        for student_username, _student_name in members:
            history = build_attendance_history(
                cursor,
                student_username,
                group_name,
                group_days,
                login_map=login_times,
                override_map=overrides,
                excluded_dates=excluded_days.get(group_name, set()),
            )
            summary = summarize_attendance_history(history)
            total_present += summary["present_days"]
            total_absent += summary["absent_days"]
        total_counted = total_present + total_absent
        total_present_all += total_present
        total_slots_all += total_counted
        group_att.append({
            "group": group_name,
            "students": len(members),
            "att_pct": round((total_present / total_counted) * 100) if total_counted else 0
        })
    avg_att_pct = round((total_present_all / total_slots_all) * 100) if total_slots_all else 0

    low_attendance = get_low_attendance_learners(10, groups, username if role == "teacher" else None)

    practical_avg_map = fetch_student_practical_averages(cursor, username if role == "teacher" else None)
    theory_avg_map = fetch_student_theory_averages(cursor, username if role == "teacher" else None)

    practical_group_avg_map = fetch_group_practical_averages(cursor, groups)
    theory_group_avg_map = fetch_group_theory_averages(cursor, groups)

    missing_practical_map = {}
    missing_theory_map = {}
    if student_usernames:
        placeholders = ",".join("?" for _ in student_usernames)
        cursor.execute(f"""
            SELECT u.username, COUNT(*)
            FROM users u
            JOIN task_groups tg ON tg.group_name = u.group_name
            JOIN tasks t ON t.id = tg.task_id
            JOIN subjects s ON s.id = t.subject_id
            WHERE u.username IN ({placeholders})
              AND t.assign_date <= ?
              AND t.task_type = 'practical'
              AND t.is_active = 1
              AND NOT EXISTS (
                  SELECT 1 FROM results r
                  WHERE r.username = u.username AND r.subject = s.name AND r.task = t.name
              )
            GROUP BY u.username
        """, (*student_usernames, today))
        missing_practical_map = {row[0]: row[1] for row in cursor.fetchall()}

        cursor.execute(f"""
            SELECT u.username, COUNT(DISTINCT tt.id)
            FROM users u
            JOIN theory_tests tt ON tt.is_active = 1 AND tt.assign_date <= ?
            LEFT JOIN theory_test_groups ttg ON tt.id = ttg.test_id
            WHERE u.username IN ({placeholders})
              AND (ttg.group_name = u.group_name OR ttg.group_name IS NULL)
              AND NOT EXISTS (
                  SELECT 1 FROM theory_submissions ts
                  WHERE ts.username = u.username AND ts.test_id = tt.id
              )
            GROUP BY u.username
        """, (today, *student_usernames))
        missing_theory_map = {row[0]: row[1] for row in cursor.fetchall()}

    at_risk_students = []
    combined_students = []
    for uname, full_name, group_name in student_rows:
        group_days = [day for day in filtered_days if day not in excluded_days.get(group_name, set())]
        attendance_history = build_attendance_history(
            cursor,
            uname,
            group_name,
            group_days,
            login_map=login_times,
            override_map=overrides,
            excluded_dates=excluded_days.get(group_name, set()),
        )
        attendance_pct = summarize_attendance_history(attendance_history)["attendance_pct"]

        practical_avg = practical_avg_map.get(uname)
        theory_avg = theory_avg_map.get(uname)
        if practical_avg is not None and theory_avg is not None:
            combined_avg = round((practical_avg + theory_avg) / 2, 1)
        else:
            combined_avg = practical_avg if practical_avg is not None else theory_avg
        combined_avg = combined_avg if combined_avg is not None else 0

        missing_total = (missing_practical_map.get(uname, 0) or 0) + (missing_theory_map.get(uname, 0) or 0)
        assigned_total = missing_total
        if practical_avg is not None or theory_avg is not None:
            assigned_total = max(assigned_total, 1)
        missing_pct = round((missing_total / assigned_total) * 100) if assigned_total else 0

        risk_score = round((100 - attendance_pct) * 0.4 + (100 - combined_avg) * 0.4 + missing_pct * 0.2)
        reasons = []
        if attendance_pct < 60:
            reasons.append(f"Attendance {attendance_pct}%")
        if combined_avg < 40:
            reasons.append(f"Average {combined_avg}%")
        if missing_total > 0:
            reasons.append(f"Missing {missing_total}")
        if risk_score > 40:
            at_risk_students.append({
                "username": uname,
                "name": full_name or uname,
                "group": group_name or "—",
                "score": risk_score,
                "status": "High Risk" if risk_score > 70 else "At Risk",
                "reason": " | ".join(reasons) if reasons else "Multiple factors"
            })
        combined_students.append((full_name or uname, group_name or "—", combined_avg))

    at_risk_students.sort(key=lambda item: item["score"], reverse=True)
    at_risk_students = at_risk_students[:10]
    combined_students = [row for row in combined_students if row[2] is not None]
    top_performers = sorted(combined_students, key=lambda item: item[2], reverse=True)[:5]
    bottom_performers = sorted(combined_students, key=lambda item: item[2])[:5]

    # ─ Recent activity (last 20)
    if role == 'teacher':
        cursor.execute("""
            SELECT a.username, a.action, a.timestamp
            FROM activities a
            LEFT JOIN users u ON u.username = a.username
            WHERE a.username = ?
              OR (u.role = 'student' AND u.teacher_username = ?)
            ORDER BY a.timestamp DESC LIMIT 20
        """, (username, username))
    else:
        cursor.execute("""
            SELECT username, action, timestamp FROM activities
            ORDER BY timestamp DESC LIMIT 20
        """)
    recent_activities = cursor.fetchall()

    # ─ Recent submissions (last 15, best score per student per task)
    if role == 'teacher':
        cursor.execute(f"""
            SELECT u.full_name, u.group_name, b.subject, b.task, b.best_score, MAX(r.timestamp)
            FROM (
                SELECT username, subject, task, MAX(score) as best_score
                FROM results GROUP BY username, subject, task
            ) b
            JOIN results r ON r.username = b.username AND r.subject = b.subject AND r.task = b.task AND r.score = b.best_score
            JOIN users u ON u.username = b.username
            WHERE u.role = 'student' {group_filter_clause}
            GROUP BY b.username, b.subject, b.task
            ORDER BY MAX(r.timestamp) DESC LIMIT 15
        """, group_filter_params)
    else:
        cursor.execute("""
            SELECT u.full_name, u.group_name, b.subject, b.task, b.best_score, MAX(r.timestamp)
            FROM (
                SELECT username, subject, task, MAX(score) as best_score
                FROM results GROUP BY username, subject, task
            ) b
            JOIN results r ON r.username = b.username AND r.subject = b.subject AND r.task = b.task AND r.score = b.best_score
            JOIN users u ON u.username = b.username
            GROUP BY b.username, b.subject, b.task
            ORDER BY MAX(r.timestamp) DESC LIMIT 15
        """)
    recent_submissions = cursor.fetchall()

    # ─ Class averages per subject per group (best score per student per task)
    # Build practical subjects per group from assigned tasks so new subjects appear even without results.
    practical_subjects = defaultdict(set)
    if role == 'teacher' and groups:
        placeholders = ",".join("?" for _ in groups)
        cursor.execute(f"""
            SELECT tg.group_name, s.name
            FROM task_groups tg
            JOIN tasks t ON t.id = tg.task_id
            JOIN subjects s ON s.id = t.subject_id
            WHERE t.task_type = 'practical' AND tg.group_name IN ({placeholders})
        """, groups)
    else:
        cursor.execute("""
            SELECT tg.group_name, s.name
            FROM task_groups tg
            JOIN tasks t ON t.id = tg.task_id
            JOIN subjects s ON s.id = t.subject_id
            WHERE t.task_type = 'practical'
        """)
    for group_name, subject in cursor.fetchall():
        practical_subjects[group_name].add(subject)

    if role == 'teacher':
        cursor.execute(f"""
            SELECT u.group_name, b.subject, ROUND(AVG(b.best_score),1), COUNT(*)
            FROM (
                SELECT username, subject, task, MAX(score) as best_score
                FROM results GROUP BY username, subject, task
            ) b
            JOIN users u ON u.username = b.username
            WHERE u.group_name IS NOT NULL AND u.role = 'student' {group_filter_clause}
            GROUP BY u.group_name, b.subject
            ORDER BY u.group_name, b.subject
        """, group_filter_params)
    else:
        cursor.execute("""
            SELECT u.group_name, b.subject, ROUND(AVG(b.best_score),1), COUNT(*)
            FROM (
                SELECT username, subject, task, MAX(score) as best_score
                FROM results GROUP BY username, subject, task
            ) b
            JOIN users u ON u.username = b.username
            WHERE u.group_name IS NOT NULL AND u.role = 'student'
            GROUP BY u.group_name, b.subject
            ORDER BY u.group_name, b.subject
        """)
    practical_avgs_raw = cursor.fetchall()

    practical_avgs = { (g, subj): (avg, cnt) for g, subj, avg, cnt in practical_avgs_raw }

    subject_avgs = defaultdict(list)
    for group_name, subjects in practical_subjects.items():
        for subject in sorted(subjects):
            avg, cnt = practical_avgs.get((group_name, subject), (None, 0))
            subject_avgs[group_name].append((subject, avg, cnt, 'Practical'))

    # Add any practical subjects that appear only in results but aren't assigned through tasks
    for group_name, subject, avg, cnt in practical_avgs_raw:
        if subject not in practical_subjects.get(group_name, set()):
            subject_avgs[group_name].append((subject, avg, cnt, 'Practical'))

    # ─ Combined theory averages per group
    if role == 'teacher':
        cursor.execute(f"""
            SELECT u.group_name, ROUND(AVG(b.best_pct),1), COUNT(*)
            FROM (
                SELECT username, test_id, MAX(percentage) as best_pct
                FROM theory_submissions GROUP BY username, test_id
            ) b
            JOIN theory_tests tt ON b.test_id = tt.id
            JOIN users u ON u.username = b.username
            WHERE u.group_name IS NOT NULL AND u.role = 'student' {group_filter_clause}
            GROUP BY u.group_name
            ORDER BY u.group_name
        """, group_filter_params)
    else:
        cursor.execute("""
            SELECT u.group_name, ROUND(AVG(b.best_pct),1), COUNT(*)
            FROM (
                SELECT username, test_id, MAX(percentage) as best_pct
                FROM theory_submissions GROUP BY username, test_id
            ) b
            JOIN theory_tests tt ON b.test_id = tt.id
            JOIN users u ON u.username = b.username
            WHERE u.group_name IS NOT NULL AND u.role = 'student'
            GROUP BY u.group_name
            ORDER BY u.group_name
        """)
    theory_avgs_raw = cursor.fetchall()

    theory_avgs = {group_name: (avg, cnt) for group_name, avg, cnt in theory_avgs_raw}

    # Add Theory row for groups with theory assignments or results
    theory_groups = set()
    if role == 'teacher' and groups:
        placeholders = ",".join("?" for _ in groups)
        cursor.execute(f"""
            SELECT DISTINCT ttg.group_name
            FROM theory_tests tt
            JOIN theory_test_groups ttg ON tt.id = ttg.test_id
            WHERE tt.assign_date IS NOT NULL AND ttg.group_name IN ({placeholders})
        """, groups)
        theory_groups.update(row[0] for row in cursor.fetchall() if row[0])

        cursor.execute("""
            SELECT 1
            FROM theory_tests tt
            LEFT JOIN theory_test_groups ttg ON tt.id = ttg.test_id
            WHERE tt.assign_date IS NOT NULL AND ttg.group_name IS NULL
            LIMIT 1
        """)
        if cursor.fetchone():
            theory_groups.update(groups)
    else:
        cursor.execute("""
            SELECT DISTINCT ttg.group_name
            FROM theory_tests tt
            JOIN theory_test_groups ttg ON tt.id = ttg.test_id
            WHERE tt.assign_date IS NOT NULL AND ttg.group_name IS NOT NULL
        """)
        theory_groups.update(row[0] for row in cursor.fetchall() if row[0])

        cursor.execute("""
            SELECT 1
            FROM theory_tests tt
            LEFT JOIN theory_test_groups ttg ON tt.id = ttg.test_id
            WHERE tt.assign_date IS NOT NULL AND ttg.group_name IS NULL
            LIMIT 1
        """)
        if cursor.fetchone():
            cursor.execute("SELECT DISTINCT group_name FROM users WHERE group_name IS NOT NULL")
            theory_groups.update(row[0] for row in cursor.fetchall() if row[0])

    for group_name in sorted(theory_groups):
        avg, cnt = theory_avgs.get(group_name, (None, 0))
        subject_avgs[group_name].append(("Theory", avg, cnt, 'Theory'))

    subject_avgs = dict(subject_avgs)

    # ─ Recent theory submissions (best per student per test)
    if role == 'teacher':
        cursor.execute(f"""
            SELECT u.full_name, u.group_name, tt.subject, tt.title, b.best_pct, MAX(ts.submitted_at)
            FROM (
                SELECT username, test_id, MAX(percentage) as best_pct
                FROM theory_submissions GROUP BY username, test_id
            ) b
            JOIN theory_submissions ts ON ts.username = b.username AND ts.test_id = b.test_id AND ts.percentage = b.best_pct
            JOIN theory_tests tt ON b.test_id = tt.id
            JOIN users u ON u.username = b.username
            WHERE u.role = 'student' {group_filter_clause}
            GROUP BY b.username, b.test_id
            ORDER BY MAX(ts.submitted_at) DESC LIMIT 15
        """, group_filter_params)
    else:
        cursor.execute("""
            SELECT u.full_name, u.group_name, tt.subject, tt.title, b.best_pct, MAX(ts.submitted_at)
            FROM (
                SELECT username, test_id, MAX(percentage) as best_pct
                FROM theory_submissions GROUP BY username, test_id
            ) b
            JOIN theory_submissions ts ON ts.username = b.username AND ts.test_id = b.test_id AND ts.percentage = b.best_pct
            JOIN theory_tests tt ON b.test_id = tt.id
            JOIN users u ON u.username = b.username
            GROUP BY b.username, b.test_id
            ORDER BY MAX(ts.submitted_at) DESC LIMIT 15
        """)
    recent_theory_submissions = cursor.fetchall()

    # ─ Students without assigned classes or teachers
    cursor.execute("""
        SELECT username, full_name
        FROM users
        WHERE role = 'student' AND (
            (group_name IS NULL OR group_name = '') OR
            (teacher_username IS NULL OR teacher_username = '')
        )
        ORDER BY full_name, username
    """)
    students_without_classes = cursor.fetchall()

    # ─ Missing tasks count per group
    if role == 'teacher':
        cursor.execute(f"""
            SELECT u.group_name, COUNT(*) as missing
            FROM users u
            JOIN task_groups tg ON tg.group_name = u.group_name
            JOIN tasks t ON t.id = tg.task_id
            JOIN subjects s ON s.id = t.subject_id
            WHERE u.role = 'student'
              AND u.group_name IS NOT NULL
              AND t.assign_date <= ?
              AND t.task_type = 'practical'
              AND NOT EXISTS (
                  SELECT 1 FROM results r
                  WHERE r.username = u.username AND r.subject = s.name AND r.task = t.name
              )
              {group_filter_clause}
            GROUP BY u.group_name
        """, (today, *group_filter_params))
    else:
        cursor.execute("""
            SELECT u.group_name, COUNT(*) as missing
            FROM users u
            JOIN task_groups tg ON tg.group_name = u.group_name
            JOIN tasks t ON t.id = tg.task_id
            JOIN subjects s ON s.id = t.subject_id
            WHERE u.role = 'student'
              AND t.assign_date <= ?
              AND t.task_type = 'practical'
              AND NOT EXISTS (
                  SELECT 1 FROM results r
                  WHERE r.username = u.username AND r.subject = s.name AND r.task = t.name
              )
            GROUP BY u.group_name
        """, (today,))
    missing_by_group = cursor.fetchall()

    group_att_summary = []
    for item in group_att:
        practical_group_avg = practical_group_avg_map.get(item["group"])
        theory_group_avg = theory_group_avg_map.get(item["group"])
        if practical_group_avg is not None and theory_group_avg is not None:
            combined_group_avg = round((practical_group_avg + theory_group_avg) / 2, 1)
        else:
            combined_group_avg = practical_group_avg if practical_group_avg is not None else theory_group_avg
        group_att_summary.append({
            **item,
            "practical_pct": practical_group_avg,
            "theory_pct": theory_group_avg,
            "combined_pct": combined_group_avg
        })

    conn.close()
    return render_template(
        "teacher_dashboard.html",
        username=username,
        total_students=total_students,
        active_today=active_today,
        avg_att_pct=avg_att_pct,
        groups=groups,
        group_att=group_att_summary,
        low_attendance=low_attendance,
        recent_activities=recent_activities,
        recent_submissions=recent_submissions,
        recent_theory_submissions=recent_theory_submissions,
        subject_avgs=subject_avgs,
        top_performers=top_performers,
        bottom_performers=bottom_performers,
        at_risk_students=at_risk_students,
        students_without_classes=students_without_classes,
        missing_by_group=missing_by_group,
        quick_actions=quick_actions,
        available_quick_actions=available_quick_actions,
        active_term=get_active_term_range(),
        days_in_period=len(days_21)
    )


@app.route("/dashboard")
def legacy_dashboard():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    role = get_user_role(username)
    if role in ["teacher", "admin"]:
        return redirect(url_for("teacher_dashboard"))
    return redirect(url_for("student_dashboard"))


@app.route("/manage_tests_simple")
def legacy_manage_tests_simple():
    return redirect(url_for("manage_tests"))


@app.route("/manage_test_questions_simple/<int:test_id>")
def legacy_manage_test_questions_simple(test_id):
    return redirect(url_for("manage_test_questions", test_id=test_id))


@app.route("/take_test_simple/<int:test_id>", methods=["GET", "POST"])
def legacy_take_test_simple(test_id):
    return redirect(url_for("take_test", test_id=test_id), code=307 if request.method == "POST" else 302)


@app.route("/response_review")
def response_review():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    role = get_user_role(username)
    if role not in ["teacher", "admin"]:
        return "Access denied", 403

    review_type = (request.args.get("type") or "theory").strip()
    selected_group = (request.args.get("group") or "").strip()
    selected_item = (request.args.get("item") or "").strip()
    groups = get_groups(username)

    conn = get_db()
    conn.row_factory = sqlite3.Row
    cursor = conn.cursor()

    available_items = []
    if review_type == "theory":
        query = """
            SELECT DISTINCT tt.id, COALESCE(tt.subject || ' - ', '') || tt.title
            FROM theory_submissions s
            JOIN theory_tests tt ON tt.id = s.test_id
            JOIN users u ON u.username = s.username
            WHERE u.role = 'student'
        """
        params = []
        if role == "teacher":
            query += " AND u.teacher_username = ?"
            params.append(username)
        if selected_group:
            query += " AND u.group_name = ?"
            params.append(selected_group)
        query += " ORDER BY tt.title"
        cursor.execute(query, params)
        available_items = [(str(row[0]), row[1]) for row in cursor.fetchall()]

    if not selected_item and available_items:
        selected_item = available_items[0][0]

    question_summaries = []
    item_label = ""
    selected_test_id = int(selected_item) if selected_item and selected_item.isdigit() else None

    if review_type == "theory" and selected_test_id:
        cursor.execute("SELECT title, subject FROM theory_tests WHERE id = ?", (selected_test_id,))
        test_row = cursor.fetchone()
        if test_row:
            item_label = f"{test_row['subject']} - {test_row['title']}" if test_row["subject"] else test_row["title"]

        submission_query = """
            SELECT s.id, s.username
            FROM theory_submissions s
            JOIN (
                SELECT username, MAX(submitted_at) AS latest_submitted_at
                FROM theory_submissions
                WHERE test_id = ?
                GROUP BY username
            ) latest ON latest.username = s.username AND latest.latest_submitted_at = s.submitted_at
            JOIN users u ON u.username = s.username
            WHERE s.test_id = ?
        """
        submission_params = [selected_test_id, selected_test_id]
        if role == "teacher":
            submission_query += " AND u.teacher_username = ?"
            submission_params.append(username)
        if selected_group:
            submission_query += " AND u.group_name = ?"
            submission_params.append(selected_group)
        cursor.execute(submission_query, submission_params)
        latest_submissions = cursor.fetchall()
        submission_ids = [row["id"] for row in latest_submissions]

        cursor.execute("""
            SELECT id, question_text, question_type, marks
            FROM theory_questions
            WHERE test_id = ?
            ORDER BY order_index
        """, (selected_test_id,))
        questions = cursor.fetchall()

        answers_by_question = defaultdict(list)
        if submission_ids:
            placeholders = ",".join("?" for _ in submission_ids)
            cursor.execute(f"""
                SELECT question_id, answer_text, is_correct, marks_awarded
                FROM theory_answers
                WHERE submission_id IN ({placeholders})
            """, submission_ids)
            for row in cursor.fetchall():
                answers_by_question[row["question_id"]].append(dict(row))

        for question in questions:
            q_id = question["id"]
            q_type = question["question_type"]
            marks = question["marks"] or 0
            cursor.execute("""
                SELECT id, option_text, is_correct, match_pair
                FROM theory_options
                WHERE question_id = ?
                ORDER BY id
            """, (q_id,))
            options = cursor.fetchall()
            answers = answers_by_question.get(q_id, [])
            accepted_answers = []
            option_rows = []
            correct_count = sum(1 for answer in answers if answer["is_correct"])

            if q_type == "true_false":
                correct_choice, accepted_corrections = get_true_false_option_data(options)
                accepted_answers = ([correct_choice] if correct_choice else []) + accepted_corrections
                buckets = defaultdict(int)
                for answer in answers:
                    selected, correction = parse_true_false_answer_text(answer["answer_text"])
                    label = f"{selected} | {correction}" if correction else (selected or "No answer")
                    buckets[label] += 1
                total_answers = len(answers)
                for label, count in sorted(buckets.items(), key=lambda item: (-item[1], item[0])):
                    correction_value = label.split(" | ", 1)[1] if " | " in label else ""
                    option_rows.append({
                        "label": label,
                        "count": count,
                        "pct": round((count / total_answers) * 100) if total_answers else 0,
                        "is_correct": label == correct_choice,
                        "acceptable_answer": correction_value if correction_value and normalize_review_text(correction_value) not in {normalize_review_text(item) for item in accepted_corrections} else "",
                        "remove_answer": correction_value if correction_value and normalize_review_text(correction_value) in {normalize_review_text(item) for item in accepted_corrections} else "",
                    })
            elif q_type == "fill_in":
                accepted_answers = get_fill_in_accepted_answers(options)
                buckets = defaultdict(int)
                for answer in answers:
                    buckets[answer["answer_text"] or "No answer"] += 1
                total_answers = len(answers)
                accepted_keys = {normalize_review_text(item) for item in accepted_answers}
                for label, count in sorted(buckets.items(), key=lambda item: (-item[1], item[0])):
                    normalized = normalize_review_text(label)
                    option_rows.append({
                        "label": label,
                        "count": count,
                        "pct": round((count / total_answers) * 100) if total_answers else 0,
                        "is_correct": normalized in accepted_keys,
                        "acceptable_answer": label if label and label != "No answer" and normalized not in accepted_keys else "",
                        "remove_answer": label if label and normalized in accepted_keys else "",
                    })
            elif q_type == "match":
                option_rows = []
                for option in options:
                    left = option["option_text"]
                    accepted = option["match_pair"]
                    responses = []
                    for answer in answers:
                        pieces = [part.strip() for part in (answer["answer_text"] or "").split(";") if "=" in part]
                        selected_match = ""
                        for piece in pieces:
                            col_a, submitted = piece.split("=", 1)
                            if col_a.strip() == left:
                                selected_match = submitted.strip()
                                break
                        if selected_match:
                            responses.append(selected_match)
                    option_rows.append({"left": left, "accepted": accepted, "responses": responses})
            else:
                correct_values = {option["option_text"] for option in options if option["is_correct"] == 1}
                accepted_answers = list(correct_values)
                buckets = defaultdict(int)
                for answer in answers:
                    buckets[answer["answer_text"] or "No answer"] += 1
                total_answers = len(answers)
                for label, count in sorted(buckets.items(), key=lambda item: (-item[1], item[0])):
                    option_rows.append({
                        "label": label,
                        "count": count,
                        "pct": round((count / total_answers) * 100) if total_answers else 0,
                        "is_correct": label in correct_values,
                        "acceptable_answer": "",
                        "remove_answer": "",
                    })

            question_summaries.append({
                "question_id": q_id,
                "question": question["question_text"],
                "type": q_type,
                "marks": marks,
                "responses": len(answers),
                "correct_pct": round((correct_count / len(answers)) * 100) if answers else 0,
                "accepted_answers": accepted_answers,
                "options": option_rows,
            })

    conn.close()
    return render_template(
        "response_review.html",
        review_type=review_type,
        groups=groups,
        selected_group=selected_group,
        selected_item=selected_item,
        available_items=available_items,
        item_label=item_label,
        question_summaries=question_summaries,
        selected_test_id=selected_test_id,
        can_manage_reviews=role in ["teacher", "admin"]
    )


@app.route("/response_review/learner")
def response_review_learner():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    role = get_user_role(username)
    if role not in ["teacher", "admin"]:
        return "Access denied", 403

    learner_username = (request.args.get("learner") or "").strip()
    test_id = request.args.get("item", type=int)
    if not learner_username or not test_id:
        return redirect(url_for("response_review"))

    conn = get_db()
    conn.row_factory = sqlite3.Row
    cursor = conn.cursor()
    cursor.execute("""
        SELECT full_name, group_name
        FROM users
        WHERE username = ?
    """, (learner_username,))
    learner_row = cursor.fetchone()
    if not learner_row:
        conn.close()
        return "Learner not found", 404

    cursor.execute("""
        SELECT s.id, s.score, s.total, s.percentage, s.submitted_at, tt.title, tt.subject
        FROM theory_submissions s
        JOIN theory_tests tt ON tt.id = s.test_id
        WHERE s.username = ? AND s.test_id = ?
        ORDER BY s.submitted_at DESC
        LIMIT 1
    """, (learner_username, test_id))
    submission = cursor.fetchone()
    if not submission:
        conn.close()
        return render_template(
            "response_review_learner.html",
            learner_name=learner_row["full_name"] or learner_username,
            learner_username=learner_username,
            learner_group=learner_row["group_name"] or "—",
            item_label="No submission found",
            review_rows=[],
            score_summary=None,
            can_manage_reviews=True
        )

    cursor.execute("""
        SELECT q.question_text, q.question_type, q.marks, a.answer_text, a.is_correct, a.marks_awarded, a.question_id
        FROM theory_answers a
        JOIN theory_questions q ON q.id = a.question_id
        WHERE a.submission_id = ?
        ORDER BY q.order_index
    """, (submission["id"],))
    answer_rows = cursor.fetchall()

    review_rows = []
    for row in answer_rows:
        cursor.execute("""
            SELECT id, option_text, is_correct, match_pair
            FROM theory_options
            WHERE question_id = ?
            ORDER BY id
        """, (row["question_id"],))
        options = cursor.fetchall()
        correct_answers = []
        match_rows = []
        if row["question_type"] == "true_false":
            correct_choice, accepted_corrections = get_true_false_option_data(options)
            correct_answers = ([correct_choice] if correct_choice else []) + accepted_corrections
        elif row["question_type"] == "fill_in":
            correct_answers = get_fill_in_accepted_answers(options)
        elif row["question_type"] == "match":
            expected_map = {option["option_text"]: option["match_pair"] for option in options if option["match_pair"] and option["match_pair"] != "correction"}
            learner_map = {}
            for pair in (row["answer_text"] or "").split(";"):
                if "=" in pair:
                    left, chosen = pair.split("=", 1)
                    learner_map[left.strip()] = chosen.strip()
            for left, accepted in expected_map.items():
                chosen = learner_map.get(left, "")
                match_rows.append({
                    "left": left,
                    "learner_match": chosen or "No answer",
                    "correct_match": accepted,
                    "is_correct": chosen == accepted
                })
        else:
            correct_answers = [option["option_text"] for option in options if option["is_correct"] == 1]

        review_rows.append({
            "question_id": row["question_id"],
            "question": row["question_text"],
            "type": row["question_type"],
            "marks": row["marks"],
            "answer": row["answer_text"] or "No answer",
            "correct": row["is_correct"],
            "marks_awarded": row["marks_awarded"],
            "correct_answers": correct_answers,
            "match_rows": match_rows
        })

    conn.close()
    return render_template(
        "response_review_learner.html",
        learner_name=learner_row["full_name"] or learner_username,
        learner_username=learner_username,
        learner_group=learner_row["group_name"] or "—",
        item_label=f"{submission['subject']} - {submission['title']}" if submission["subject"] else submission["title"],
        review_rows=review_rows,
        score_summary={
            "score": submission["score"],
            "total": submission["total"],
            "percentage": submission["percentage"],
            "submitted_at": submission["submitted_at"][:16] if submission["submitted_at"] else "—"
        },
        can_manage_reviews=True
    )


@app.route("/response_review/true_false_accept", methods=["POST"])
def response_review_true_false_accept():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    test_id = request.form.get("test_id", type=int)
    question_id = request.form.get("question_id", type=int)
    acceptable_answer = (request.form.get("acceptable_answer") or "").strip()
    selected_group = (request.form.get("selected_group") or "").strip()
    learner_username = (request.form.get("learner_username") or "").strip()
    if not (test_id and question_id and acceptable_answer):
        flash("Accepted answer could not be added.", "error")
        return redirect(url_for("response_review"))

    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("SELECT question_type FROM theory_questions WHERE id = ?", (question_id,))
    row = cursor.fetchone()
    if not row:
        conn.close()
        flash("Question not found.", "error")
        return redirect(url_for("response_review"))
    q_type = row[0]

    if q_type == "true_false":
        cursor.execute("""
            INSERT INTO theory_options (question_id, option_text, is_correct, match_pair)
            VALUES (?, ?, 0, 'correction')
        """, (question_id, acceptable_answer))
    elif q_type == "fill_in":
        cursor.execute("""
            INSERT INTO theory_options (question_id, option_text, is_correct)
            VALUES (?, ?, 1)
        """, (question_id, acceptable_answer))
    conn.commit()
    conn.close()

    updated = regrade_theory_question_answers(test_id, question_id, selected_group or None, learner_username or None)
    flash(f"Accepted answer added. {updated} learner answer(s) regraded.", "success")
    if learner_username:
        return redirect(url_for("response_review_learner", learner=learner_username, item=test_id))
    return redirect(url_for("response_review", type="theory", group=selected_group, item=test_id))


@app.route("/response_review/accepted_answer_remove", methods=["POST"])
def response_review_accepted_answer_remove():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    test_id = request.form.get("test_id", type=int)
    question_id = request.form.get("question_id", type=int)
    remove_answer = (request.form.get("remove_answer") or "").strip()
    selected_group = (request.form.get("selected_group") or "").strip()
    if not (test_id and question_id and remove_answer):
        flash("Accepted answer could not be removed.", "error")
        return redirect(url_for("response_review"))

    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("""
        DELETE FROM theory_options
        WHERE question_id = ? AND option_text = ? AND (is_correct = 1 OR match_pair = 'correction')
    """, (question_id, remove_answer))
    conn.commit()
    conn.close()

    updated = regrade_theory_question_answers(test_id, question_id, selected_group or None)
    flash(f"Accepted answer removed. {updated} learner answer(s) regraded.", "success")
    return redirect(url_for("response_review", type="theory", group=selected_group, item=test_id))


@app.route("/teacher_dashboard/quick_actions", methods=["POST"])
def update_teacher_quick_actions():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    role = get_user_role(username)
    if role not in ["teacher", "admin"]:
        return "Access denied", 403

    selected = request.form.getlist("quick_actions")
    valid_ids = {action["key"] for action in get_teacher_quick_action_catalog()}
    selected = [action_id for action_id in selected if action_id in valid_ids]

    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("DELETE FROM teacher_quick_actions WHERE username = ?", (username,))
    for action_id in selected:
        cursor.execute("""
            INSERT INTO teacher_quick_actions (username, action_key)
            VALUES (?, ?)
        """, (username, action_id))
    conn.commit()
    conn.close()
    return redirect(url_for("teacher_dashboard"))


@app.route("/student/attendance_review_request", methods=["POST"])
def student_attendance_review_request():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    if get_user_role(username) != "student":
        return "Access denied", 403

    attendance_date = request.form.get("attendance_date", "").strip()
    message = request.form.get("message", "").strip()
    if not attendance_date:
        flash("Attendance date is required.", "error")
        return redirect(url_for("student_dashboard"))

    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("""
        SELECT teacher_username, full_name
        FROM users
        WHERE username = ?
    """, (username,))
    row = cursor.fetchone()
    if not row or not row[0]:
        conn.close()
        flash("No teacher is assigned to this learner.", "error")
        return redirect(url_for("student_dashboard"))

    teacher_username = row[0]
    thread_id = create_communication_thread(
        cursor,
        username,
        "attendance_review",
        subject_line="Attendance review request",
        attendance_date=attendance_date,
        initial_message=message or f"Please review my attendance for {attendance_date}.",
        chat_session_id=session.get("chat_session_id", "")
    )
    cursor.execute("""
        UPDATE communication_threads
        SET teacher_username = ?
        WHERE id = ?
    """, (teacher_username, thread_id))
    conn.commit()
    conn.close()
    flash("Attendance review request submitted.", "success")
    return redirect(url_for("student_dashboard"))


@app.route("/student/send_message", methods=["POST"])
def student_send_message():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    if get_user_role(username) != "student":
        return "Access denied", 403

    topic = (request.form.get("topic") or "chat").strip() or "chat"
    message = (request.form.get("message") or "").strip()
    if not message:
        flash("Message cannot be empty.", "error")
        return redirect(url_for("student_dashboard"))

    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("""
        SELECT teacher_username
        FROM users
        WHERE username = ?
    """, (username,))
    row = cursor.fetchone()
    if not row or not row[0]:
        conn.close()
        flash("No teacher is assigned to this learner.", "error")
        return redirect(url_for("student_dashboard"))

    teacher_username = row[0]
    chat_session_id = request.form.get("chat_session_id", "").strip() or session.get("chat_session_id", "")
    if not chat_session_id:
        chat_session_id = datetime.now().strftime("%Y%m%d%H%M%S")
        session["chat_session_id"] = chat_session_id

    cursor.execute("""
        SELECT id
        FROM communication_threads
        WHERE student_username = ? AND topic = ? AND status = 'open' AND chat_session_id = ?
        ORDER BY id DESC
        LIMIT 1
    """, (username, topic, chat_session_id))
    existing = cursor.fetchone()
    if existing:
        add_communication_message(cursor, existing[0], username, "student", message)
    else:
        thread_id = create_communication_thread(
            cursor,
            username,
            topic,
            subject_line="Chat" if topic == "chat" else topic.replace("_", " ").title(),
            initial_message=message,
            chat_session_id=chat_session_id
        )
        cursor.execute("""
            UPDATE communication_threads
            SET teacher_username = ?
            WHERE id = ?
        """, (teacher_username, thread_id))
    conn.commit()
    conn.close()
    flash("Message sent.", "success")
    return redirect(url_for("student_dashboard"))


@app.route("/communications")
def communications():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    role = get_user_role(username)
    if role not in ["teacher", "admin"]:
        return "Access denied", 403

    selected_group = (request.args.get("group") or "").strip()
    selected_topic = (request.args.get("topic") or "").strip()
    selected_date = (request.args.get("date") or "").strip()
    groups = get_groups(username)

    conn = get_db()
    conn.row_factory = sqlite3.Row
    cursor = conn.cursor()

    query = """
        SELECT t.*,
               u.full_name AS student_name,
               u.group_name
        FROM communication_threads t
        JOIN users u ON u.username = t.student_username
        WHERE 1=1
    """
    params = []
    if role == "teacher":
        query += " AND (t.teacher_username = ? OR u.teacher_username = ?)"
        params.extend([username, username])
    if selected_group:
        query += " AND u.group_name = ?"
        params.append(selected_group)
    if selected_topic:
        query += " AND t.topic = ?"
        params.append(selected_topic)
    if selected_date:
        query += " AND (t.attendance_date = ? OR substr(COALESCE(t.updated_at, t.created_at), 1, 10) = ?)"
        params.extend([selected_date, selected_date])
    query += " ORDER BY COALESCE(t.updated_at, t.created_at) DESC, t.id DESC"
    cursor.execute(query, params)

    threads = []
    for row in cursor.fetchall():
        thread = dict(row)
        cursor.execute("""
            SELECT *
            FROM communication_messages
            WHERE thread_id = ?
            ORDER BY COALESCE(created_at, '') ASC, id ASC
        """, (thread["id"],))
        thread["messages"] = [dict(message_row) for message_row in cursor.fetchall()]
        threads.append(thread)

    conn.close()
    return render_template(
        "communications.html",
        threads=threads,
        groups=groups,
        selected_group=selected_group,
        selected_topic=selected_topic,
        selected_date=selected_date
    )


@app.route("/communications/reply", methods=["POST"])
def communications_reply():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    role = get_user_role(username)
    if role not in ["teacher", "admin"]:
        return "Access denied", 403

    thread_id = request.form.get("thread_id", type=int)
    message = (request.form.get("message") or "").strip()
    if not thread_id or not message:
        flash("Reply could not be sent.", "error")
        return redirect(url_for("communications"))

    conn = get_db()
    cursor = conn.cursor()
    add_communication_message(cursor, thread_id, username, role, message)
    conn.commit()
    conn.close()
    flash("Reply sent.", "success")
    return redirect(url_for("communications", group=request.form.get("group", ""), topic=request.form.get("topic", ""), date=request.form.get("date", "")))


@app.route("/communications/mark_read", methods=["POST"])
def communications_mark_read():
    username = session.get("username")
    if not username:
        return jsonify({"ok": False}), 401
    role = get_user_role(username)
    if role not in ["teacher", "admin"]:
        return jsonify({"ok": False}), 403

    thread_id = request.form.get("thread_id", type=int)
    if not thread_id:
        return jsonify({"ok": False}), 400

    now_text = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("""
        UPDATE communication_threads
        SET teacher_read_at = ?
        WHERE id = ?
    """, (now_text, thread_id))
    conn.commit()
    conn.close()
    return jsonify({"ok": True})


@app.route("/communications/attendance_action", methods=["POST"])
def communications_attendance_action():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    role = get_user_role(username)
    if role not in ["teacher", "admin"]:
        return "Access denied", 403

    thread_id = request.form.get("thread_id", type=int)
    student_username = (request.form.get("student_username") or "").strip()
    attendance_date = (request.form.get("attendance_date") or "").strip()
    action = (request.form.get("action") or "").strip()
    status_map = {
        "mark_present": "present",
        "mark_late": "late",
        "mark_absent": "absent"
    }
    status = status_map.get(action)
    if not (thread_id and student_username and attendance_date and status):
        flash("Attendance action could not be completed.", "error")
        return redirect(url_for("communications"))

    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("""
        INSERT INTO attendance_override (username, date, status)
        VALUES (?, ?, ?)
        ON CONFLICT(username, date) DO UPDATE SET status = excluded.status
    """, (student_username, attendance_date, status))
    add_communication_message(
        cursor,
        thread_id,
        username,
        role,
        f"Attendance updated to {status.title()} for {attendance_date}."
    )
    conn.commit()
    conn.close()
    flash(f"Attendance marked {status}.", "success")
    return redirect(url_for("communications", group=request.form.get("group", ""), topic=request.form.get("topic", ""), date=request.form.get("date", "")))


@app.route("/api/message_status")
def api_message_status():
    username = session.get("username")
    if not username:
        return jsonify({"ok": False}), 401
    role = get_user_role(username)
    if role == "student":
        if request.args.get("mark_read") == "1":
            mark_student_threads_read(username)
        return jsonify({
            "ok": True,
            "unread": get_student_unread_message_count(username),
            "auto_open": student_has_fresh_teacher_reply(username)
        })
    if role in ["teacher", "admin"]:
        return jsonify({
            "ok": True,
            "unread": get_teacher_unread_message_count(username)
        })
    return jsonify({"ok": False}), 403



@app.route("/promote/<username>")
def promote(username):
    admin_user = session.get("username")
    
    if not admin_user:
        return redirect(url_for("login"))
    
    admin_role = get_user_role(admin_user)
    if admin_role not in ["teacher", "admin"]:
        return "Access denied", 403
    
    conn = get_db()
    cursor = conn.cursor()

    cursor.execute("""
    UPDATE users SET role = 'teacher' WHERE username = ?
    """, (username,))

    conn.commit()
    log_activity(admin_user, f"promoted {username} to teacher")
    conn.close()

    return redirect(url_for("admin_panel"))

@app.route("/demote/<username>")
def demote(username):
    admin_user = session.get("username")
    
    if not admin_user:
        return redirect(url_for("login"))
    
    admin_role = get_user_role(admin_user)
    if admin_role not in ["teacher", "admin"]:
        return "Access denied", 403
    
    conn = get_db()
    cursor = conn.cursor()

    cursor.execute("""
    UPDATE users SET role = 'student' WHERE username = ?
    """, (username,))

    conn.commit()
    log_activity(admin_user, f"demoted {username} to student")
    conn.close()

    return redirect(url_for("admin_panel"))

@app.route("/delete_user/<username>", methods=["POST"])
def delete_user(username):
    admin_user = session.get("username")
    if not admin_user:
        return redirect(url_for("login"))
    if get_user_role(admin_user) not in ["teacher", "admin"]:
        return "Access denied", 403
    if username == admin_user:
        return "Cannot delete your own account", 400

    conn = get_db()
    cursor = conn.cursor()

    # Delete all data associated with this user
    cursor.execute("DELETE FROM results WHERE username = ?", (username,))
    cursor.execute("DELETE FROM weaknesses WHERE username = ?", (username,))
    cursor.execute("DELETE FROM login_history WHERE username = ?", (username,))
    cursor.execute("DELETE FROM attendance_override WHERE username = ?", (username,))
    cursor.execute("DELETE FROM activities WHERE username = ?", (username,))
    cursor.execute("DELETE FROM learner_notes WHERE username = ?", (username,))
    cursor.execute("DELETE FROM result_removals WHERE username = ?", (username,))
    cursor.execute("""
        DELETE FROM theory_answers WHERE submission_id IN
            (SELECT id FROM theory_submissions WHERE username = ?)
    """, (username,))
    cursor.execute("DELETE FROM theory_submissions WHERE username = ?", (username,))
    cursor.execute("DELETE FROM users WHERE username = ?", (username,))

    conn.commit()
    conn.close()
    log_activity(admin_user, f"deleted user {username}")
    return redirect(url_for("admin_panel"))


def calculate_attendance_percentage(data, days):
    total_cells = len(data) * len(days)
    present = 0

    for row in data:
        for d in days:
            if row["days"].get(d):
                present += 1

    return round((present / total_cells) * 100) if total_cells else 0

@app.route("/term_dates", methods=["GET", "POST"])
def term_dates():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    if request.method == "POST":
        conn = get_db()
        cursor = conn.cursor()
        for term_num in range(1, 5):
            start = request.form.get(f"term{term_num}_start", "").strip()
            end = request.form.get(f"term{term_num}_end", "").strip()
            if start and end:
                cursor.execute("""
                    INSERT INTO term_dates (term, start_date, end_date)
                    VALUES (?, ?, ?)
                    ON CONFLICT(term) DO UPDATE SET start_date = excluded.start_date, end_date = excluded.end_date
                """, (term_num, start, end))
            else:
                cursor.execute("DELETE FROM term_dates WHERE term = ?", (term_num,))
        conn.commit()
        conn.close()
        log_activity(username, "updated term dates")
        return redirect(url_for("attendance"))

    return redirect(url_for("attendance"))


@app.route("/attendance")
def attendance():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    role = get_user_role(username)
    if role not in ["teacher", "admin"]:
        return "Access denied", 403

    selected_group = request.args.get("group")
    edit_mode = request.args.get("edit") == "1"   # ✅ key line
    range_param = request.args.get("range", "week")

    # Calculate date range
    today = datetime.now().strftime("%Y-%m-%d")
    if range_param == "week":
        start_date = (datetime.now() - timedelta(days=7)).strftime("%Y-%m-%d")
        end_date = today
    elif range_param == "2weeks":
        start_date = (datetime.now() - timedelta(days=14)).strftime("%Y-%m-%d")
        end_date = today
    elif range_param == "month":
        start_date = (datetime.now() - timedelta(days=30)).strftime("%Y-%m-%d")
        end_date = today
    else:
        start_date = (datetime.now() - timedelta(days=7)).strftime("%Y-%m-%d")
        end_date = today

    groups = get_groups(username) if role == 'teacher' else get_groups()

    if role == 'teacher' and selected_group and selected_group not in groups:
        selected_group = None

    days = []
    data = []
    daily_present_counts = {}
    daily_absent_counts = {}
    attendance_summary = []

    if selected_group:
        days, data = get_attendance_data(selected_group, start_date, end_date, teacher_username=username if role == 'teacher' else None)
        daily_present_counts = {
            day: sum(1 for row in data if row["days"].get(day))
            for day in days
        }
        daily_absent_counts = {
            day: len(data) - daily_present_counts[day]
            for day in days
        }
    else:
        conn = get_db()
        cursor = conn.cursor()
        attendance_summary = build_attendance_group_summary(
            cursor,
            groups,
            teacher_username=username if role == "teacher" else None,
            days=get_last_21_days(),
        )
        conn.close()

    # Get excluded dates
    conn = get_db()
    cursor = conn.cursor()
    if selected_group:
        cursor.execute("""
        SELECT date, group_name, reason, created_by, created_at
        FROM excluded_dates
        WHERE group_name IS NULL OR group_name = ?
        ORDER BY date DESC
        """, (selected_group,))
    else:
        cursor.execute("""
        SELECT date, group_name, reason, created_by, created_at
        FROM excluded_dates
        ORDER BY date DESC
        """)
    excluded_dates = cursor.fetchall()
    conn.close()

    terms = get_term_dates()

    return render_template(
        "attendance.html",
        groups=groups,
        selected_group=selected_group,
        days=days,
        data=data,
        daily_present_counts=daily_present_counts,
        daily_absent_counts=daily_absent_counts,
        edit_mode=edit_mode,
        today=datetime.now().strftime("%Y-%m-%d"),
        excluded_dates=excluded_dates,
        terms=terms,
        attendance_summary=attendance_summary,
    )

@app.route("/admin_panel")
@app.route("/admin")
def admin_panel():
    username = session.get("username")

    if not username:
        return redirect(url_for("login"))

    role = get_user_role(username)
    if role not in ["teacher", "admin"]:
        return "Access denied", 403

    search = request.args.get("search", "").strip()
    group = request.args.get("group", "").strip()
    sort = request.args.get("sort", "last_active").strip()
    order = request.args.get("order", "desc").lower()

    valid_sorts = {
        "username": "username",
        "full_name": "full_name",
        "group_name": "group_name",
        "teacher_username": "teacher_username",
        "role": "role",
        "last_active": "last_active"
    }

    if sort not in valid_sorts:
        sort = "last_active"
    if order not in ["asc", "desc"]:
        order = "desc"

    conn = get_db()
    cursor = conn.cursor()

    query = """
    SELECT username, full_name, group_name, teacher_username, role, last_active
    FROM users
    WHERE 1=1
    """
    params = []

    # Teachers only see their own students (plus themselves)
    if role == 'teacher':
        query += " AND (teacher_username = ? OR username = ?)"
        params.extend([username, username])

    if group:
        query += " AND group_name = ?"
        params.append(group)

    query += f" ORDER BY {valid_sorts[sort]} {order.upper()}"

    cursor.execute(query, params)
    users = cursor.fetchall()

    # Get group list for dropdown — teachers only see their own groups
    if role == 'teacher':
        my_groups = get_groups(username)
        groups = my_groups
    else:
        cursor.execute("SELECT DISTINCT group_name FROM users WHERE group_name IS NOT NULL")
        groups = [g[0] for g in cursor.fetchall()]

    conn.close()

    return render_template(
        "admin.html",
        users=users,
        groups=groups,
        search=search,
        selected_group=group,
        sort=sort,
        order=order
    )

@app.route("/import_users", methods=["POST"])
def import_users():
    username = session.get("username")
    
    if not username:
        return redirect(url_for("login"))
    
    role = get_user_role(username)
    if role not in ["teacher", "admin"]:
        return "Access denied", 403
    
    # Check if a file was uploaded
    if 'excel_file' not in request.files:
        flash("No file uploaded", "error")
        return redirect(url_for("admin_panel"))
    
    file = request.files['excel_file']
    if file.filename == '':
        flash("No file selected", "error")
        return redirect(url_for("admin_panel"))
    
    filename = file.filename or ""
    if not filename.lower().endswith(('.xlsx', '.xls')):
        flash("Please upload an Excel file (.xlsx or .xls)", "error")
        return redirect(url_for("admin_panel"))
    
    try:
        # Read the uploaded Excel file
        df = pd.read_excel(file)
        
        # Validate required columns
        required_columns = ['username', 'full_name', 'group']
        missing_columns = [col for col in required_columns if col not in df.columns]
        if missing_columns:
            flash(f"Missing required columns: {', '.join(missing_columns)}", "error")
            return redirect(url_for("admin_panel"))
        
        conn = get_db()
        cursor = conn.cursor()
        
        imported_count = 0
        updated_count = 0
        
        # Check which optional columns are present
        teacher_username_present = 'teacher_username' in df.columns
        role_present = 'role' in df.columns
        
        for _, row in df.iterrows():
            username_val = str(row["username"]).strip().upper()
            full_name = str(row["full_name"]).strip()
            group_name = str(row["group"]).strip()
            
            # Handle optional columns
            teacher_username = str(row["teacher_username"]).strip() if teacher_username_present else None
            if teacher_username == '':
                teacher_username = None
            
            role_value = str(row["role"]).strip().lower() if role_present else 'student'
            if role_value == '' or role_value not in ['student', 'teacher', 'admin']:
                role_value = 'student'
            
            # Check if user already exists
            cursor.execute("SELECT username FROM users WHERE username = ?", (username_val,))
            existing = cursor.fetchone()
            
            # Build dynamic SQL for UPSERT
            columns = ['username', 'full_name', 'group_name', 'teacher_username', 'role']
            values = [username_val, full_name, group_name, teacher_username, role_value]
            update_parts = ['full_name = excluded.full_name', 'group_name = excluded.group_name']
            
            if teacher_username_present:
                update_parts.append('teacher_username = excluded.teacher_username')
            
            if role_present:
                update_parts.append('role = excluded.role')
            
            update_clause = ', '.join(update_parts)
            
            cursor.execute(f"""
            INSERT INTO users ({', '.join(columns)})
            VALUES ({', '.join(['?'] * len(columns))})
            ON CONFLICT(username) DO UPDATE SET
                {update_clause}
            """, values)
            
            if existing:
                updated_count += 1
            else:
                imported_count += 1
        
        conn.commit()
        conn.close()
        
        flash(f"Successfully imported {imported_count} new users and updated {updated_count} existing users.", "success")
    
    except Exception as e:
        flash(f"Error importing users: {str(e)}", "error")
    
    return redirect(url_for("admin_panel"))

@app.route("/download_user_template")
def download_user_template():
    username = session.get("username")
    
    if not username:
        return redirect(url_for("login"))
    
    role = get_user_role(username)
    if role not in ["teacher", "admin"]:
        return "Access denied", 403
    
    # Create a sample DataFrame with the required columns
    sample_data = {
        'username': ['STUDENT001', 'STUDENT002', 'STUDENT003'],
        'full_name': ['Smith, John', 'Doe, Jane', 'Johnson, Bob'],
        'group': ['12A', '12A', '12B'],
        'teacher_username': ['TEACHER1', 'TEACHER1', 'TEACHER2'],
        'role': ['student', 'student', 'student']
    }
    df = pd.DataFrame(sample_data)
    
    # Create a BytesIO buffer to store the Excel file
    from io import BytesIO
    buffer = BytesIO()
    df.to_excel(buffer, index=False)
    buffer.seek(0)
    
    return send_file(
        buffer,
        as_attachment=True,
        download_name='user_import_template.xlsx',
        mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    )

@app.route("/recent_activity")
def recent_activity():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    role = get_user_role(username)
    if role not in ["teacher", "admin"]:
        return "Access denied", 403

    conn = get_db()
    cursor = conn.cursor()

    if role == 'teacher':
        cursor.execute("""
            SELECT a.username, a.action, strftime('%Y-%m-%d %H:%M:%S', a.timestamp)
            FROM activities a
            LEFT JOIN users u ON u.username = a.username
            WHERE a.username = ?
              OR (u.role = 'student' AND u.teacher_username = ?)
            ORDER BY a.timestamp DESC LIMIT 100
        """, (username, username))
    else:
        cursor.execute("""
            SELECT username, action, strftime('%Y-%m-%d %H:%M:%S', timestamp)
            FROM activities ORDER BY timestamp DESC LIMIT 100
        """)
    activities = cursor.fetchall()
    conn.close()
    return render_template("recent_activity.html", activities=activities)

@app.route("/edit_user/<username>", methods=["GET", "POST"])
def edit_user(username):
    admin_user = session.get("username")
    if not admin_user:
        return redirect(url_for("login"))

    if get_user_role(admin_user) not in ["teacher", "admin"]:
        return "Access denied", 403

    next_url = request.args.get("next") or request.form.get("next")

    conn = get_db()
    cursor = conn.cursor()

    cursor.execute("SELECT username, full_name, group_name, teacher_username, role FROM users WHERE username = ?", (username,))
    user = cursor.fetchone()
    if not user:
        conn.close()
        return "User not found", 404

    if get_user_role(admin_user) == "teacher":
        user_group = user[2] or ""
        assigned_teacher = user[3] or ""
        if assigned_teacher and assigned_teacher != admin_user:
            conn.close()
            return "Access denied", 403
        if not assigned_teacher and user[4] != "student":
            conn.close()
            return "Access denied", 403

    if request.method == "POST":
        full_name = request.form.get("full_name")
        group_name = request.form.get("group_name")
        teacher_username = request.form.get("teacher_username") or None
        role_value = request.form.get("role") or "student"

        if role_value == "admin" and get_user_role(admin_user) != "admin":
            return "Access denied", 403

        cursor.execute("""
        UPDATE users
        SET full_name = ?, group_name = ?, teacher_username = ?, role = ?
        WHERE username = ?
        """, (full_name, group_name, teacher_username, role_value, username))

        conn.commit()
        log_activity(admin_user, f"edited user {username}")
        conn.close()
        if next_url:
            return redirect(next_url)
        return redirect(url_for("admin_panel"))

    all_teachers = get_teachers()
    current_role = get_user_role(admin_user)
    conn.close()

    return render_template("edit_user.html", user=user, next_url=next_url, all_teachers=all_teachers, current_role=current_role)

@app.route("/learner_record/<username>")
def learner_record(username):
    admin_user = session.get("username")
    if not admin_user:
        return redirect(url_for("login"))
    admin_role = get_user_role(admin_user)
    if admin_role not in ["teacher", "admin"]:
        return "Access denied", 403

    conn = get_db()
    cursor = conn.cursor()

    cursor.execute("SELECT username, full_name, group_name, role, last_active FROM users WHERE username = ?", (username,))
    user = cursor.fetchone()
    if not user:
        conn.close()
        return "User not found", 404

    # Teachers can only view records of their own students
    if admin_role == 'teacher':
        cursor.execute("SELECT teacher_username FROM users WHERE username = ?", (username,))
        row = cursor.fetchone()
        if not row or row[0] != admin_user:
            conn.close()
            return "Access denied", 403

    # ── Attendance ────────────────────────────────────────────────────
    attendance_year = datetime.now().year
    days = get_current_year_attendance_days()
    cursor.execute("SELECT group_name FROM users WHERE username = ?", (username,))
    group_row = cursor.fetchone()
    user_group = group_row[0] if group_row else None
    auto_exclude_empty_attendance_days(cursor, [user_group], created_by=admin_user, days=days)
    cursor.execute("""
        SELECT date FROM excluded_dates
        WHERE group_name IS NULL OR group_name = ?
    """, (user_group,))
    excluded_dates = {row[0] for row in cursor.fetchall()}
    cursor.execute("""
        SELECT DISTINCT lh.date
        FROM login_history lh
        JOIN users u ON u.username = lh.username
        WHERE u.group_name = ? AND lh.username != ?
    """, (user_group, username))
    class_checked_dates = {row[0] for row in cursor.fetchall()}

    login_map = fetch_first_login_times(cursor, [username], days)
    override_map = fetch_attendance_override_statuses(cursor, [username], days)
    late_cutoffs = fetch_group_late_thresholds(cursor, user_group, days) if user_group else {}
    history = build_attendance_history(
        cursor,
        username,
        user_group,
        days,
        login_map=login_map,
        override_map=override_map,
        late_cutoffs=late_cutoffs,
        class_checked_dates=class_checked_dates,
        excluded_dates=excluded_dates,
    )
    attendance_summary = summarize_attendance_history(history)
    total_days = attendance_summary["total_days"]
    present_days = attendance_summary["present_days"]
    absent_days = attendance_summary["absent_days"]
    late_days = attendance_summary["late_days"]
    attendance_pct = attendance_summary["attendance_pct"]
    recent_history = history[-10:]  # last 10 days
    attendance_months = build_attendance_months(history)

    # ── Academic results ──────────────────────────────────────────────
    cursor.execute("""
        SELECT subject, task, MAX(score) as score, feedback, MAX(timestamp) as timestamp
        FROM results WHERE username = ? GROUP BY subject, task ORDER BY timestamp DESC
    """, (username,))
    all_results = cursor.fetchall()

    cursor.execute("""
        SELECT subject, ROUND(AVG(best_score),1)
        FROM (SELECT subject, task, MAX(score) as best_score FROM results WHERE username = ? GROUP BY subject, task)
        GROUP BY subject
    """, (username,))
    subject_avgs = cursor.fetchall()

    cursor.execute("""
        SELECT ROUND(AVG(best_score),1)
        FROM (SELECT subject, task, MAX(score) as best_score FROM results WHERE username = ? GROUP BY subject, task)
    """, (username,))
    overall_avg_row = cursor.fetchone()
    practical_avg = overall_avg_row[0] if overall_avg_row and overall_avg_row[0] else 0

    cursor.execute("""
        SELECT ROUND(AVG(best_pct),1)
        FROM (SELECT test_id, MAX(percentage) as best_pct FROM theory_submissions WHERE username = ? GROUP BY test_id)
    """, (username,))
    theory_avg_row = cursor.fetchone()
    theory_avg = theory_avg_row[0] if theory_avg_row and theory_avg_row[0] else None

    overall_avg = round((practical_avg + theory_avg) / 2, 1) if practical_avg and theory_avg else (practical_avg or theory_avg or 0)

    recent_results = all_results[:10]

    # Trend: compare last 5 vs previous 5
    scores = [r[2] for r in all_results if r[2] is not None]
    if len(scores) >= 6:
        recent_avg = sum(scores[:3]) / 3
        older_avg = sum(scores[3:6]) / 3
        if recent_avg > older_avg + 2:
            trend = "improving"
        elif recent_avg < older_avg - 2:
            trend = "dropping"
        else:
            trend = "stable"
    else:
        trend = "not enough data"

    # Task status rows - DYNAMICALLY fetch only assigned tasks
    cursor.execute("SELECT group_name FROM users WHERE username = ?", (username,))
    user_group_row = cursor.fetchone()
    user_group = user_group_row[0] if user_group_row else None
    
    # Get practical tasks assigned to this user's group
    cursor.execute("""
        SELECT DISTINCT s.name as subject, t.name as task_name
        FROM tasks t
        JOIN subjects s ON t.subject_id = s.id
        JOIN task_groups tg ON t.id = tg.task_id
        WHERE tg.group_name = ? AND t.task_type = 'practical'
        ORDER BY s.name, t.name
    """, (user_group,))
    assigned_practical_tasks = cursor.fetchall()
    
    # Get theory tests assigned to this user's group
    cursor.execute("""
        SELECT DISTINCT tt.id, tt.title, COALESCE(tt.subject, ''),
               COALESCE(q_stats.total_items, 0),
               COALESCE(q_stats.test_questions, 0)
        FROM theory_tests tt
        LEFT JOIN theory_test_groups ttg ON tt.id = ttg.test_id
        LEFT JOIN (
            SELECT test_id,
                   COUNT(*) as total_items,
                   SUM(CASE WHEN question_type IN ('content_slide', 'title_slide', 'heading_slide') THEN 0 ELSE 1 END) as test_questions
            FROM theory_questions
            GROUP BY test_id
        ) q_stats ON q_stats.test_id = tt.id
        WHERE (ttg.group_name = ? OR ttg.group_name IS NULL)
        ORDER BY tt.title
    """, (user_group,))
    assigned_theory_tests = cursor.fetchall()
    
    # Build task rows from assigned tasks only
    results_map = {(r[0], r[1]): {"score": r[2], "feedback": r[3], "timestamp": r[4]} for r in all_results}
    task_rows = []
    for subject, task in assigned_practical_tasks:
        row = results_map.get((subject, task))
        task_rows.append({
            "subject": subject, "task": task,
            "score": row["score"] if row else None,
            "feedback": row["feedback"] if row else None,
            "timestamp": row["timestamp"] if row else None,
            "status": "Submitted" if row else "Not submitted",
            "type": "practical"
        })
    
    # Get theory test submissions
    theory_submissions_map = {}
    for test_id, title, subject, total_items, test_questions in assigned_theory_tests:
        cursor.execute("""
            SELECT MAX(percentage), MAX(submitted_at), COUNT(*), COALESCE(SUM(time_spent_seconds), 0),
                   COALESCE(MAX(submission_type), 'test')
            FROM theory_submissions
            WHERE username = ? AND test_id = ?
        """, (username, test_id))
        result = cursor.fetchone()
        if result and result[0] is not None:
            theory_submissions_map[test_id] = {
                "score": result[0],
                "timestamp": result[1],
                "attempts": result[2] or 0,
                "time_spent": result[3] or 0,
                "submission_type": result[4] or "test"
            }
    
    # Add theory tests to task rows
    for test_id, title, subject, total_items, test_questions in assigned_theory_tests:
        row = theory_submissions_map.get(test_id)
        task_rows.append({
            "subject": "Theory", "task": title,
            "score": row["score"] if row else None,
            "feedback": "",
            "timestamp": row["timestamp"] if row else None,
            "status": "Submitted" if row else "Not submitted",
            "type": "theory"
        })
    
    average = round(sum(t["score"] for t in task_rows if t["score"] is not None) /
                    max(1, sum(1 for t in task_rows if t["score"] is not None)), 1) \
              if any(t["score"] is not None for t in task_rows) else 0

    # ── Theory test results ───────────────────────────────────────────
    cursor.execute("""
        SELECT tt.title, ts.score, ts.total, ts.percentage, ts.submitted_at,
               COALESCE(ts.time_spent_seconds, 0), COALESCE(ts.submission_type, 'test')
        FROM theory_submissions ts
        JOIN theory_tests tt ON ts.test_id = tt.id
        JOIN (
            SELECT test_id, MAX(submitted_at) as latest
            FROM theory_submissions
            WHERE username = ?
            GROUP BY test_id
        ) latest ON latest.test_id = ts.test_id AND latest.latest = ts.submitted_at
        WHERE ts.username = ?
        ORDER BY ts.submitted_at DESC LIMIT 10
    """, (username, username))
    theory_results = cursor.fetchall()

    cursor.execute("""
        SELECT test_id, current_slide, max_slide, COALESCE(time_spent_seconds, 0),
               COALESCE(completed, 0), updated_at
        FROM theory_progress
        WHERE username = ?
    """, (username,))
    theory_progress_map = {
        row[0]: {
            "current_slide": row[1] or 0,
            "max_slide": row[2] or 0,
            "time_spent": row[3] or 0,
            "completed": bool(row[4]),
            "updated_at": row[5]
        }
        for row in cursor.fetchall()
    }

    practical_total = len(assigned_practical_tasks)
    practical_done = sum(1 for row in task_rows if row["type"] == "practical" and row["score"] is not None)
    theory_total = len(assigned_theory_tests)
    theory_done = 0
    theory_in_progress = 0
    missing_items = []
    learning_progress = []
    theory_time_seconds = 0

    for test_id, title, subject, total_items, test_questions in assigned_theory_tests:
        total_items = total_items or 0
        test_questions = test_questions or 0
        is_lesson = total_items > 0 and test_questions == 0
        submission = theory_submissions_map.get(test_id)
        progress = theory_progress_map.get(test_id, {})
        completed = bool(submission) or bool(progress.get("completed"))
        has_progress = test_id in theory_progress_map
        viewed = min(total_items, (progress.get("max_slide", 0) + 1) if has_progress and total_items else 0)
        progress_pct = 100 if completed else (round((viewed / total_items) * 100) if total_items else 0)
        time_spent = submission["time_spent"] if submission else progress.get("time_spent", 0)
        theory_time_seconds += time_spent or 0

        if completed:
            theory_done += 1
            status = "Completed"
        elif progress_pct > 0:
            theory_in_progress += 1
            status = "In progress"
        else:
            status = "Not started"
            missing_items.append({"type": "Lesson" if is_lesson else "Test", "title": title})

        learning_progress.append({
            "id": test_id,
            "title": title,
            "subject": subject or "Theory",
            "kind": "Lesson" if is_lesson else "Test",
            "status": status,
            "progress_pct": progress_pct,
            "viewed": viewed,
            "has_progress": has_progress,
            "total": total_items,
            "score": submission["score"] if submission else None,
            "attempts": submission["attempts"] if submission else 0,
            "time_spent": time_spent or 0,
            "updated_at": (submission["timestamp"] if submission else progress.get("updated_at"))
        })

    practical_missing = [
        {"type": "Practical", "title": f"{row['subject']} - {row['task']}"}
        for row in task_rows
        if row["type"] == "practical" and row["score"] is None
    ]
    missing_items = practical_missing + missing_items
    learning_summary = {
        "practical_total": practical_total,
        "practical_done": practical_done,
        "theory_total": theory_total,
        "theory_done": theory_done,
        "theory_in_progress": theory_in_progress,
        "missing_total": len(missing_items),
        "theory_time_seconds": theory_time_seconds,
        "overall_completion": round(((practical_done + theory_done) / max(1, practical_total + theory_total)) * 100, 1)
    }

    # ── Weaknesses ────────────────────────────────────────────────────
    cursor.execute("""
        SELECT skill, count FROM weaknesses
        WHERE username = ? ORDER BY count DESC LIMIT 10
    """, (username,))
    weaknesses = cursor.fetchall()

    # ── Recent activity ───────────────────────────────────────────────
    cursor.execute("""
        SELECT action, timestamp FROM activities
        WHERE username = ? ORDER BY timestamp DESC LIMIT 10
    """, (username,))
    recent_activity = cursor.fetchall()

    # ── Teacher notes ─────────────────────────────────────────────────
    cursor.execute("SELECT id, note, flag, created_by, created_at FROM learner_notes WHERE username = ? ORDER BY created_at DESC", (username,))
    notes = cursor.fetchall()

    conn.close()
    return render_template(
        "learner_record.html",
        user=user,
        history=recent_history,
        attendance_calendar=history,
        attendance_months=attendance_months,
        attendance_year=attendance_year,
        attendance_pct=attendance_pct,
        present_days=present_days,
        absent_days=absent_days,
        late_days=late_days,
        total_days=total_days,
        task_rows=task_rows,
        average=average,
        overall_avg=overall_avg,
        practical_avg=practical_avg,
        theory_avg=theory_avg,
        subject_avgs=subject_avgs,
        recent_results=recent_results,
        theory_results=theory_results,
        learning_summary=learning_summary,
        learning_progress=learning_progress,
        missing_items=missing_items[:12],
        trend=trend,
        weaknesses=weaknesses,
        recent_activity=recent_activity,
        notes=notes
    )

@app.route("/learner_record/<username>/add_note", methods=["POST"])
def add_learner_note(username):
    admin_user = session.get("username")
    if not admin_user:
        return redirect(url_for("login"))
    if get_user_role(admin_user) not in ["teacher", "admin"]:
        return "Access denied", 403
    note = request.form.get("note", "").strip()
    flag = request.form.get("flag", "").strip()
    if note:
        conn = get_db()
        cursor = conn.cursor()
        cursor.execute("""
            INSERT INTO learner_notes (username, note, flag, created_by, created_at)
            VALUES (?, ?, ?, ?, ?)
        """, (username, note, flag, admin_user, datetime.now().isoformat()))
        conn.commit()
        conn.close()
    return redirect(url_for("learner_record", username=username))


@app.route("/learner_record/<username>/delete_note/<int:note_id>", methods=["POST"])
def delete_learner_note(username, note_id):
    admin_user = session.get("username")
    if not admin_user:
        return redirect(url_for("login"))
    if get_user_role(admin_user) not in ["teacher", "admin"]:
        return "Access denied", 403
    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("DELETE FROM learner_notes WHERE id = ? AND username = ?", (note_id, username))
    conn.commit()
    conn.close()
    return redirect(url_for("learner_record", username=username))


@app.route("/remove_results", methods=["POST"])
def remove_results():
    teacher = session.get("username")
    if not teacher or get_user_role(teacher) not in ["teacher", "admin"]:
        return "Access denied", 403

    task_type  = request.form.get("task_type")   # 'practical' or 'theory'
    subject    = request.form.get("subject", "")
    task_name  = request.form.get("task_name", "")
    test_id    = request.form.get("test_id", "")
    target     = request.form.get("target")       # 'all' or a specific username
    reason     = request.form.get("reason", "").strip()
    group      = request.form.get("group", "")

    if not reason:
        reason = "Removed by teacher"

    conn = get_db()
    cursor = conn.cursor()

    # Determine which students are affected
    if target == "all":
        if task_type == "practical":
            cursor.execute("SELECT DISTINCT username FROM results WHERE subject = ? AND task = ?",
                           (subject, task_name))
        else:
            cursor.execute("SELECT DISTINCT username FROM theory_submissions WHERE test_id = ?",
                           (test_id,))
        affected = [row[0] for row in cursor.fetchall()]
    else:
        affected = [target]

    now = datetime.now().isoformat()

    for username in affected:
        # Log the removal
        cursor.execute("""
            INSERT INTO result_removals
                (username, task_type, subject, task_name, test_id, removed_by, reason, removed_at)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
        """, (username, task_type, subject, task_name,
               test_id if task_type == "theory" else None,
               teacher, reason, now))

        # Also add a learner note so it shows on their record
        note_text = (f"⚠️ Marks removed by {teacher}: "
                     f"{'[Theory] ' + task_name if task_type == 'theory' else subject + ' - ' + task_name}. "
                     f"Reason: {reason}")
        cursor.execute("""
            INSERT INTO learner_notes (username, note, flag, created_by, created_at)
            VALUES (?, ?, 'warning', ?, ?)
        """, (username, note_text, teacher, now))

        # Delete the actual results
        if task_type == "practical":
            cursor.execute("DELETE FROM results WHERE username = ? AND subject = ? AND task = ?",
                           (username, subject, task_name))
        else:
            cursor.execute("""
                DELETE FROM theory_answers WHERE submission_id IN
                    (SELECT id FROM theory_submissions WHERE username = ? AND test_id = ?)
            """, (username, test_id))
            cursor.execute("DELETE FROM theory_submissions WHERE username = ? AND test_id = ?",
                           (username, test_id))

    conn.commit()
    log_activity(teacher, f"removed {'all' if target == 'all' else target}'s results for "
                          f"{'[Theory] ' + task_name if task_type == 'theory' else subject + ' - ' + task_name}. "
                          f"Reason: {reason}")
    conn.close()
    return redirect(url_for("group_results", group=group))


@app.route("/export/results")
def export_results():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    role = get_user_role(username)
    if role not in ["teacher", "admin"]:
        return "Access denied", 403
    all_groups = get_groups(username) if role == 'teacher' else get_groups()
    return render_template("export_results.html", all_groups=all_groups)


@app.route("/export_results_multi", methods=["POST"])
def export_results_multi():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    selected_groups = request.form.getlist("groups")

    if not selected_groups:
        return "No groups selected", 400

    conn = get_db()
    cursor = conn.cursor()

    file_path = "multi_group_results_export.xlsx"

    with pd.ExcelWriter(file_path, engine='openpyxl') as writer:
        for group in selected_groups:
            # Get all students in the group
            cursor.execute("""
                SELECT username, full_name FROM users
                WHERE group_name = ? AND role = 'student'
                ORDER BY full_name
            """, (group,))
            students_raw = cursor.fetchall()

            # Get practical tasks assigned to this group
            cursor.execute("""
                SELECT DISTINCT s.name as subject, t.name as task_name, t.id
                FROM tasks t
                JOIN subjects s ON t.subject_id = s.id
                JOIN task_groups tg ON t.id = tg.task_id
                WHERE tg.group_name = ? AND t.task_type = 'practical'
                ORDER BY s.name, t.name
            """, (group,))
            practical_tasks = cursor.fetchall()

            # Get theory tests assigned to this group
            cursor.execute("""
                SELECT DISTINCT tt.id, tt.title, tt.subject
                FROM theory_tests tt
                LEFT JOIN theory_test_groups ttg ON tt.id = ttg.test_id
                WHERE (ttg.group_name = ? OR ttg.group_name IS NULL)
                ORDER BY tt.subject, tt.title
            """, (group,))
            theory_tasks = cursor.fetchall()

            # Build rows for Excel
            rows = []
            for student_username, full_name in students_raw:
                row = {
                    "Username": student_username,
                    "Name": full_name or student_username,
                    "Group": group
                }

                # Get practical results
                for subject, task_name, task_id in practical_tasks:
                    cursor.execute("""
                        SELECT MAX(score)
                        FROM results
                        WHERE username = ? AND subject = ? AND task = ?
                    """, (student_username, subject, task_name))
                    result = cursor.fetchone()
                    col_name = f"{subject} - {task_name}"
                    row[col_name] = result[0] if result and result[0] is not None else ""

                # Get theory results
                for test_id, title, subject in theory_tasks:
                    cursor.execute("""
                        SELECT MAX(percentage)
                        FROM theory_submissions
                        WHERE username = ? AND test_id = ?
                    """, (student_username, test_id))
                    result = cursor.fetchone()
                    col_name = f"[Theory] {title}"
                    row[col_name] = result[0] if result and result[0] is not None else ""

                # Calculate overall average
                all_scores = []
                for subject, task_name, task_id in practical_tasks:
                    cursor.execute("""
                        SELECT MAX(score)
                        FROM results
                        WHERE username = ? AND subject = ? AND task = ?
                    """, (student_username, subject, task_name))
                    result = cursor.fetchone()
                    if result and result[0] is not None:
                        all_scores.append(result[0])
                
                for test_id, title, subject in theory_tasks:
                    cursor.execute("""
                        SELECT MAX(percentage)
                        FROM theory_submissions
                        WHERE username = ? AND test_id = ?
                    """, (student_username, test_id))
                    result = cursor.fetchone()
                    if result and result[0] is not None:
                        all_scores.append(result[0])
                
                row["Overall Average"] = round(sum(all_scores) / len(all_scores), 1) if all_scores else ""

                rows.append(row)

            # Create DataFrame and write to Excel
            df = pd.DataFrame(rows)
            # Clean sheet name (remove special characters)
            sheet_name = group.replace("/", "_").replace("\\", "_")[:31]  # Excel sheet name limit
            df.to_excel(writer, sheet_name=sheet_name, index=False)

    log_activity(username, "exported results by group")
    conn.close()

    response = send_file(file_path, as_attachment=True)
    response.headers["HX-Redirect"] = url_for("teacher_dashboard")
    return response

@app.route("/export/attendance")
def export_attendance():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403
    
    conn = get_db()

    group = request.args.get("group")
    role = get_user_role(username)
    days, data = get_attendance_data(group, teacher_username=username if role == 'teacher' else None)

    rows = []

    for row in data:
        base = {
            "Username": row["username"],
            "Name": row["name"],
            "Group": row["group"],
            "Attendance %": row["attendance_pct"]
        }

        for d in days:
            val = row["days"].get(d)
            base[d] = val["time"] if val else "A"

        rows.append(base)

    df = pd.DataFrame(rows)
    
    file_path = "attendance_export.xlsx"
    df.to_excel(file_path, index=False)

    log_activity(username, "exported attendance")

    conn.close()

    return send_file(file_path, as_attachment=True)

@app.route("/risk_learners")
def risk_learners():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    role = get_user_role(username)
    if role not in ["teacher", "admin"]:
        return "Access denied", 403

    selected_group = request.args.get("group")
    groups = get_groups(username) if role == 'teacher' else get_groups()

    # Teachers can only view their own groups
    if role == 'teacher' and selected_group and selected_group not in groups:
        selected_group = None

    conn = get_db()
    cursor = conn.cursor()
    if not selected_group:
        practical_avg_map = fetch_student_practical_averages(cursor, username if role == "teacher" else None)
        theory_avg_map = fetch_student_theory_averages(cursor, username if role == "teacher" else None)
        attendance_summary = build_attendance_group_summary(
            cursor,
            groups,
            teacher_username=username if role == "teacher" else None,
            days=get_last_21_days(),
        )
        group_summaries = []
        for item in attendance_summary:
            group_name = item["group"]
            if role == "teacher":
                cursor.execute("""
                    SELECT username
                    FROM users
                    WHERE group_name = ? AND role = 'student' AND teacher_username = ?
                """, (group_name, username))
            else:
                cursor.execute("""
                    SELECT username
                    FROM users
                    WHERE group_name = ? AND role = 'student'
                """, (group_name,))
            usernames = [row[0] for row in cursor.fetchall()]
            combined_scores = []
            for uname in usernames:
                practical_avg = practical_avg_map.get(uname)
                theory_avg = theory_avg_map.get(uname)
                if practical_avg is not None and theory_avg is not None:
                    combined_scores.append(round((practical_avg + theory_avg) / 2, 1))
                elif practical_avg is not None:
                    combined_scores.append(practical_avg)
                elif theory_avg is not None:
                    combined_scores.append(theory_avg)
            avg_combined = round(sum(combined_scores) / len(combined_scores), 1) if combined_scores else 0
            risk_value = round((100 - item["attendance_pct"]) * 0.4 + (100 - avg_combined) * 0.4)
            group_summaries.append({
                "group": group_name,
                "students": item["students"],
                "attendance_pct": item["attendance_pct"],
                "avg_combined": avg_combined,
                "risk_score": risk_value,
            })
        conn.close()
        return render_template("Riks_learners.html", groups=groups, selected_group=None, summary_cards=[], group_summaries=group_summaries)

    # Get students for this group
    if role == 'teacher':
        cursor.execute("""
            SELECT username, full_name
            FROM users
            WHERE group_name = ? AND role = 'student' AND teacher_username = ?
            ORDER BY full_name
        """, (selected_group, username))
    else:
        cursor.execute("""
            SELECT username, full_name
            FROM users
            WHERE group_name = ? AND role = 'student'
            ORDER BY full_name
        """, (selected_group,))
    students_raw = cursor.fetchall()

    days_21 = get_last_21_days()
    auto_exclude_empty_attendance_days(cursor, [selected_group], created_by=username, days=get_current_year_attendance_days())
    excluded_days = fetch_group_excluded_dates(cursor, [selected_group]).get(selected_group, set())
    filtered_days = [day for day in days_21 if day not in excluded_days]
    student_usernames = [student_username for student_username, _ in students_raw]
    login_times = fetch_first_login_times(cursor, student_usernames, filtered_days)
    overrides = fetch_attendance_override_statuses(cursor, student_usernames, filtered_days)
    practical_avg_map = fetch_student_practical_averages(cursor, username if role == "teacher" else None)
    theory_avg_map = fetch_student_theory_averages(cursor, username if role == "teacher" else None)
    today = datetime.now().strftime("%Y-%m-%d")
    risk_students = []

    for student_username, full_name in students_raw:
        practical_avg = practical_avg_map.get(student_username)
        theory_avg = theory_avg_map.get(student_username)
        if practical_avg is not None and theory_avg is not None:
            avg_score = round((practical_avg + theory_avg) / 2, 1)
        else:
            avg_score = practical_avg if practical_avg is not None else theory_avg
        avg_score = avg_score if avg_score is not None else 0

        history = build_attendance_history(
            cursor,
            student_username,
            selected_group,
            filtered_days,
            login_map=login_times,
            override_map=overrides,
            excluded_dates=excluded_days,
        )
        attendance_pct = summarize_attendance_history(history)["attendance_pct"]

        cursor.execute("""
            SELECT COUNT(DISTINCT t.id)
            FROM task_groups tg
            JOIN tasks t ON t.id = tg.task_id
            WHERE tg.group_name = ?
              AND t.assign_date <= ?
        """, (selected_group, today))
        total_assigned = cursor.fetchone()[0] or 0

        cursor.execute("""
            SELECT COUNT(DISTINCT t.id)
            FROM task_groups tg
            JOIN tasks t ON t.id = tg.task_id
            JOIN subjects s ON s.id = t.subject_id
            WHERE tg.group_name = ?
              AND t.assign_date <= ?
              AND NOT EXISTS (
                  SELECT 1 FROM results r
                  WHERE r.username = ?
                    AND r.subject = s.name
                    AND r.task = t.name
              )
        """, (selected_group, today, student_username))
        missing_practical = cursor.fetchone()[0] or 0

        cursor.execute("""
            SELECT COUNT(DISTINCT tt.id)
            FROM theory_tests tt
            LEFT JOIN theory_test_groups ttg ON tt.id = ttg.test_id
            WHERE (ttg.group_name = ? OR ttg.group_name IS NULL)
              AND tt.assign_date <= ?
        """, (selected_group, today))
        total_theory = cursor.fetchone()[0] or 0

        cursor.execute("""
            SELECT COUNT(DISTINCT tt.id)
            FROM theory_tests tt
            LEFT JOIN theory_test_groups ttg ON tt.id = ttg.test_id
            WHERE (ttg.group_name = ? OR ttg.group_name IS NULL)
              AND tt.assign_date <= ?
              AND NOT EXISTS (
                  SELECT 1 FROM theory_submissions ts
                  WHERE ts.username = ?
                    AND ts.test_id = tt.id
              )
        """, (selected_group, today, student_username))
        missing_theory = cursor.fetchone()[0] or 0

        total_tasks = total_assigned + total_theory
        missing_pct = round((missing_practical + missing_theory) / total_tasks * 100) if total_tasks else 0

        risk_score = round((100 - attendance_pct) * 0.4 + (100 - avg_score) * 0.4 + missing_pct * 0.2)
        if risk_score <= 40:
            status = 'Safe'
        elif risk_score <= 70:
            status = 'At Risk'
        else:
            status = 'High Risk'

        reasons = []
        if attendance_pct < 60:
            reasons.append(f"A {attendance_pct}%")
        if avg_score < 40:
            reasons.append(f"AVG{avg_score}%")
        if missing_pct > 70:
            reasons.append(f"Missing {missing_pct}%")
        if not reasons:
            reasons.append("balanced risk factors")

        risk_students.append({
            'username': student_username,
            'name': full_name or student_username,
            'group': selected_group,
            'attendance_pct': attendance_pct,
            'avg_score': avg_score,
            'missing_pct': missing_pct,
            'score': risk_score,
            'status': status,
            'reason': ' + '.join(reasons)
        })

    safe_count = sum(1 for student in risk_students if student["status"] == "Safe")
    at_risk_count = sum(1 for student in risk_students if student["status"] == "At Risk")
    high_risk_count = sum(1 for student in risk_students if student["status"] == "High Risk")
    avg_attendance = round(sum(student["attendance_pct"] for student in risk_students) / len(risk_students), 1) if risk_students else 0
    avg_combined = round(sum(student["avg_score"] for student in risk_students) / len(risk_students), 1) if risk_students else 0
    summary_cards = [
        {"label": "Safe", "value": safe_count, "tone": "ok"},
        {"label": "At Risk", "value": at_risk_count, "tone": "mid"},
        {"label": "High Risk", "value": high_risk_count, "tone": "risk"},
        {"label": "Avg Attendance", "value": f"{avg_attendance}%", "tone": "info"},
        {"label": "Avg Combined", "value": f"{avg_combined}%", "tone": "info"},
    ]

    conn.commit()
    conn.close()
    return render_template("Riks_learners.html", groups=groups, selected_group=selected_group,
                           risk_students=risk_students, summary_cards=summary_cards, group_summaries=[])


@app.route("/group_results")
def group_results():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    role = get_user_role(username)
    if role not in ["teacher", "admin"]:
        return "Access denied", 403

    selected_group = request.args.get("group")
    groups = get_groups(username) if role == 'teacher' else get_groups()

    # Teachers cannot access groups that aren't theirs
    if role == 'teacher' and selected_group and selected_group not in groups:
        selected_group = None

    conn = get_db()
    cursor = conn.cursor()
    if not selected_group:
        practical_group_avg_map = fetch_group_practical_averages(cursor, groups)
        theory_group_avg_map = fetch_group_theory_averages(cursor, groups)
        result_summaries = []
        for group_name in groups:
            if role == 'teacher':
                cursor.execute("""
                    SELECT COUNT(*)
                    FROM users
                    WHERE group_name = ? AND role = 'student' AND teacher_username = ?
                """, (group_name, username))
            else:
                cursor.execute("""
                    SELECT COUNT(*)
                    FROM users
                    WHERE group_name = ? AND role = 'student'
                """, (group_name,))
            student_count = cursor.fetchone()[0] or 0
            practical_avg = practical_group_avg_map.get(group_name)
            theory_avg = theory_group_avg_map.get(group_name)
            if practical_avg is not None and theory_avg is not None:
                combined_avg = round((practical_avg + theory_avg) / 2, 1)
            else:
                combined_avg = practical_avg if practical_avg is not None else theory_avg
            result_summaries.append({
                "group": group_name,
                "students": student_count,
                "practical_avg": practical_avg,
                "theory_avg": theory_avg,
                "combined_avg": combined_avg or 0,
            })
        conn.close()
        return render_template("group_results.html", groups=groups, selected_group=None, result_summaries=result_summaries)

    # Get students — teachers only see their own students
    if role == 'teacher':
        cursor.execute("""
            SELECT username, full_name FROM users
            WHERE group_name = ? AND role = 'student' AND teacher_username = ?
            ORDER BY full_name
        """, (selected_group, username))
    else:
        cursor.execute("""
            SELECT username, full_name FROM users
            WHERE group_name = ? AND role = 'student'
            ORDER BY full_name
        """, (selected_group,))
    students_raw = cursor.fetchall()

    # Get practical tasks assigned to this group
    cursor.execute("""
        SELECT DISTINCT s.name as subject, t.name as task_name, t.id
        FROM tasks t
        JOIN subjects s ON t.subject_id = s.id
        JOIN task_groups tg ON t.id = tg.task_id
        WHERE tg.group_name = ? AND t.task_type = 'practical'
        ORDER BY s.name, t.name
    """, (selected_group,))
    practical_tasks = [{"subject": row[0], "name": row[1], "id": row[2]} for row in cursor.fetchall()]

    # Get theory tests assigned to this group (include inactive tests if they have submissions)
    cursor.execute("""
        SELECT DISTINCT tt.id, tt.title, tt.subject
        FROM theory_tests tt
        LEFT JOIN theory_test_groups ttg ON tt.id = ttg.test_id
        WHERE (ttg.group_name = ? OR ttg.group_name IS NULL)
        ORDER BY tt.subject, tt.title
    """, (selected_group,))
    theory_tasks = [{"test_id": row[0], "title": row[1], "subject": row[2]} for row in cursor.fetchall()]

    # Build student results
    students = []
    total_averages = []

    for username, full_name in students_raw:
        student = {
            "username": username,
            "full_name": full_name,
            "practical_results": {},
            "theory_results": {},
            "overall_avg": None
        }

        # Get practical results for this student
        for task in practical_tasks:
            cursor.execute("""
                SELECT score, timestamp
                FROM results
                WHERE username = ? AND subject = ? AND task = ?
                ORDER BY timestamp DESC
            """, (username, task["subject"], task["name"]))
            scores = cursor.fetchall()
            
            if scores:
                all_scores = [s[0] for s in scores]
                best_score = max(all_scores)
                student["practical_results"][(task["subject"], task["name"])] = {
                    "best_score": best_score,
                    "attempts": len(scores),
                    "all_scores": all_scores
                }

        # Get theory results for this student
        for task in theory_tasks:
            cursor.execute("""
                SELECT percentage, submitted_at
                FROM theory_submissions
                WHERE username = ? AND test_id = ?
                ORDER BY submitted_at DESC
            """, (username, task["test_id"]))
            scores = cursor.fetchall()
            
            if scores:
                all_scores = [s[0] for s in scores]
                best_score = max(all_scores)
                student["theory_results"][task["test_id"]] = {
                    "best_score": best_score,
                    "attempts": len(scores),
                    "all_scores": all_scores
                }

        # Calculate overall average (best scores only)
        all_best_scores = []
        all_best_scores.extend([r["best_score"] for r in student["practical_results"].values()])
        all_best_scores.extend([r["best_score"] for r in student["theory_results"].values()])
        
        if all_best_scores:
            student["overall_avg"] = round(sum(all_best_scores) / len(all_best_scores), 1)
            total_averages.append(student["overall_avg"])

        students.append(student)

    # Calculate group average
    group_average = round(sum(total_averages) / len(total_averages), 1) if total_averages else 0

    conn.close()

    return render_template(
        "group_results.html",
        groups=groups,
        selected_group=selected_group,
        students=students,
        practical_tasks=practical_tasks,
        theory_tasks=theory_tasks,
        group_average=group_average,
        result_summaries=[]
    )


@app.route("/export_attendance_form")
def export_attendance_form():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    role = get_user_role(username)
    if role not in ["teacher", "admin"]:
        return "Access denied", 403
    all_groups = get_groups(username) if role == 'teacher' else get_groups()
    return render_template("export_attendance.html", all_groups=all_groups)

@app.route("/export_attendance_multi", methods=["POST"])
def export_attendance_multi():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    selected_groups = request.form.getlist("groups")
    start_date = request.form.get("start_date")
    end_date = request.form.get("end_date")

    if not selected_groups:
        return "No groups selected", 400

    if not start_date or not end_date:
        return "Date range is required", 400

    conn = get_db()

    file_path = "multi_group_attendance_export.xlsx"

    role = get_user_role(username)
    with pd.ExcelWriter(file_path, engine='openpyxl') as writer:
        for group in selected_groups:
            days, data = get_attendance_data(group, start_date, end_date, teacher_username=username if role == 'teacher' else None)

            # Filter days to the selected date range
            filtered_days = [d for d in days if start_date <= d <= end_date]

            rows = []

            for row in data:
                base = {
                    "Username": row["username"],
                    "Name": row["name"],
                    "Group": row["group"],
                    "Attendance %": row["attendance_pct"]
                }

                for d in filtered_days:
                    val = row["days"].get(d)
                    base[d] = val["time"] if val else "A"

                rows.append(base)

            df = pd.DataFrame(rows)
            # Clean sheet name (remove special characters)
            sheet_name = group.replace("/", "_").replace("\\", "_")[:31]  # Excel sheet name limit
            df.to_excel(writer, sheet_name=sheet_name, index=False)

    log_activity(username, "exported attendance by group")

    conn.close()

    response = send_file(file_path, as_attachment=True)
    response.headers["HX-Redirect"] = url_for("teacher_dashboard")
    return response

@app.route("/reset_attendance", methods=["POST"])
def reset_attendance():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    group = request.form.get("group")
    date = request.form.get("date")

    if not group or not date:
        return "Missing data", 400

    conn = get_db()
    cursor = conn.cursor()

    # 🔹 Get users in group
    cursor.execute("""
    SELECT username FROM users WHERE group_name = ? AND role = 'student'
    """, (group,))
    users = [u[0] for u in cursor.fetchall()]

    if users:
        placeholders = ",".join(["?"] * len(users))

        # 🔹 Delete login history (auto attendance)
        cursor.execute(f"""
        DELETE FROM login_history
        WHERE username IN ({placeholders}) AND date = ?
        """, (*users, date))

        # 🔹 Delete manual overrides
        cursor.execute(f"""
        DELETE FROM attendance_override
        WHERE username IN ({placeholders}) AND date = ?
        """, (*users, date))

        for user in users:
            add_learner_note_entry(cursor, user,
                f"Attendance reset for {date} in group {group}.", username)

    conn.commit()
    conn.close()

    return redirect(url_for("attendance", group=group))

@app.route("/mark_all_present", methods=["POST"])
def mark_all_present():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    group = request.form.get("group")
    date = request.form.get("date")

    if not group or not date:
        return "Missing data", 400

    conn = get_db()
    cursor = conn.cursor()

    # 🔹 Get users in group
    cursor.execute("""
    SELECT username FROM users WHERE group_name = ? AND role = 'student'
    """, (group,))
    users = [u[0] for u in cursor.fetchall()]

    if users:
        # 🔹 Mark all as present (insert/update overrides)
        for user in users:
            cursor.execute("SELECT status FROM attendance_override WHERE username = ? AND date = ?", (user, date))
            previous = cursor.fetchone()

            cursor.execute("""
            INSERT INTO attendance_override (username, date, status)
            VALUES (?, ?, ?)
            ON CONFLICT(username, date)
            DO UPDATE SET status = excluded.status
            """, (user, date, "present"))

            if previous is None or previous[0] != "present":
                add_learner_note_entry(cursor, user,
                    f"Marked present for {date} in group {group}.", username)

    conn.commit()
    log_activity(username, f"marked all in {group} present on {date}")
    conn.close()

    return redirect(url_for("attendance", group=group))

@app.route("/save_attendance", methods=["POST"])
def save_attendance():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    group = request.form.get("group")

    conn = get_db()
    cursor = conn.cursor()

    for key, value in request.form.items():
        if not key.startswith("att_"):
            continue

        # FIX: Safe split (prevents unpack error)
        try:
            _, data = key.split("att_")
            user, day = data.split("|")
        except ValueError:
            continue

        value = value.strip().lower()

        # 🔹 SKIP EMPTY INPUT → DO NOTHING
        if value == "":
            continue

        # 🔹 CHECK if login already exists
        cursor.execute("""
        SELECT MIN(login_time)
        FROM login_history
        WHERE username = ? AND date = ?
        """, (user, day))

        login = cursor.fetchone()[0]

        # 🔹 Determine status
        if value == "x":
            status = "absent"
        else:
            status = "present"

        # 🔹 Insert override ONLY if needed
        cursor.execute("SELECT status FROM attendance_override WHERE username = ? AND date = ?", (user, day))
        previous = cursor.fetchone()

        cursor.execute("""
        INSERT INTO attendance_override (username, date, status)
        VALUES (?, ?, ?)
        ON CONFLICT(username, date)
        DO UPDATE SET status = excluded.status
        """, (user, day, status))

        if previous is None or previous[0] != status:
            add_learner_note_entry(cursor, user,
                f"Attendance manually set to {status.upper()} for {day}.", username)

    conn.commit()
    log_activity(username, f"saved attendance for {group}")
    conn.close()

    return redirect(url_for("attendance", group=group))

@app.route("/exclude_date", methods=["POST"])
def exclude_date():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    date = request.form.get("date")
    group = request.form.get("group")  # Can be empty string for global, or specific group
    reason = request.form.get("reason", "")

    if not date:
        return "Missing date", 400

    # Convert empty string to None for global exclusions
    if group == "":
        group = None

    conn = get_db()
    cursor = conn.cursor()

    timestamp = datetime.now().isoformat()
    if group is None:
        cursor.execute("""
        SELECT rowid
        FROM excluded_dates
        WHERE date = ? AND group_name IS NULL
        ORDER BY rowid
        LIMIT 1
        """, (date,))
        existing = cursor.fetchone()
        if existing:
            cursor.execute("""
            UPDATE excluded_dates
            SET reason = ?, created_by = ?, created_at = ?
            WHERE rowid = ?
            """, (reason, username, timestamp, existing[0]))
            cursor.execute("""
            DELETE FROM excluded_dates
            WHERE date = ? AND group_name IS NULL AND rowid != ?
            """, (date, existing[0]))
        else:
            cursor.execute("""
            INSERT INTO excluded_dates (date, group_name, reason, created_by, created_at)
            VALUES (?, NULL, ?, ?, ?)
            """, (date, reason, username, timestamp))
    else:
        cursor.execute("""
        INSERT INTO excluded_dates (date, group_name, reason, created_by, created_at)
        VALUES (?, ?, ?, ?, ?)
        ON CONFLICT(date, group_name)
        DO UPDATE SET reason = excluded.reason, created_by = excluded.created_by, created_at = excluded.created_at
        """, (date, group, reason, username, timestamp))

    conn.commit()
    log_activity(username, f"excluded date {date} for group {group or 'all groups'} ({reason})")
    conn.close()

    return redirect(url_for("attendance", group=request.form.get("selected_group")))

@app.route("/include_date", methods=["POST"])
def include_date():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    date = request.form.get("date")
    group = request.form.get("group")  # Can be empty string for global, or specific group

    if not date:
        return "Missing date", 400

    # Convert empty string to None for global exclusions
    if group == "":
        group = None

    conn = get_db()
    cursor = conn.cursor()

    # If group is specified, only remove group-specific exclusion
    # If group is None, remove global exclusion
    if group:
        cursor.execute("DELETE FROM excluded_dates WHERE date = ? AND group_name = ?", (date, group))
    else:
        cursor.execute("DELETE FROM excluded_dates WHERE date = ? AND group_name IS NULL", (date,))

    conn.commit()
    log_activity(username, f"included date {date} for group {group or 'all groups'}")
    conn.close()

    return redirect(url_for("attendance", group=request.form.get("selected_group")))

@app.route("/mark_all_absent", methods=["POST"])
def mark_all_absent():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    group = request.form.get("group")
    date = request.form.get("date")

    if not group or not date:
        return "Missing data", 400

    conn = get_db()
    cursor = conn.cursor()

    # 🔹 Get users in group
    cursor.execute("""
    SELECT username FROM users WHERE group_name = ? AND role = 'student'
    """, (group,))
    users = [u[0] for u in cursor.fetchall()]

    if users:
        # 🔹 Mark all as absent (insert/update overrides)
        for user in users:
            cursor.execute("SELECT status FROM attendance_override WHERE username = ? AND date = ?", (user, date))
            previous = cursor.fetchone()

            cursor.execute("""
            INSERT INTO attendance_override (username, date, status)
            VALUES (?, ?, ?)
            ON CONFLICT(username, date)
            DO UPDATE SET status = excluded.status
            """, (user, date, "absent"))

            if previous is None or previous[0] != "absent":
                add_learner_note_entry(cursor, user,
                    f"Marked absent for {date} in group {group}.", username)

    conn.commit()
    log_activity(username, f"marked all in {group} absent on {date}")
    conn.close()

    return redirect(url_for("attendance", group=group))

@app.route("/api/excluded_dates", methods=["GET"])
def api_excluded_dates():
    """API endpoint to fetch excluded dates as JSON"""
    username = session.get("username")
    if not username:
        return {"error": "Unauthorized"}, 401

    role = get_user_role(username)
    if role not in ["teacher", "admin"]:
        return {"error": "Access denied"}, 403

    conn = get_db()
    cursor = conn.cursor()

    cursor.execute("""
    SELECT date, group_name, reason, created_by, created_at
    FROM excluded_dates
    ORDER BY date DESC
    """)

    excluded_dates = []
    for row in cursor.fetchall():
        excluded_dates.append({
            "date": row[0],
            "group_name": row[1],
            "reason": row[2],
            "created_by": row[3],
            "created_at": row[4]
        })

    conn.close()

    return {"excluded_dates": excluded_dates}

@app.route("/api/attendance_data", methods=["GET"])
def api_attendance_data():
    """API endpoint to fetch attendance data as JSON"""
    username = session.get("username")
    if not username:
        return {"error": "Unauthorized"}, 401

    role = get_user_role(username)
    if role not in ["teacher", "admin"]:
        return {"error": "Access denied"}, 403

    group = request.args.get("group")
    start_date = request.args.get("start_date")
    end_date = request.args.get("end_date")

    if not group:
        return {"error": "Group required"}, 400

    days, data = get_attendance_data(group, start_date, end_date, teacher_username=username if role == 'teacher' else None)

    return {
        "days": days,
        "data": data,
        "success": True
    }

@app.route("/manage_subjects", methods=["GET", "POST"])
def manage_subjects():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    role = get_user_role(username)
    if role not in ["teacher", "admin"]:
        return "Access denied", 403

    if request.method == "POST":
        action = request.form.get("action")
        if action == "create":
            subject_name = request.form.get("subject_name")
            if subject_name:
                conn = get_db()
                cursor = conn.cursor()
                try:
                    cursor.execute("INSERT INTO subjects (name, created_by, created_at) VALUES (?, ?, ?)",
                                   (subject_name, username, datetime.now().isoformat()))
                    conn.commit()
                    log_activity(username, f"created subject {subject_name}")
                except sqlite3.IntegrityError:
                    pass  # Subject already exists
                conn.close()
        elif action == "delete":
            subject_id = request.form.get("subject_id")
            if subject_id:
                conn = get_db()
                cursor = conn.cursor()
                cursor.execute("SELECT name FROM subjects WHERE id = ?", (subject_id,))
                subj = cursor.fetchone()
                if subj:
                    # Delete all results for tasks in this subject
                    cursor.execute("""
                        DELETE FROM results WHERE subject = ?
                    """, (subj[0],))
                    # Delete task groups and tasks
                    cursor.execute("DELETE FROM task_groups WHERE task_id IN (SELECT id FROM tasks WHERE subject_id = ?)", (subject_id,))
                    cursor.execute("DELETE FROM tasks WHERE subject_id = ?", (subject_id,))
                    cursor.execute("DELETE FROM subjects WHERE id = ?", (subject_id,))
                    conn.commit()
                    log_activity(username, f"deleted subject {subj[0]} and all related tasks and results")
                conn.close()

        # After POST we also render template with fresh subject list
        conn = get_db()
        cursor = conn.cursor()
        cursor.execute("SELECT id, name FROM subjects ORDER BY name")
        subjects = cursor.fetchall()
        conn.close()
        return render_template("manage_subjects.html", subjects=subjects)


    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("SELECT id, name FROM subjects ORDER BY name")
    subjects = cursor.fetchall()
    conn.close()

    return render_template("manage_subjects.html", subjects=subjects)


@app.route("/manage_tasks/<subject_id>", methods=["GET", "POST"])

def manage_tasks(subject_id):
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    role = get_user_role(username)
    if role not in ["teacher", "admin"]:
        return "Access denied", 403

    conn = get_db()
    cursor = conn.cursor()

    cursor.execute("SELECT name FROM subjects WHERE id = ?", (subject_id,))
    subj = cursor.fetchone()
    if not subj:
        conn.close()
        return "Subject not found", 404
    subject_name = subj[0]

    if request.method == "POST":
        action = request.form.get("action")
        if action == "create":
            task_name = request.form.get("task_name")
            assign_date = request.form.get("assign_date")
            marking_script = request.form.get("marking_script")
            marking_setup_id = request.form.get("marking_setup_id") or None
            allow_multiple = 1 if request.form.get("allow_multiple") else 0
            max_attempts = int(request.form.get("max_attempts", 1)) if allow_multiple else 1
            groups = request.form.getlist("groups")
            teachers = request.form.getlist("teachers")
            if task_name and assign_date:
                question_text = request.form.get("question_text", "").strip()
                cursor.execute(
                    "INSERT INTO tasks (subject_id, name, assign_date, marking_script, marking_setup_id, "
                    "question_text, task_type, allow_multiple, max_attempts, created_by, created_at) "
                    "VALUES (?, ?, ?, ?, ?, ?, 'practical', ?, ?, ?, ?)",
                    (subject_id, task_name, assign_date, marking_script, marking_setup_id, question_text,
                     allow_multiple, max_attempts, username, datetime.now().isoformat()),
                )
                task_id = cursor.lastrowid
                # Handle sample file upload
                if 'sample_file' in request.files:
                    file = request.files['sample_file']
                    if file.filename:
                        file_content = file.read()
                        file_name = file.filename
                        cursor.execute("UPDATE tasks SET sample_file = ?, sample_file_name = ? WHERE id = ?", (file_content, file_name, task_id))
                for group in groups:
                    cursor.execute("INSERT INTO task_groups (task_id, group_name) VALUES (?, ?)", (task_id, group))
                for teacher in teachers:
                    if teacher.strip():
                        cursor.execute("INSERT INTO task_teachers (task_id, teacher_username) VALUES (?, ?)", (task_id, teacher))
                conn.commit()
                log_activity(username, f"created task {task_name} in {subject_name}")
        elif action == "reuse":
            source_task_id = request.form.get("source_task_id")
            new_task_name = request.form.get("task_name", "").strip()
            new_assign_date = request.form.get("assign_date")
            new_marking_script = request.form.get("marking_script") or None
            new_marking_setup_id = request.form.get("marking_setup_id") or None
            new_allow_multiple = 1 if request.form.get("allow_multiple") else 0
            new_max_attempts = int(request.form.get("max_attempts", 1)) if new_allow_multiple else 1
            new_groups = request.form.getlist("groups")
            new_teachers = request.form.getlist("teachers")
            if source_task_id and new_task_name and new_assign_date:
                cursor.execute(
                    "SELECT question_text, sample_file, sample_file_name FROM tasks WHERE id = ?",
                    (source_task_id,)
                )
                src = cursor.fetchone()
                src_question = src[0] if src else ""
                src_file = src[1] if src else None
                src_file_name = src[2] if src else None
                cursor.execute(
                    "INSERT INTO tasks (subject_id, name, assign_date, marking_script, marking_setup_id, "
                    "question_text, task_type, allow_multiple, max_attempts, sample_file, sample_file_name, "
                    "created_by, created_at) VALUES (?, ?, ?, ?, ?, ?, 'practical', ?, ?, ?, ?, ?, ?)",
                    (subject_id, new_task_name, new_assign_date, new_marking_script, new_marking_setup_id, src_question,
                     new_allow_multiple, new_max_attempts, src_file, src_file_name,
                     username, datetime.now().isoformat())
                )
                new_task_id = cursor.lastrowid
                for g in new_groups:
                    if g.strip():
                        cursor.execute("INSERT INTO task_groups (task_id, group_name) VALUES (?, ?)", (new_task_id, g))
                for t in new_teachers:
                    if t.strip():
                        cursor.execute("INSERT INTO task_teachers (task_id, teacher_username) VALUES (?, ?)", (new_task_id, t))
                conn.commit()
                log_activity(username, f"reused task {source_task_id} as '{new_task_name}' in {subject_name}")
        elif action == "assign_theory":
            task_name = request.form.get("task_name")
            assign_date = request.form.get("assign_date")
            theory_test_id = request.form.get("theory_test_id")
            groups = request.form.getlist("groups")
            teachers = request.form.getlist("teachers")
            if task_name and assign_date and theory_test_id:
                cursor.execute("INSERT INTO tasks (subject_id, name, assign_date, theory_test_id, task_type, created_by, created_at) VALUES (?, ?, ?, ?, 'theory', ?, ?)",
                               (subject_id, task_name, assign_date, theory_test_id, username, datetime.now().isoformat()))
                task_id = cursor.lastrowid
                for group in groups:
                    cursor.execute("INSERT INTO task_groups (task_id, group_name) VALUES (?, ?)", (task_id, group))
                for teacher in teachers:
                    if teacher.strip():
                        cursor.execute("INSERT INTO task_teachers (task_id, teacher_username) VALUES (?, ?)", (task_id, teacher))
                conn.commit()
                log_activity(username, f"assigned theory test as task {task_name} in {subject_name}")
        elif action == "delete":
            task_id = request.form.get("task_id")
            if task_id:
                cursor.execute("SELECT name, subject_id FROM tasks WHERE id = ?", (task_id,))
                tsk = cursor.fetchone()
                if tsk:
                    task_name = tsk[0]
                    subject_id_val = tsk[1]
                    # Get subject name to delete results
                    cursor.execute("SELECT name FROM subjects WHERE id = ?", (subject_id_val,))
                    subj_row = cursor.fetchone()
                    if subj_row:
                        subject_name = subj_row[0]
                        # Delete all results for this task
                        cursor.execute("""
                            DELETE FROM results WHERE subject = ? AND task = ?
                        """, (subject_name, task_name))
                    # Delete task groups and task
                    cursor.execute("DELETE FROM task_groups WHERE task_id = ?", (task_id,))
                    cursor.execute("DELETE FROM tasks WHERE id = ?", (task_id,))
                    conn.commit()
                    log_activity(username, f"deleted task {task_name} from {subject_name} and all related results")

        conn.close()
        return redirect(url_for("manage_tasks", subject_id=subject_id))


    # Get all groups — teachers can assign tasks to any group
    cursor.execute("SELECT DISTINCT group_name FROM users WHERE group_name IS NOT NULL ORDER BY group_name")
    all_groups = [row[0] for row in cursor.fetchall()]

    # Get available marking scripts
    available_scripts = get_marking_scripts()
    marking_conn = get_marking_db()
    marking_cursor = marking_conn.cursor()
    marking_cursor.execute("SELECT id, title FROM marking_setups ORDER BY title, id")
    available_marking_setups = marking_cursor.fetchall()
    marking_conn.close()

    # Get available theory tests
    cursor.execute("SELECT id, title, subject FROM theory_tests ORDER BY title")
    available_theory_tests = cursor.fetchall()

    teachers = get_teachers()
    teacher_checkboxes = ''.join(
        f'<label style="display:inline-flex;align-items:center;gap:5px;">'
        f'<input type="checkbox" name="teachers" value="{escape(t[0])}"> {escape(t[1] or t[0])}</label>'
        for t in teachers
    )

    # Get tasks — all tasks visible, teachers can manage any
    cursor.execute("""
    SELECT t.id, t.name, t.assign_date, t.marking_script, t.task_type, t.theory_test_id,
           t.allow_multiple, t.max_attempts, t.is_active, t.question_text, t.sample_file_name,
           t.marking_setup_id,
           GROUP_CONCAT(DISTINCT tg.group_name),
           GROUP_CONCAT(DISTINCT tt.teacher_username)
    FROM tasks t
    LEFT JOIN task_groups tg ON t.id = tg.task_id
    LEFT JOIN task_teachers tt ON t.id = tt.task_id
    WHERE t.subject_id = ?
    GROUP BY t.id
    ORDER BY t.assign_date, t.name
    """, (subject_id,))
    tasks = cursor.fetchall()
    conn.close()

    task_list = ""
    for task_id, task_name, assign_date, marking_script, task_type, theory_test_id, allow_multiple, max_attempts, is_active, question_text, sample_file_name, marking_setup_id, group_list, teacher_list in tasks:
        if task_type == "theory":
            type_label = '<span style="background:#0078D4;color:white;padding:2px 6px;border-radius:10px;font-size:0.8em;">📝 Theory</span>'
            script_label = f'Test ID: {theory_test_id}'
        else:
            type_label = '<span style="background:#107C10;color:white;padding:2px 6px;border-radius:10px;font-size:0.8em;">📁 Practical</span>'
            script_label = marking_script if marking_script else '<span style="color:red;">None assigned</span>'
            if marking_setup_id:
                setup_name = next((title for setup_id, title in available_marking_setups if setup_id == marking_setup_id), None)
                setup_label = escape(setup_name or f"Setup {marking_setup_id}")
                script_label += f'<br><span style="color:#005a9e;">Setup: {setup_label}</span>'

        status_badge = '<span style="background:#c8f7c5;color:#107C10;padding:2px 8px;border-radius:10px;font-size:0.8em;">Active</span>' if is_active else '<span style="background:#f7c5c5;color:#A4262C;padding:2px 8px;border-radius:10px;font-size:0.8em;">Inactive</span>'
        toggle_label = '⏸ Deactivate' if is_active else '▶ Activate'
        toggle_style = 'background:#ff8c00;color:white;' if is_active else 'background:#107C10;color:white;'

        if task_type == 'practical':
            attempts_label = 'Single' if not allow_multiple else f'Multiple ({max_attempts})'
        else:
            attempts_label = 'Theory task'

        task_list += f"""
        <tr>
            <td>{escape(task_name)} {type_label}</td>
            <td>{assign_date}</td>
            <td>{group_list or 'None'}</td>
            <td>{teacher_list or 'None'}</td>
            <td>{script_label}</td>
            <td>{f'<a href="/tasks/{task_id}/sample_file" target="_blank">{escape(sample_file_name)}</a>' if sample_file_name else 'None'}</td>
            <td>{attempts_label}</td>
            <td>{status_badge}</td>
            <td style="white-space:nowrap; vertical-align:middle;">
                {'<a href="/tasks/' + str(task_id) + '/edit" title="Edit task" class="btn btn-primary">✏️</a>' if task_type == 'practical' else ''}
                <a href="/tasks/{task_id}/preview" class="icon-btn" title="Preview learner view">👁</a>
                {'<button type="button" class="btn btn-success" title="Reuse: copy into a new task" onclick="openReuseTaskModal(' + str(task_id) + ', ' + repr(task_name) + ', ' + repr(marking_script or '') + ', ' + str(allow_multiple) + ', ' + str(max_attempts) + ', ' + repr(marking_setup_id or '') + ')">📋</button>' if task_type == 'practical' else ''}
                <form method="post" action="/tasks/{task_id}/toggle" style="display:inline-flex; margin:0;">
                    <input type="hidden" name="subject_id" value="{subject_id}">
                    <button type="submit" title="{toggle_label}" class="btn" style="{toggle_style}">{'⏸' if is_active else '▶'}</button>
                </form>
                <form method="post" action="/tasks/{task_id}/clear_uploads" style="display:inline-flex; margin:0;"
                      onsubmit="return confirm('Clear ALL uploads for {escape(task_name)}? This cannot be undone.')">
                    <input type="hidden" name="subject_id" value="{subject_id}">
                    <button type="submit" title="Clear uploads" class="btn btn-danger">🗑</button>
                </form>
                <form method="post" style="display:inline-flex; margin:0;" onsubmit="return confirm('⚠️ WARNING: Delete task {escape(task_name)} and ALL STUDENT SUBMISSIONS?\n\nThis will permanently remove:\n- All student uploads and scores\n- Task from Group Results\n\nThis action CANNOT be undone!')">
                    <input type="hidden" name="action" value="delete">
                    <input type="hidden" name="task_id" value="{task_id}">
                    <button type="submit" title="Delete task" class="btn" style="background:#555;color:white;">🗑</button>
                </form>
            </td>
        </tr>
        """

    group_checkboxes = ""
    for group in all_groups:
        group_checkboxes += f'<label style="display:inline-flex;align-items:center;gap:5px;margin-right:10px;"><input type="checkbox" name="groups" value="{escape(group)}"> {escape(group)}</label>'

    script_options = '<option value="">-- No marking script --</option>'
    for script in available_scripts:
        script_options += f'<option value="{escape(script)}">{escape(script)}</option>'

    marking_setup_options = '<option value="">-- No marking setup --</option>'
    for setup_id, setup_title in available_marking_setups:
        marking_setup_options += f'<option value="{setup_id}">{escape(setup_title)}</option>'

    theory_test_options = '<option value="">-- Select Theory Test --</option>'
    for tt_id, tt_title, tt_subject in available_theory_tests:
        label = f"{tt_title}" + (f" ({tt_subject})" if tt_subject else "")
        theory_test_options += f'<option value="{tt_id}">{escape(label)}</option>'

    return render_template(
        "manage_tasks.html",
        subject_name=subject_name,
        script_options=script_options,
        marking_setup_options=marking_setup_options,
        group_checkboxes=group_checkboxes,
        teacher_checkboxes=teacher_checkboxes,
        task_list=task_list,
    )


@app.route("/tasks/<int:task_id>/preview")
def task_preview(task_id):
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    role = get_user_role(username)
    if role not in ["teacher", "admin"]:
        return "Access denied", 403

    conn = get_db()
    cursor = conn.cursor()
    cursor.execute(
        "SELECT t.name, t.assign_date, t.marking_script, t.question_text, t.task_type, t.theory_test_id, t.sample_file_name, s.name "
        "FROM tasks t JOIN subjects s ON t.subject_id = s.id WHERE t.id = ?",
        (task_id,),
    )
    task_row = cursor.fetchone()
    if not task_row:
        conn.close()
        return "Task not found", 404

    task_name, assign_date, marking_script, question_text, task_type, theory_test_id, sample_file_name, subject_name = task_row
    theory_test_title = None
    if task_type == "theory" and theory_test_id:
        cursor.execute("SELECT title FROM theory_tests WHERE id = ?", (theory_test_id,))
        test_row = cursor.fetchone()
        theory_test_title = test_row[0] if test_row else None

    conn.close()
    return render_template(
        "task_preview.html",
        subject_name=subject_name,
        task_name=task_name,
        assign_date=assign_date,
        task_type=task_type,
        marking_script=marking_script,
        question_text=question_text,
        sample_file_name=sample_file_name,
        sample_url=f"/tasks/{task_id}/sample_file" if sample_file_name else None,
        theory_test_title=theory_test_title,
        theory_test_id=theory_test_id,
    )


@app.route("/tasks/<int:task_id>/sample_file")
def download_sample_file(task_id):
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("SELECT sample_file, sample_file_name FROM tasks WHERE id = ?", (task_id,))
    row = cursor.fetchone()
    conn.close()

    if not row or not row[0]:
        return "Sample file not found", 404

    file_content, file_name = row
    return send_file(
        io.BytesIO(file_content),
        as_attachment=True,
        download_name=file_name,
        mimetype='application/octet-stream'
    )


@app.route("/interactive_learning_files/<path:relative_path>")
def download_interactive_learning_file(relative_path):
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    file_path = resolve_interactive_learning_path(relative_path)
    if not file_path:
        return "File not found", 404

    return send_file(file_path, as_attachment=True, download_name=os.path.basename(file_path))


@app.route("/manage_lessons")
def manage_lessons():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    role = get_user_role(username)
    if role not in ["teacher", "admin"]:
        return "Access denied", 403

    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("""
        SELECT
            t.id, t.title, t.subject, t.assign_date, t.time_limit, t.is_active,
            GROUP_CONCAT(DISTINCT tg.group_name),
            COUNT(DISTINCT q.id),
            t.allow_multiple, t.max_attempts, t.show_answers,
            GROUP_CONCAT(DISTINCT tt.teacher_username),
            SUM(CASE WHEN q.question_type IN ('content_slide', 'title_slide', 'heading_slide') THEN 1 ELSE 0 END) as content_count
        FROM theory_tests t
        LEFT JOIN theory_questions q ON t.id = q.test_id
        LEFT JOIN theory_test_groups tg ON t.id = tg.test_id
        LEFT JOIN theory_test_teachers tt ON t.id = tt.test_id
        GROUP BY t.id
        HAVING COALESCE(content_count, 0) > 0
        ORDER BY t.created_at DESC
    """)
    lessons = cursor.fetchall()
    groups = get_groups(username) if role == "teacher" else get_groups()
    teachers = get_teachers()
    teacher_checkboxes = ''.join(
        f'<label style="font-weight:normal;display:inline-flex;align-items:center;gap:5px;">'
        f'<input type="checkbox" name="teachers" value="{escape(t[0])}"> {escape(t[1] or t[0])}</label>'
        for t in teachers
    )
    lesson_list = ""
    for lesson in lessons:
        lesson_id = lesson[0]
        lesson_title = escape(lesson[1] or "")
        lesson_subject = escape(lesson[2] or "")
        assign_date = lesson[3] or '—'
        time_limit_val = lesson[4] or 0
        is_active = bool(lesson[5])
        groups_text = escape(lesson[6] or 'All Groups')
        total_items = lesson[7] or 0
        show_answers = bool(lesson[10])
        teachers_text = escape(lesson[11] or 'All Teachers')
        content_count = lesson[12] or 0
        test_question_count = max(0, total_items - content_count)
        status_badge = '<span class="badge-active">Active</span>' if is_active else '<span class="badge-inactive">Inactive</span>'
        toggle_label = 'Deactivate' if is_active else 'Activate'
        toggle_class = 'btn-warning' if is_active else 'btn-success'
        lesson_list += f"""
        <tr>
            <td>{lesson_title}</td>
            <td>{lesson_subject or '—'}</td>
            <td>{groups_text}</td>
            <td>{teachers_text}</td>
            <td>{content_count}</td>
            <td>{test_question_count}</td>
            <td>{assign_date}</td>
            <td>{time_limit_val if time_limit_val else 'No limit'}</td>
            <td>{'✔ Yes' if show_answers else '✘ No'}</td>
            <td>{status_badge}</td>
            <td style="white-space:nowrap; vertical-align:middle;">
                <div class="action-cell">
                <a href="/manage_lessons/{lesson_id}/questions" class="btn btn-primary" title="Edit lesson">✏️</a>
                <a href="/manage_tests/{lesson_id}/edit" class="btn btn-warning" title="Edit settings">⚙️</a>
                <form method="post" action="/manage_lessons/{lesson_id}/toggle" style="display:inline-flex; margin:0;">
                    <button type="submit" class="btn {toggle_class}" title="{toggle_label}">{'⏸' if is_active else '▶'}</button>
                </form>
                <form method="post" action="/manage_lessons/{lesson_id}/delete" style="display:inline-flex; margin:0;"
                      onsubmit="return confirm('Delete this lesson setup, its questions, and all learner submissions?')">
                    <button type="submit" class="btn btn-danger" title="Delete lesson">🗑</button>
                </form>
                </div>
            </td>
        </tr>
        """
    conn.close()

    return render_template(
        "manage_lessons.html",
        groups=groups,
        teacher_checkboxes=teacher_checkboxes,
        test_list=lesson_list,
        page_title="Lesson Setup",
        page_intro="Create and manage slide-based lessons here. Theory Tests remain the plain question-only tests.",
    )


@app.route("/manage_lesson_tests")
def manage_lesson_tests():
    return redirect(url_for("manage_lessons"))


@app.route("/manage_lessons/create", methods=["POST"])
def create_lesson():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    title = request.form.get("title", "").strip()
    subject = request.form.get("subject", "").strip()
    assign_date = request.form.get("assign_date", "").strip()
    time_limit = safe_int(request.form.get("time_limit"), 0)
    allow_multiple = 1 if request.form.get("allow_multiple") else 0
    max_attempts = safe_int(request.form.get("max_attempts"), 1)
    show_answers = 1 if request.form.get("show_answers") else 0
    groups = request.form.getlist("groups")
    teachers = request.form.getlist("teachers")

    if not title or not assign_date:
        return "Title and assign date are required", 400

    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("""
        INSERT INTO theory_tests
            (title, subject, assign_date, time_limit, allow_multiple, max_attempts, show_answers, created_by, created_at)
        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
    """, (title, subject, assign_date, time_limit, allow_multiple, max_attempts, show_answers, username, datetime.now().isoformat()))
    lesson_id = cursor.lastrowid

    for group_name in groups:
        if group_name.strip():
            cursor.execute("INSERT INTO theory_test_groups (test_id, group_name) VALUES (?, ?)", (lesson_id, group_name))
    for teacher_username in teachers:
        if teacher_username.strip():
            cursor.execute("INSERT INTO theory_test_teachers (test_id, teacher_username) VALUES (?, ?)", (lesson_id, teacher_username))

    conn.commit()
    conn.close()
    log_activity(username, f"created theory lesson '{title}'")
    return redirect(url_for("manage_lesson_questions", test_id=lesson_id))


@app.route("/manage_lessons/<int:lesson_id>/toggle", methods=["POST"])
def toggle_lesson(lesson_id):
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("SELECT is_active FROM theory_tests WHERE id = ?", (lesson_id,))
    row = cursor.fetchone()
    if row:
        cursor.execute("UPDATE theory_tests SET is_active = ? WHERE id = ?", (0 if row[0] else 1, lesson_id))
        conn.commit()
    conn.close()
    return redirect(url_for("manage_lessons"))


@app.route("/manage_lessons/<int:lesson_id>/delete", methods=["POST"])
def delete_lesson(lesson_id):
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("SELECT title FROM theory_tests WHERE id = ?", (lesson_id,))
    row = cursor.fetchone()
    lesson_title = row[0] if row else f"Lesson {lesson_id}"
    cursor.execute("DELETE FROM theory_answers WHERE submission_id IN (SELECT id FROM theory_submissions WHERE test_id = ?)", (lesson_id,))
    cursor.execute("DELETE FROM theory_submissions WHERE test_id = ?", (lesson_id,))
    cursor.execute("DELETE FROM theory_progress WHERE test_id = ?", (lesson_id,))
    cursor.execute("DELETE FROM theory_options WHERE question_id IN (SELECT id FROM theory_questions WHERE test_id = ?)", (lesson_id,))
    cursor.execute("DELETE FROM theory_questions WHERE test_id = ?", (lesson_id,))
    cursor.execute("DELETE FROM theory_test_groups WHERE test_id = ?", (lesson_id,))
    cursor.execute("DELETE FROM theory_test_teachers WHERE test_id = ?", (lesson_id,))
    cursor.execute("DELETE FROM theory_tests WHERE id = ?", (lesson_id,))
    conn.commit()
    conn.close()
    log_activity(username, f"deleted lesson setup '{lesson_title}'")
    return redirect(url_for("manage_lessons"))


@app.route("/manage_lessons/<int:lesson_id>/checkpoints", methods=["GET", "POST"])
def manage_lesson_checkpoints(lesson_id):
    return redirect(url_for("manage_lesson_questions", test_id=lesson_id))


@app.route("/lessons")
def learner_lessons():
    return redirect(url_for("lesson_tests"))


@app.route("/lesson/<int:lesson_id>", methods=["GET", "POST"])
def lesson_view(lesson_id):
    return redirect(url_for("take_test", test_id=lesson_id), code=307 if request.method == "POST" else 302)


# ── Theory Test Routes ───────────────────────────────────────────────────────

@app.route("/manage_tests")
def manage_tests():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    role = get_user_role(username)
    if role not in ["teacher", "admin"]:
        return "Access denied", 403

    conn = get_db()
    cursor = conn.cursor()
    # All teachers see all tests and can assign to any group
    cursor.execute("""
        SELECT
            t.id, t.title, t.subject, t.assign_date, t.time_limit, t.is_active,
            GROUP_CONCAT(DISTINCT tg.group_name),
            COUNT(DISTINCT q.id),
            t.allow_multiple, t.max_attempts, t.show_answers,
            GROUP_CONCAT(DISTINCT tt.teacher_username),
            SUM(CASE WHEN q.question_type IN ('content_slide', 'title_slide', 'heading_slide') THEN 1 ELSE 0 END) as content_count
        FROM theory_tests t
        LEFT JOIN theory_questions q ON t.id = q.test_id
        LEFT JOIN theory_test_groups tg ON t.id = tg.test_id
        LEFT JOIN theory_test_teachers tt ON t.id = tt.test_id
        GROUP BY t.id
        HAVING COALESCE(content_count, 0) = 0
        ORDER BY t.created_at DESC
    """)
    tests = cursor.fetchall()
    groups = get_groups(username) if role == 'teacher' else get_groups()
    teachers = get_teachers()
    teacher_checkboxes = ''.join(
        f'<label style="font-weight:normal;display:inline-flex;align-items:center;gap:5px;">'
        f'<input type="checkbox" name="teachers" value="{escape(t[0])}"> {escape(t[1] or t[0])}</label>'
        for t in teachers
    )

    test_list = ""
    for test in tests:
        test_id = test[0]
        test_title = escape(test[1] or "")
        test_subject = escape(test[2] or "")
        assign_date = test[3] or '—'
        time_limit_val = test[4] or 0
        is_active = bool(test[5])
        groups_text = escape(test[6] or 'All Groups')
        question_count = test[7] or 0
        allow_multiple = bool(test[8])
        max_attempts = test[9] or 1
        show_answers = bool(test[10])
        teachers_text = escape(test[11] or 'All Teachers')
        category_badge = '<span class="badge-inactive" style="background:#e3f2fd;color:#0078D4;">Test</span>'

        status_badge = '<span class="badge-active">Active</span>' if is_active else '<span class="badge-inactive">Inactive</span>'
        attempt_text = f'{max_attempts} max' if allow_multiple else '1 (single)'
        toggle_label = 'Deactivate' if is_active else 'Activate'
        toggle_class = 'btn-warning' if is_active else 'btn-success'

        test_list += f"""
        <tr>
            <td>{test_title}</td>
            <td>{test_subject or '—'}</td>
            <td>{category_badge}</td>
            <td>{groups_text}</td>
            <td>{teachers_text}</td>
            <td>{question_count}</td>
            <td>{assign_date}</td>
            <td>{time_limit_val if time_limit_val else 'No limit'}</td>
            <td>{attempt_text}</td>
            <td>{'✔ Yes' if show_answers else '✘ No'}</td>
            <td>{status_badge}</td>
            <td style="white-space:nowrap; vertical-align:middle;">
                <div class="action-cell">
                <a href="/manage_tests/{test_id}/questions" class="btn btn-primary" title="Edit questions">✏️</a>
                <a href="/manage_tests/{test_id}/edit" class="btn btn-warning" title="Edit settings">⚙️</a>
                <button type="button" class="btn btn-success" title="Reuse: copy questions into a new test"
                    onclick="openReuseModal(this)"
                    data-test-id="{test_id}"
                    data-title="{test_title}"
                    data-subject="{test_subject}"
                    data-time-limit="{time_limit_val}"
                    data-allow-multiple="{1 if allow_multiple else 0}"
                    data-max-attempts="{max_attempts}"
                    data-show-answers="{1 if show_answers else 0}">📋</button>
                <form method="post" action="/manage_tests/{test_id}/toggle" style="display:inline-flex; margin:0;">
                    <button type="submit" class="btn {toggle_class}" title="{toggle_label}">{'⏸' if is_active else '▶'}</button>
                </form>
                <form method="post" action="/manage_tests/{test_id}/delete" style="display:inline-flex; margin:0;"
                      onsubmit="return confirm('⚠️ WARNING: Delete this test, all questions, AND ALL STUDENT SUBMISSIONS?\n\nThis will permanently remove:\n- All test questions\n- All student attempts and scores\n- Test from Group Results\n\nThis action CANNOT be undone!')">
                    <button type="submit" class="btn btn-danger" title="Delete test">🗑</button>
                </form>
                </div>
            </td>
        </tr>
        """

    conn.close()
    return render_template(
        "manage_tests.html",
        tests=tests,
        groups=groups,
        teacher_checkboxes=teacher_checkboxes,
        test_list=test_list,
        linked_test_count=0,
        standalone_test_count=len(tests),
        page_title="Theory Tests",
        page_intro="This page manages the simple question-only part of Theory.",
    )


@app.route("/question_bank", methods=["GET", "POST"])
def question_bank():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    conn = get_db()
    cursor = conn.cursor()
    message = ""
    error = ""

    if request.method == "POST":
        action = request.form.get("action", "").strip()
        if action == "delete_question":
            question_id = request.form.get("question_id", type=int)
            if question_id:
                cursor.execute("DELETE FROM question_bank_options WHERE bank_question_id = ?", (question_id,))
                cursor.execute("DELETE FROM question_bank_questions WHERE id = ?", (question_id,))
                conn.commit()
                message = "Question removed from bank."
        elif action == "bulk_delete_questions":
            question_ids = [safe_int(item, 0) for item in request.form.getlist("question_ids")]
            question_ids = [item for item in question_ids if item > 0]
            if question_ids:
                placeholders = ",".join("?" for _ in question_ids)
                cursor.execute(f"DELETE FROM question_bank_options WHERE bank_question_id IN ({placeholders})", question_ids)
                cursor.execute(f"DELETE FROM question_bank_questions WHERE id IN ({placeholders})", question_ids)
                conn.commit()
                message = f"Removed {len(question_ids)} bank question(s)."
        elif action == "import_bank_csv":
            upload = request.files.get("questions_csv")
            try:
                if not upload or not upload.filename:
                    raise ValueError("Please choose a CSV file.")
                content = upload.stream.read().decode("utf-8-sig")
                reader = csv.DictReader(io.StringIO(content))
                imported = 0
                skipped = 0
                for item in reader:
                    q_text = (item.get("question_text") or item.get("question") or "").strip()
                    q_type = (item.get("question_type") or "").strip().lower()
                    if not q_text or q_type not in QUESTION_BANK_SUPPORTED_TYPES:
                        continue
                    marks = safe_int(item.get("marks"), 1)
                    if q_type == "mcq_single":
                        marks = 1
                    subject = (item.get("subject") or "").strip()
                    modules = ", ".join(parse_module_names(item.get("modules") or item.get("module") or ""))
                    option_payload = []
                    if q_type == "mcq_single":
                        option_texts = []
                        for idx in range(1, 7):
                            option_text = (item.get(f"option_{idx}") or "").strip()
                            if option_text:
                                option_texts.append(option_text)
                        correct_raw = (item.get("correct_answer") or "").strip()
                        correct_index = safe_int(correct_raw, 0) - 1 if correct_raw.isdigit() else -1
                        normalized_correct = correct_raw.lower()
                        for idx, option_text in enumerate(option_texts):
                            option_payload.append((option_text, 1 if idx == correct_index or option_text.lower() == normalized_correct else 0, None))
                    elif q_type == "true_false":
                        correct_answer = (item.get("correct_answer") or "True").strip().title()
                        option_payload.extend([
                            ("True", 1 if correct_answer == "True" else 0, None),
                            ("False", 1 if correct_answer == "False" else 0, None),
                        ])
                        correction = (item.get("correction") or "").strip()
                        if correction:
                            option_payload.append((correction, 0, "correction"))
                    elif q_type == "fill_in":
                        answer = (item.get("answer") or item.get("correct_answer") or "").strip()
                        if answer:
                            option_payload.append((answer, 1, None))
                    elif q_type == "match":
                        left = ""
                        right = ""
                        for idx in range(1, 7):
                            candidate_left = (item.get(f"match_a_{idx}") or "").strip()
                            candidate_right = (item.get(f"match_b_{idx}") or "").strip()
                            if candidate_left and candidate_right:
                                left = candidate_left
                                right = candidate_right
                                break
                        if left and right:
                            option_payload.append((left, 1, right))
                    if bank_question_exists(cursor, q_text, q_type, subject, modules, option_payload):
                        skipped += 1
                        continue
                    cursor.execute("""
                        INSERT INTO question_bank_questions (question_text, question_type, marks, subject, modules, created_by, created_at)
                        VALUES (?, ?, ?, ?, ?, ?, ?)
                    """, (q_text, q_type, marks, subject, modules, username, datetime.now().isoformat()))
                    bank_question_id = cursor.lastrowid
                    for option_text, is_correct, match_pair in option_payload:
                        cursor.execute("""
                            INSERT INTO question_bank_options (bank_question_id, option_text, is_correct, match_pair)
                            VALUES (?, ?, ?, ?)
                        """, (bank_question_id, option_text, is_correct, match_pair))
                    imported += 1
                conn.commit()
                message = f"Imported {imported} bank question(s)."
                if skipped:
                    message += f" Skipped {skipped} duplicate(s)."
            except Exception as exc:
                error = f"Import failed: {exc}"
        elif action == "add_question":
            q_text = (request.form.get("question_text") or "").strip()
            q_type = (request.form.get("question_type") or "").strip()
            marks = safe_int(request.form.get("marks"), 1)
            if q_type == "mcq_single":
                marks = 1
            subject = (request.form.get("subject") or "").strip()
            modules = ", ".join(parse_module_names(request.form.get("modules") or ""))
            if not q_text or q_type not in QUESTION_BANK_SUPPORTED_TYPES:
                error = "Question text and a supported type are required."
            else:
                option_payload = []
                if q_type == "mcq_single":
                    option_texts = request.form.getlist("option_text")
                    correct_index = safe_int(request.form.get("is_correct"), -1)
                    for idx, option_text in enumerate(option_texts):
                        option_text = option_text.strip()
                        if option_text:
                            option_payload.append((option_text, 1 if idx == correct_index else 0, None))
                elif q_type == "true_false":
                    tf_correct = (request.form.get("tf_correct") or "True").strip()
                    correction = (request.form.get("correction_term") or "").strip()
                    option_payload.extend([
                        ("True", 1 if tf_correct == "True" else 0, None),
                        ("False", 1 if tf_correct == "False" else 0, None),
                    ])
                    if correction:
                        option_payload.append((correction, 0, "correction"))
                elif q_type == "fill_in":
                    answer = (request.form.get("fill_answer") or "").strip()
                    if answer:
                        option_payload.append((answer, 1, None))
                elif q_type == "match":
                    match_a = request.form.getlist("match_a")
                    match_b = request.form.getlist("match_b")
                    for left, right in zip(match_a, match_b):
                        left = left.strip()
                        right = right.strip()
                        if left and right:
                            option_payload.append((left, 1, right))
                            break
                if bank_question_exists(cursor, q_text, q_type, subject, modules, option_payload):
                    error = "That question already exists in the bank."
                else:
                    cursor.execute("""
                        INSERT INTO question_bank_questions (question_text, question_type, marks, subject, modules, created_by, created_at)
                        VALUES (?, ?, ?, ?, ?, ?, ?)
                    """, (q_text, q_type, marks, subject, modules, username, datetime.now().isoformat()))
                    bank_question_id = cursor.lastrowid
                    for option_text, is_correct, match_pair in option_payload:
                        cursor.execute("""
                            INSERT INTO question_bank_options (bank_question_id, option_text, is_correct, match_pair)
                            VALUES (?, ?, ?, ?)
                        """, (bank_question_id, option_text, is_correct, match_pair))
                    conn.commit()
                    message = "Question added to bank."

    selected_subject = (request.args.get("subject") or "").strip()
    selected_module = (request.args.get("module") or "").strip()
    selected_question_type = (request.args.get("question_type") or "").strip()
    where_parts = []
    params = []
    if selected_subject:
        where_parts.append("LOWER(COALESCE(subject, '')) = ?")
        params.append(selected_subject.lower())
    if selected_module:
        where_parts.append("LOWER(COALESCE(modules, '')) LIKE ?")
        params.append(f"%{selected_module.lower()}%")
    if selected_question_type and selected_question_type in QUESTION_BANK_SUPPORTED_TYPES:
        where_parts.append("question_type = ?")
        params.append(selected_question_type)
    where_clause = f"WHERE {' AND '.join(where_parts)}" if where_parts else ""
    cursor.execute(f"""
        SELECT id, question_text, question_type, marks, subject, modules, created_by, created_at
        FROM question_bank_questions
        {where_clause}
        ORDER BY id DESC
    """, params)
    question_rows = cursor.fetchall()
    questions = []
    for row in question_rows:
        cursor.execute("""
            SELECT id, option_text, is_correct, match_pair
            FROM question_bank_options
            WHERE bank_question_id = ?
            ORDER BY id
        """, (row[0],))
        questions.append({"q": row, "options": cursor.fetchall()})

    cursor.execute("SELECT modules FROM question_bank_questions WHERE COALESCE(modules, '') != ''")
    module_names = []
    for (modules_text,) in cursor.fetchall():
        module_names.extend(parse_module_names(modules_text))
    module_names = sorted({item for item in module_names}, key=str.lower)
    cursor.execute("SELECT DISTINCT TRIM(subject) FROM question_bank_questions WHERE COALESCE(TRIM(subject), '') != '' ORDER BY LOWER(TRIM(subject))")
    subject_names = [row[0] for row in cursor.fetchall() if row[0]]
    conn.close()
    return render_template(
        "question_bank.html",
        questions=questions,
        module_names=module_names,
        subject_names=subject_names,
        selected_subject=selected_subject,
        selected_module=selected_module,
        selected_question_type=selected_question_type,
        question_types=QUESTION_BANK_SUPPORTED_TYPES,
        message=message,
        error=error,
    )


@app.route("/question_bank_template.csv")
def question_bank_template():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    output = io.StringIO()
    fieldnames = [
        "question_text",
        "question_type",
        "marks",
        "subject",
        "modules",
        "option_1",
        "option_2",
        "option_3",
        "option_4",
        "correct_answer",
        "correction",
        "match_a_1",
        "match_b_1",
        "match_a_2",
        "match_b_2",
        "match_a_3",
        "match_b_3",
        "match_a_4",
        "match_b_4",
    ]
    writer = csv.DictWriter(output, fieldnames=fieldnames)
    writer.writeheader()
    writer.writerows([
        {
            "question_text": "What does CPU stand for?",
            "question_type": "mcq_single",
            "marks": "1",
            "subject": "Computer Basics",
            "modules": "Hardware, Fundamentals",
            "option_1": "Central Process Unit",
            "option_2": "Central Processing Unit",
            "option_3": "Computer Personal Unit",
            "option_4": "Central Processor Utility",
            "correct_answer": "2",
        },
        {
            "question_text": "The desktop is the main screen area you see after logging in.",
            "question_type": "true_false",
            "marks": "1",
            "subject": "Computer Basics",
            "modules": "Desktop",
            "correct_answer": "True",
            "correction": "",
        },
        {
            "question_text": "The shortcut key for copy is ____.",
            "question_type": "fill_in",
            "marks": "1",
            "subject": "Keyboard",
            "modules": "Shortcuts",
            "correct_answer": "Ctrl+C",
        },
        {
            "question_text": "Match the shortcut to the action.",
            "question_type": "match",
            "marks": "1",
            "subject": "Keyboard",
            "modules": "Shortcuts",
            "match_a_1": "Ctrl+C",
            "match_b_1": "Copy",
        },
        {
            "question_text": "Match the shortcut to the action.",
            "question_type": "match",
            "marks": "1",
            "subject": "Keyboard",
            "modules": "Shortcuts",
            "match_a_1": "Ctrl+V",
            "match_b_1": "Paste",
        },
    ])

    response = app.response_class(output.getvalue(), mimetype="text/csv")
    response.headers["Content-Disposition"] = "attachment; filename=question_bank_template.csv"
    return response


@app.route("/generate_theory_test", methods=["GET", "POST"])
def generate_theory_test():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    conn = get_db()
    cursor = conn.cursor()
    groups = get_groups(username) if get_user_role(username) == "teacher" else get_groups()
    teachers = get_teachers()
    teacher_checkboxes = ''.join(
        f'<label style="font-weight:normal;display:inline-flex;align-items:center;gap:5px;">'
        f'<input type="checkbox" name="teachers" value="{escape(t[0])}"> {escape(t[1] or t[0])}</label>'
        for t in teachers
    )

    cursor.execute("SELECT modules FROM question_bank_questions WHERE COALESCE(modules, '') != ''")
    module_names = []
    for (modules_text,) in cursor.fetchall():
        module_names.extend(parse_module_names(modules_text))
    module_names = sorted({item for item in module_names}, key=str.lower)

    cursor.execute("SELECT DISTINCT TRIM(subject) FROM question_bank_questions WHERE COALESCE(TRIM(subject), '') != '' ORDER BY LOWER(TRIM(subject))")
    subject_names = [row[0] for row in cursor.fetchall() if row[0]]
    cursor.execute("SELECT COALESCE(TRIM(subject), ''), COALESCE(modules, '') FROM question_bank_questions")
    subject_module_map = {}
    for subject_name, modules_text in cursor.fetchall():
        normalized_subject = (subject_name or "").strip()
        if not normalized_subject:
            continue
        subject_module_map.setdefault(normalized_subject, [])
        subject_module_map[normalized_subject].extend(parse_module_names(modules_text))
    subject_module_map = {
        subject_name: sorted({module for module in modules if module}, key=str.lower)
        for subject_name, modules in subject_module_map.items()
    }
    counts = get_question_bank_counts(cursor, module_names, [])

    error = ""
    message = ""
    if request.method == "POST":
        title = (request.form.get("title") or "").strip()
        subject = (request.form.get("subject") or "").strip()
        assign_date = (request.form.get("assign_date") or "").strip()
        time_limit = safe_int(request.form.get("time_limit"), 0)
        allow_multiple = 1 if request.form.get("allow_multiple") else 0
        max_attempts = safe_int(request.form.get("max_attempts"), 1)
        show_answers = 1 if request.form.get("show_answers") else 0
        selected_groups = request.form.getlist("groups")
        selected_teachers = request.form.getlist("teachers")
        selected_modules = parse_module_names(",".join(request.form.getlist("modules")))
        selected_subjects = [item.strip() for item in request.form.getlist("bank_subjects") if item.strip()]
        request_counts = {q_type: max(0, safe_int(request.form.get(f"count_{q_type}"), 0)) for q_type in QUESTION_BANK_SUPPORTED_TYPES}

        if not title or not assign_date:
            error = "Title and assign date are required."
        elif not selected_modules:
            error = "Select at least one module."
        elif not any(request_counts.values()):
            error = "Choose at least one question count."
        else:
            selected_questions = []
            selected_match_pairs = []
            for q_type, needed in request_counts.items():
                if needed <= 0:
                    continue
                where_parts = ["question_type = ?"]
                params = [q_type]
                module_filters = " OR ".join(["LOWER(COALESCE(modules, '')) LIKE ?" for _ in selected_modules])
                where_parts.append(f"({module_filters})")
                params.extend(f"%{module.lower()}%" for module in selected_modules)
                if selected_subjects:
                    subject_filters = " OR ".join(["LOWER(COALESCE(subject, '')) = ?" for _ in selected_subjects])
                    where_parts.append(f"({subject_filters})")
                    params.extend(subject.lower() for subject in selected_subjects)
                cursor.execute(f"""
                    SELECT id
                    FROM question_bank_questions
                    WHERE {' AND '.join(where_parts)}
                    ORDER BY RANDOM()
                    LIMIT ?
                """, (*params, needed))
                picked = [row[0] for row in cursor.fetchall()]
                if len(picked) < needed:
                    scope_text = "selected modules"
                    if selected_subjects:
                        scope_text += " and subjects"
                    error = f"Not enough {q_type.replace('_', ' ')} questions in the {scope_text}."
                    break
                if q_type == "match":
                    selected_match_pairs = picked
                else:
                    selected_questions.extend(picked)

            if not error:
                cursor.execute("""
                    INSERT INTO theory_tests
                        (title, subject, assign_date, time_limit, allow_multiple, max_attempts, show_answers, created_by, created_at)
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
                """, (title, subject, assign_date, time_limit, allow_multiple, max_attempts, show_answers, username, datetime.now().isoformat()))
                test_id = cursor.lastrowid
                for group_name in selected_groups:
                    if group_name.strip():
                        cursor.execute("INSERT INTO theory_test_groups (test_id, group_name) VALUES (?, ?)", (test_id, group_name))
                for teacher_username in selected_teachers:
                    if teacher_username.strip():
                        cursor.execute("INSERT INTO theory_test_teachers (test_id, teacher_username) VALUES (?, ?)", (test_id, teacher_username))
                order_index = 0
                for bank_question_id in selected_questions:
                    clone_bank_question_to_test(cursor, bank_question_id, test_id, order_index)
                    order_index += 1
                if selected_match_pairs:
                    create_generated_match_question(cursor, selected_match_pairs, test_id, order_index)
                conn.commit()
                conn.close()
                log_activity(username, f"generated theory test '{title}' from question bank")
                return redirect(url_for("manage_test_questions", test_id=test_id))

    conn.close()
    return render_template(
        "generate_theory_test.html",
        groups=groups,
        teachers=teachers,
        module_names=module_names,
        subject_names=subject_names,
        subject_module_map=subject_module_map,
        counts=counts,
        error=error,
        message=message,
    )


@app.route("/question_bank_counts")
def question_bank_counts():
    username = session.get("username")
    if not username:
        return jsonify({"error": "unauthorized"}), 401
    if get_user_role(username) not in ["teacher", "admin"]:
        return jsonify({"error": "forbidden"}), 403

    modules = request.args.getlist("modules")
    subjects = request.args.getlist("subjects")
    conn = get_db()
    cursor = conn.cursor()
    counts = get_question_bank_counts(cursor, modules, subjects)
    conn.close()
    return jsonify({"counts": counts})


@app.route("/manage_tests/create", methods=["POST"])
def create_test():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    title = request.form.get("title", "").strip()
    subject = request.form.get("subject", "").strip()
    groups = request.form.getlist("groups")
    teachers = request.form.getlist("teachers")
    time_limit = safe_int(request.form.get("time_limit"), 0)
    assign_date = request.form.get("assign_date")  # YYYY-MM-DD
    allow_multiple = 1 if request.form.get("allow_multiple") else 0
    max_attempts = safe_int(request.form.get("max_attempts"), 1)
    show_answers = 1 if request.form.get("show_answers") else 0

    if not title:
        return "Test title is required", 400

    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("""
        INSERT INTO theory_tests
            (title, subject, assign_date, time_limit, allow_multiple, max_attempts, show_answers, created_by, created_at)
        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
    """, (title, subject, assign_date, time_limit, allow_multiple, max_attempts, show_answers, username, datetime.now().isoformat()))
    test_id = cursor.lastrowid
    for g in groups:
        if g.strip():
            cursor.execute("INSERT INTO theory_test_groups (test_id, group_name) VALUES (?, ?)", (test_id, g))
    for teacher in teachers:
        if teacher.strip():
            cursor.execute("INSERT INTO theory_test_teachers (test_id, teacher_username) VALUES (?, ?)", (test_id, teacher))
    conn.commit()
    conn.close()
    log_activity(username, f"created theory test '{title}'")
    return redirect(url_for("manage_test_questions", test_id=test_id))


@app.route("/manage_tests/<int:test_id>/edit", methods=["GET", "POST"])
def edit_test(test_id):
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    conn = get_db()
    cursor = conn.cursor()

    cursor.execute("SELECT id, title, subject, assign_date, time_limit, allow_multiple, max_attempts, show_answers FROM theory_tests WHERE id = ?", (test_id,))
    test = cursor.fetchone()
    if not test:
        conn.close()
        return "Test not found", 404

    if request.method == "POST":
        allow_multiple = 1 if request.form.get("allow_multiple") else 0
        max_attempts = safe_int(request.form.get("max_attempts"), 1)
        show_answers = 1 if request.form.get("show_answers") else 0
        groups = request.form.getlist("groups")
        teachers = request.form.getlist("teachers")

        assign_date = request.form.get("assign_date")
      #  time_limit = request.form.get("time_limit", 0)

        cursor.execute("""
            UPDATE theory_tests
            SET assign_date = ?, allow_multiple = ?, max_attempts = ?, show_answers = ?
            WHERE id = ?
        """, (assign_date, allow_multiple, max_attempts, show_answers, test_id))

        cursor.execute("DELETE FROM theory_test_groups WHERE test_id = ?", (test_id,))
        for g in groups:
            if g.strip():
                cursor.execute("INSERT INTO theory_test_groups (test_id, group_name) VALUES (?, ?)", (test_id, g))

        cursor.execute("DELETE FROM theory_test_teachers WHERE test_id = ?", (test_id,))
        for teacher in teachers:
            if teacher.strip():
                cursor.execute("INSERT INTO theory_test_teachers (test_id, teacher_username) VALUES (?, ?)", (test_id, teacher))

        conn.commit()
        conn.close()
        log_activity(username, f"edited theory test settings for test {test_id}")
        return redirect(url_for("manage_tests"))

    # GET — load current groups
    cursor.execute("SELECT group_name FROM theory_test_groups WHERE test_id = ?", (test_id,))
    current_groups = {row[0] for row in cursor.fetchall()}
    cursor.execute("SELECT teacher_username FROM theory_test_teachers WHERE test_id = ?", (test_id,))
    current_teachers = {row[0] for row in cursor.fetchall()}
    all_groups = get_groups()
    all_teachers = get_teachers()
    conn.close()

    return render_template("edit_test.html", test=test, current_groups=current_groups, all_groups=all_groups, all_teachers=all_teachers, current_teachers=current_teachers)


@app.route("/manage_tests/<int:test_id>/toggle", methods=["POST"])
def toggle_test(test_id):
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("SELECT is_active FROM theory_tests WHERE id = ?", (test_id,))
    row = cursor.fetchone()
    if row:
        new_state = 0 if row[0] else 1
        cursor.execute("UPDATE theory_tests SET is_active = ? WHERE id = ?", (new_state, test_id))
        conn.commit()
    conn.close()
    return redirect(url_for("manage_tests"))


@app.route("/manage_tests/<int:test_id>/delete", methods=["POST"])
def delete_test(test_id):
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    conn = get_db()
    cursor = conn.cursor()
    
    # Get test title for logging
    cursor.execute("SELECT title FROM theory_tests WHERE id = ?", (test_id,))
    test_row = cursor.fetchone()
    test_title = test_row[0] if test_row else f"Test {test_id}"
    
    # Delete in order: answers → submissions → options → questions → test_groups → tasks → test
    cursor.execute("DELETE FROM theory_answers WHERE submission_id IN (SELECT id FROM theory_submissions WHERE test_id = ?)", (test_id,))
    cursor.execute("DELETE FROM theory_submissions WHERE test_id = ?", (test_id,))
    cursor.execute("DELETE FROM theory_options WHERE question_id IN (SELECT id FROM theory_questions WHERE test_id = ?)", (test_id,))
    cursor.execute("DELETE FROM theory_questions WHERE test_id = ?", (test_id,))
    cursor.execute("DELETE FROM theory_test_groups WHERE test_id = ?", (test_id,))
    # Also delete any tasks that reference this theory test
    cursor.execute("DELETE FROM task_groups WHERE task_id IN (SELECT id FROM tasks WHERE theory_test_id = ?)", (test_id,))
    cursor.execute("DELETE FROM tasks WHERE theory_test_id = ?", (test_id,))
    cursor.execute("DELETE FROM theory_tests WHERE id = ?", (test_id,))
    conn.commit()
    conn.close()
    log_activity(username, f"deleted theory test '{test_title}' and all related submissions")
    return redirect(url_for("manage_tests"))


@app.route("/manage_tests/<int:test_id>/questions", methods=["GET", "POST"])
def manage_test_questions(test_id):
    return _manage_theory_questions(test_id, "test")


@app.route("/manage_lessons/<int:test_id>/questions", methods=["GET", "POST"])
def manage_lesson_questions(test_id):
    return _manage_theory_questions(test_id, "lesson")


def _manage_theory_questions(test_id, builder_mode):
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    conn = get_db()
    cursor = conn.cursor()

    cursor.execute("SELECT id, title, subject, background_image, COALESCE(background_fit, 'cover') FROM theory_tests WHERE id = ?", (test_id,))
    test = cursor.fetchone()
    if not test:
        conn.close()
        return "Test not found", 404

    if request.method == "POST":
        action = request.form.get("action")
        if action == "add_question":
            q_text = request.form.get("question_text", "").strip()
            if builder_mode == "lesson":
                q_text = externalize_data_uri_images(q_text)
            q_type = request.form.get("question_type", "")
            marks = int(request.form.get("marks", 1))
            if q_type in LESSON_SLIDE_TYPES:
                marks = 0

            cursor.execute("SELECT COUNT(*) FROM theory_questions WHERE test_id = ?", (test_id,))
            order_index = cursor.fetchone()[0]

            cursor.execute("""
                INSERT INTO theory_questions (test_id, question_text, question_type, marks, order_index)
                VALUES (?, ?, ?, ?, ?)
            """, (test_id, q_text, q_type, marks, order_index))
            q_id = cursor.lastrowid

            # Handle options per question type
            if q_type in ["mcq_single", "mcq_multi"]:
                options = request.form.getlist("option_text")
                correct = request.form.getlist("is_correct")
                for i, opt in enumerate(options):
                    if opt.strip():
                        is_correct = 1 if str(i) in correct else 0
                        cursor.execute("""
                            INSERT INTO theory_options (question_id, option_text, is_correct)
                            VALUES (?, ?, ?)
                        """, (q_id, opt.strip(), is_correct))

            elif q_type == "true_false":
                correct_answer = request.form.get("tf_correct", "True")
                correction_term = request.form.get("correction_term", "").strip()
                cursor.execute("INSERT INTO theory_options (question_id, option_text, is_correct) VALUES (?, 'True', ?)",
                               (q_id, 1 if correct_answer == "True" else 0))
                cursor.execute("INSERT INTO theory_options (question_id, option_text, is_correct) VALUES (?, 'False', ?)",
                               (q_id, 1 if correct_answer == "False" else 0))
                if correction_term:
                    cursor.execute("INSERT INTO theory_options (question_id, option_text, is_correct, match_pair) VALUES (?, ?, 0, 'correction')",
                                   (q_id, correction_term))

            elif q_type == "fill_in":
                answer = request.form.get("fill_answer", "").strip()
                cursor.execute("INSERT INTO theory_options (question_id, option_text, is_correct) VALUES (?, ?, 1)",
                               (q_id, answer))

            elif q_type == "match":
                col_a = request.form.getlist("match_a")
                col_b = request.form.getlist("match_b")
                for a, b in zip(col_a, col_b):
                    if a.strip() and b.strip():
                        cursor.execute("""
                            INSERT INTO theory_options (question_id, option_text, is_correct, match_pair)
                            VALUES (?, ?, 1, ?)
                        """, (q_id, a.strip(), b.strip()))
            conn.commit()
            log_activity(username, f"added question to test {test_id}")

        elif action == "duplicate_question" and builder_mode == "lesson":
            q_id = request.form.get("question_id")
            cursor.execute("""
                SELECT question_text, question_type, marks, order_index
                FROM theory_questions
                WHERE id = ? AND test_id = ?
            """, (q_id, test_id))
            source = cursor.fetchone()
            if source:
                q_text, q_type, marks, order_index = source
                cursor.execute(
                    "UPDATE theory_questions SET order_index = order_index + 1 WHERE test_id = ? AND order_index > ?",
                    (test_id, order_index)
                )
                cursor.execute("""
                    INSERT INTO theory_questions (test_id, question_text, question_type, marks, order_index)
                    VALUES (?, ?, ?, ?, ?)
                """, (test_id, q_text, q_type, marks, order_index + 1))
                new_q_id = cursor.lastrowid
                cursor.execute("""
                    SELECT option_text, is_correct, match_pair
                    FROM theory_options
                    WHERE question_id = ?
                    ORDER BY id
                """, (q_id,))
                for option_text, is_correct, match_pair in cursor.fetchall():
                    cursor.execute("""
                        INSERT INTO theory_options (question_id, option_text, is_correct, match_pair)
                        VALUES (?, ?, ?, ?)
                    """, (new_q_id, option_text, is_correct, match_pair))
                conn.commit()
                log_activity(username, f"duplicated slide {q_id} in test {test_id}")

        elif action == "reorder_questions" and builder_mode == "lesson":
            ordered_ids = request.form.get("ordered_ids", "")
            ids = [int(item) for item in ordered_ids.split(",") if item.strip().isdigit()]
            for index, q_id in enumerate(ids):
                cursor.execute(
                    "UPDATE theory_questions SET order_index = ? WHERE id = ? AND test_id = ?",
                    (index, q_id, test_id)
                )
            conn.commit()
            log_activity(username, f"reordered slides in test {test_id}")

        elif action == "delete_question":
            q_id = request.form.get("question_id")
            cursor.execute("DELETE FROM theory_options WHERE question_id = ?", (q_id,))
            cursor.execute("DELETE FROM theory_questions WHERE id = ?", (q_id,))
            conn.commit()

        elif action == "edit_question":
            q_id = request.form.get("question_id")
            q_text = request.form.get("question_text", "").strip()
            if builder_mode == "lesson":
                q_text = externalize_data_uri_images(q_text)
            q_type = request.form.get("question_type", "")
            marks = int(request.form.get("marks", 1))
            if q_type in LESSON_SLIDE_TYPES:
                marks = 0

            cursor.execute("UPDATE theory_questions SET question_text = ?, marks = ? WHERE id = ?",
                           (q_text, marks, q_id))

            # Replace all options
            cursor.execute("DELETE FROM theory_options WHERE question_id = ?", (q_id,))

            if q_type in ["mcq_single", "mcq_multi"]:
                options = request.form.getlist("option_text")
                correct = request.form.getlist("is_correct")
                for i, opt in enumerate(options):
                    if opt.strip():
                        is_correct = 1 if str(i) in correct else 0
                        cursor.execute("INSERT INTO theory_options (question_id, option_text, is_correct) VALUES (?, ?, ?)",
                                       (q_id, opt.strip(), is_correct))

            elif q_type == "true_false":
                correct_answer = request.form.get("tf_correct", "True")
                correction_term = request.form.get("correction_term", "").strip()
                cursor.execute("INSERT INTO theory_options (question_id, option_text, is_correct) VALUES (?, 'True', ?)",
                               (q_id, 1 if correct_answer == "True" else 0))
                cursor.execute("INSERT INTO theory_options (question_id, option_text, is_correct) VALUES (?, 'False', ?)",
                               (q_id, 1 if correct_answer == "False" else 0))
                if correction_term:
                    cursor.execute("INSERT INTO theory_options (question_id, option_text, is_correct, match_pair) VALUES (?, ?, 0, 'correction')",
                                   (q_id, correction_term))

            elif q_type == "fill_in":
                answer = request.form.get("fill_answer", "").strip()
                cursor.execute("INSERT INTO theory_options (question_id, option_text, is_correct) VALUES (?, ?, 1)",
                               (q_id, answer))

            elif q_type == "match":
                col_a = request.form.getlist("match_a")
                col_b = request.form.getlist("match_b")
                for a, b in zip(col_a, col_b):
                    if a.strip() and b.strip():
                        cursor.execute("INSERT INTO theory_options (question_id, option_text, is_correct, match_pair) VALUES (?, ?, 1, ?)",
                                       (q_id, a.strip(), b.strip()))
            conn.commit()
            log_activity(username, f"edited question {q_id} in test {test_id}")

        elif action == "autosave_question" and builder_mode == "lesson":
            q_id = request.form.get("question_id")
            q_text = externalize_data_uri_images(request.form.get("question_text", "").strip())
            marks = safe_int(request.form.get("marks"), 0)
            q_type = request.form.get("question_type", "")
            if q_type in LESSON_SLIDE_TYPES:
                marks = 0

            cursor.execute("SELECT 1 FROM theory_questions WHERE id = ? AND test_id = ?", (q_id, test_id))
            if not cursor.fetchone():
                conn.close()
                return {"ok": False, "error": "not_found"}, 404

            cursor.execute("UPDATE theory_questions SET question_text = ?, marks = ? WHERE id = ?", (q_text, marks, q_id))
            cursor.execute("DELETE FROM theory_options WHERE question_id = ?", (q_id,))

            if q_type in ["mcq_single", "mcq_multi"]:
                options = request.form.getlist("option_text")
                correct = request.form.getlist("is_correct")
                for i, opt in enumerate(options):
                    if opt.strip():
                        cursor.execute(
                            "INSERT INTO theory_options (question_id, option_text, is_correct) VALUES (?, ?, ?)",
                            (q_id, opt.strip(), 1 if str(i) in correct else 0)
                        )
            elif q_type == "true_false":
                correct_answer = request.form.get("tf_correct", "True")
                correction_term = request.form.get("correction_term", "").strip()
                cursor.execute("INSERT INTO theory_options (question_id, option_text, is_correct) VALUES (?, 'True', ?)",
                               (q_id, 1 if correct_answer == "True" else 0))
                cursor.execute("INSERT INTO theory_options (question_id, option_text, is_correct) VALUES (?, 'False', ?)",
                               (q_id, 1 if correct_answer == "False" else 0))
                if correction_term:
                    cursor.execute("INSERT INTO theory_options (question_id, option_text, is_correct, match_pair) VALUES (?, ?, 0, 'correction')",
                                   (q_id, correction_term))
            elif q_type == "fill_in":
                answer = request.form.get("fill_answer", "").strip()
                cursor.execute("INSERT INTO theory_options (question_id, option_text, is_correct) VALUES (?, ?, 1)", (q_id, answer))
            elif q_type == "match":
                col_a = request.form.getlist("match_a")
                col_b = request.form.getlist("match_b")
                for a, b in zip(col_a, col_b):
                    if a.strip() and b.strip():
                        cursor.execute(
                            "INSERT INTO theory_options (question_id, option_text, is_correct, match_pair) VALUES (?, ?, 1, ?)",
                            (q_id, a.strip(), b.strip())
                        )

            if "background_image" in request.form:
                background_image = externalize_data_uri_images(request.form.get("background_image", "").strip() or None)
                background_fit = request.form.get("background_fit", "cover").strip() or "cover"
                if background_fit not in ("cover", "contain", "stretch"):
                    background_fit = "cover"
                cursor.execute(
                    "UPDATE theory_tests SET background_image = ?, background_fit = ? WHERE id = ?",
                    (background_image, background_fit, test_id)
                )

            conn.commit()
            return {"ok": True, "saved_at": datetime.now().strftime("%H:%M:%S")}

        elif action == "save_lesson_background" and builder_mode == "lesson":
            background_image = externalize_data_uri_images(request.form.get("background_image", "").strip() or None)
            background_fit = request.form.get("background_fit", "cover").strip() or "cover"
            if background_fit not in ("cover", "contain", "stretch"):
                background_fit = "cover"
            cursor.execute(
                "UPDATE theory_tests SET background_image = ?, background_fit = ? WHERE id = ?",
                (background_image, background_fit, test_id)
            )
            conn.commit()
            log_activity(username, f"updated lesson background for test {test_id}")

        elif action == "import_pptx" and builder_mode == "lesson":
            import_success = None
            import_error = None
            pptx_file = request.files.get("pptx_file")
            append = request.form.get("append") is not None

            try:
                if not pptx_file or not pptx_file.filename:
                    raise ValueError("Please choose a PowerPoint .pptx file.")
                if not pptx_file.filename.lower().endswith(".pptx"):
                    raise ValueError("Only .pptx files can be imported.")

                if not append:
                    cursor.execute("""
                        DELETE FROM theory_options
                        WHERE question_id IN (SELECT id FROM theory_questions WHERE test_id = ?)
                    """, (test_id,))
                    cursor.execute("DELETE FROM theory_questions WHERE test_id = ?", (test_id,))
                    order_index = 0
                else:
                    cursor.execute("""
                        SELECT COALESCE(MAX(order_index), -1) + 1
                        FROM theory_questions
                        WHERE test_id = ?
                    """, (test_id,))
                    order_index = cursor.fetchone()[0]

                slide_html = pptx_to_content_slide_html(pptx_file)
                if not slide_html:
                    raise ValueError("No text or pictures were found in that PowerPoint file.")

                for html in slide_html:
                    cursor.execute("""
                        INSERT INTO theory_questions (test_id, question_text, question_type, marks, order_index)
                        VALUES (?, ?, 'content_slide', 0, ?)
                    """, (test_id, externalize_data_uri_images(html), order_index))
                    order_index += 1

                conn.commit()
                log_activity(username, f"imported PPTX slides into test {test_id}")
                import_success = f"Imported {len(slide_html)} PowerPoint slide(s)."
            except Exception as e:
                conn.rollback()
                import_error = str(e)

            cursor.execute("""
                SELECT id, question_text, question_type, marks, order_index
                FROM theory_questions WHERE test_id = ? ORDER BY order_index
            """, (test_id,))
            questions = cursor.fetchall()
            questions_with_options = []
            for q in questions:
                cursor.execute("SELECT id, option_text, is_correct, match_pair FROM theory_options WHERE question_id = ?", (q[0],))
                options = cursor.fetchall()
                questions_with_options.append({"q": q, "options": options})

            conn.close()
            return render_template(
                "manage_lesson_questions.html",
                test=test,
                questions=questions_with_options,
                import_success=import_success,
                import_error=import_error,
                builder_mode=builder_mode
            )

        elif action == "import_questions_json":
            import_success = None
            import_error = None

            from theory_json_importer import insert_theory_test_from_json

            payload = request.form.get("questions_json", "").strip()
            # Append-only by default: checkbox checked => append
            append = request.form.get("append") is not None

            try:
                if not payload:
                    raise ValueError("questions_json is empty")

                if not append:
                    # Replace: clear existing questions/options for this test
                    cursor.execute("""
                        DELETE FROM theory_options
                        WHERE question_id IN (SELECT id FROM theory_questions WHERE test_id = ?)
                    """, (test_id,))
                    cursor.execute("DELETE FROM theory_questions WHERE test_id = ?", (test_id,))
                    start_order_index = 0
                else:
                    cursor.execute("""
                        SELECT COALESCE(MAX(order_index), -1) + 1
                        FROM theory_questions
                        WHERE test_id = ?
                    """, (test_id,))
                    start_order_index = cursor.fetchone()[0]

                insert_theory_test_from_json(
                    cursor,
                    test_id=test_id,
                    username=username,
                    payload=payload,
                    start_order_index=start_order_index
                )

                conn.commit()
                log_activity(username, f"imported questions into test {test_id} (append={append})")
                import_success = "Import completed successfully."

            except Exception as e:
                conn.rollback()
                import_error = str(e)

            # Reload questions and re-render (same page) with import status
            cursor.execute("""
                SELECT id, question_text, question_type, marks, order_index
                FROM theory_questions WHERE test_id = ? ORDER BY order_index
            """, (test_id,))
            questions = cursor.fetchall()

            questions_with_options = []
            for q in questions:
                cursor.execute("SELECT id, option_text, is_correct, match_pair FROM theory_options WHERE question_id = ?", (q[0],))
                options = cursor.fetchall()
                questions_with_options.append({"q": q, "options": options})

            conn.close()
            return render_template(
                "manage_lesson_questions.html" if builder_mode == "lesson" else "manage_test_questions.html",
                test=test,
                questions=questions_with_options,
                import_success=import_success,
                import_error=import_error,
                builder_mode=builder_mode
            )

        conn.close()
        if builder_mode == "lesson":
            return redirect(url_for("manage_lesson_questions", test_id=test_id))
        return redirect(url_for("manage_test_questions", test_id=test_id))

    # GET — load questions with their options
    cursor.execute("""
        SELECT id, question_text, question_type, marks, order_index
        FROM theory_questions WHERE test_id = ? ORDER BY order_index
    """, (test_id,))
    questions = cursor.fetchall()

    questions_with_options = []
    for q in questions:
        cursor.execute("SELECT id, option_text, is_correct, match_pair FROM theory_options WHERE question_id = ?", (q[0],))
        options = cursor.fetchall()
        questions_with_options.append({"q": q, "options": options})

    conn.close()
    return render_template(
        "manage_lesson_questions.html" if builder_mode == "lesson" else "manage_test_questions.html",
        test=test,
        questions=questions_with_options,
        builder_mode=builder_mode
    )


@app.route("/tests")
def learner_tests():
    """Show available tests to a learner."""
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("SELECT group_name FROM users WHERE username = ?", (username,))
    row = cursor.fetchone()
    user_group = row[0] if row else None

    cursor.execute("""
        SELECT t.id, t.title, t.subject, t.time_limit,
               best.best_percentage,
               sub.id as latest_submission_id,
               t.allow_multiple, t.max_attempts,
               COALESCE(cnt.attempt_count, 0) as attempt_count,
               COALESCE(qstats.content_count, 0) as content_count,
               COALESCE(qstats.question_count, 0) as question_count,
               COALESCE(prog.current_slide, 0) as current_slide,
               COALESCE(prog.completed, 0) as progress_completed,
               COALESCE(prog.max_slide, 0) as max_slide,
               COALESCE(prog.time_spent_seconds, 0) as progress_seconds
        FROM theory_tests t
        LEFT JOIN (
            SELECT test_id, MAX(percentage) as best_percentage
            FROM theory_submissions WHERE username = ?
            GROUP BY test_id
        ) best ON t.id = best.test_id
        LEFT JOIN (
            SELECT test_id, id, percentage,
                   ROW_NUMBER() OVER (PARTITION BY test_id ORDER BY submitted_at DESC) as rn
            FROM theory_submissions WHERE username = ?
        ) sub ON t.id = sub.test_id AND sub.rn = 1
        LEFT JOIN (
            SELECT test_id, COUNT(*) as attempt_count
            FROM theory_submissions WHERE username = ?
            GROUP BY test_id
        ) cnt ON t.id = cnt.test_id
        LEFT JOIN (
            SELECT test_id,
                   SUM(CASE WHEN question_type IN ('content_slide', 'title_slide', 'heading_slide') THEN 1 ELSE 0 END) as content_count,
                   COUNT(*) as question_count
            FROM theory_questions
            GROUP BY test_id
        ) qstats ON t.id = qstats.test_id
        LEFT JOIN theory_progress prog ON prog.test_id = t.id AND prog.username = ?
        WHERE t.is_active = 1
          AND COALESCE(qstats.content_count, 0) = 0
          AND (
              NOT EXISTS (SELECT 1 FROM theory_test_groups WHERE test_id = t.id)
              OR EXISTS (SELECT 1 FROM theory_test_groups WHERE test_id = t.id AND group_name = ?)
          )
        GROUP BY t.id
        ORDER BY t.created_at DESC
    """, (username, username, username, username, user_group))
    tests = cursor.fetchall()
    conn.close()
    return render_template("learner_tests.html", tests=tests)


@app.route("/lesson_tests")
def lesson_tests():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("SELECT group_name FROM users WHERE username = ?", (username,))
    row = cursor.fetchone()
    user_group = row[0] if row else None

    cursor.execute("""
        SELECT t.id, t.title, t.subject, t.time_limit,
               best.best_percentage,
               sub.id as latest_submission_id,
               t.allow_multiple, t.max_attempts,
               COALESCE(cnt.attempt_count, 0) as attempt_count,
               COALESCE(qstats.content_count, 0) as content_count,
               COALESCE(qstats.total_questions, 0) as total_questions,
               COALESCE(prog.current_slide, 0) as current_slide,
               COALESCE(prog.max_slide, 0) as max_slide,
               COALESCE(prog.time_spent_seconds, 0) as progress_seconds
        FROM theory_tests t
        LEFT JOIN (
            SELECT test_id, MAX(percentage) as best_percentage
            FROM theory_submissions WHERE username = ?
            GROUP BY test_id
        ) best ON t.id = best.test_id
        LEFT JOIN (
            SELECT test_id, id, percentage,
                   ROW_NUMBER() OVER (PARTITION BY test_id ORDER BY submitted_at DESC) as rn
            FROM theory_submissions WHERE username = ?
        ) sub ON t.id = sub.test_id AND sub.rn = 1
        LEFT JOIN (
            SELECT test_id, COUNT(*) as attempt_count
            FROM theory_submissions WHERE username = ?
            GROUP BY test_id
        ) cnt ON t.id = cnt.test_id
        LEFT JOIN (
            SELECT test_id,
                   SUM(CASE WHEN question_type IN ('content_slide', 'title_slide', 'heading_slide') THEN 1 ELSE 0 END) as content_count,
                   COUNT(*) as total_questions
            FROM theory_questions
            GROUP BY test_id
        ) qstats ON t.id = qstats.test_id
        LEFT JOIN theory_progress prog ON prog.test_id = t.id AND prog.username = ?
        WHERE t.is_active = 1
          AND COALESCE(qstats.content_count, 0) > 0
          AND (
              NOT EXISTS (SELECT 1 FROM theory_test_groups WHERE test_id = t.id)
              OR EXISTS (SELECT 1 FROM theory_test_groups WHERE test_id = t.id AND group_name = ?)
          )
        GROUP BY t.id
        ORDER BY t.created_at DESC
    """, (username, username, username, username, user_group))
    tests = cursor.fetchall()
    conn.close()
    return render_template("learner_lesson_tests.html", tests=tests)


@app.route("/my_tasks")
def learner_tasks():
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    conn = get_db()
    cursor = conn.cursor()
    role = get_user_role(username)
    cursor.execute("SELECT group_name FROM users WHERE username = ?", (username,))
    row = cursor.fetchone()
    user_group = row[0] if row else None
    today = datetime.now().date().isoformat()

    if role in ["teacher", "admin"]:
        cursor.execute("""
            SELECT t.id, t.name, t.assign_date, t.task_type, t.allow_multiple,
                   t.max_attempts, t.is_active, t.theory_test_id, t.subject_id, s.name
            FROM tasks t
            JOIN subjects s ON t.subject_id = s.id
            ORDER BY t.assign_date, s.name, t.name
        """)
    else:
        cursor.execute("""
            SELECT t.id, t.name, t.assign_date, t.task_type, t.allow_multiple,
                   t.max_attempts, t.is_active, t.theory_test_id, t.subject_id, s.name
            FROM tasks t
            JOIN subjects s ON t.subject_id = s.id
            JOIN task_groups tg ON t.id = tg.task_id
            WHERE tg.group_name = ? AND t.assign_date <= ? AND t.is_active = 1
            ORDER BY t.assign_date, s.name, t.name
        """, (user_group, today))

    task_rows = cursor.fetchall()
    tasks = []
    for task_id, task_name, assign_date, task_type, allow_multiple, max_attempts, is_active, theory_test_id, subject_id, subject_name in task_rows:
        cursor.execute(
            "SELECT COUNT(*), COALESCE(MAX(score), 0) FROM results WHERE username = ? AND subject = ? AND task = ?",
            (username, subject_name, task_name)
        )
        submission_count, best_score = cursor.fetchone()
        submission_count = submission_count or 0

        tasks.append({
            "id": task_id,
            "name": task_name,
            "assign_date": assign_date,
            "task_type": task_type,
            "allow_multiple": allow_multiple,
            "max_attempts": max_attempts,
            "is_active": is_active,
            "theory_test_id": theory_test_id,
            "subject_id": subject_id,
            "subject": subject_name,
            "submission_count": submission_count,
            "best_score": best_score,
        })

    conn.close()
    return render_template("learner_tasks.html", tasks=tasks, username=username)


def pptx_to_content_slide_html(uploaded_file):
    from pptx import Presentation

    prs = Presentation(uploaded_file)
    slide_w = float(prs.slide_width)
    slide_h = float(prs.slide_height)
    html_slides = []

    for slide in prs.slides:
        boxes = []
        for shape in slide.shapes:
            left = max(0, (float(shape.left) / slide_w) * 100)
            top = max(0, (float(shape.top) / slide_h) * 100)
            width = max(5, (float(shape.width) / slide_w) * 100)
            height = max(5, (float(shape.height) / slide_h) * 100)
            style = f"left:{left:.2f}%;top:{top:.2f}%;width:{width:.2f}%;height:{height:.2f}%;"

            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if not text:
                    continue
                paragraphs = "".join(
                    f"<p>{escape(line)}</p>"
                    for line in text.splitlines()
                    if line.strip()
                )
                boxes.append(f'<div class="slide-box text-box" style="{style}">{paragraphs}</div>')
                continue

            if hasattr(shape, "image"):
                image = shape.image
                ext = (image.ext or "png").lower()
                mime = mimetypes.types_map.get(f".{ext}", "image/png")
                encoded = base64.b64encode(image.blob).decode("ascii")
                src = f"data:{mime};base64,{encoded}"
                boxes.append(f'<div class="slide-box image-box" style="{style}"><img src="{src}" alt=""></div>')

        if boxes:
            html_slides.append("".join(boxes))

    return html_slides


@app.route("/preview_test/<int:test_id>")
def preview_test(test_id):
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    conn = get_db()
    cursor = conn.cursor()

    cursor.execute("SELECT id, title, subject, time_limit, allow_multiple, max_attempts, show_answers, background_image, COALESCE(background_fit, 'cover') FROM theory_tests WHERE id = ?", (test_id,))
    test = cursor.fetchone()
    if not test:
        conn.close()
        return "Test not found", 404

    cursor.execute("""
        SELECT id, question_text, question_type, marks
        FROM theory_questions WHERE test_id = ? ORDER BY order_index
    """, (test_id,))
    questions = cursor.fetchall()

    questions_with_options = []
    for q in questions:
        cursor.execute("SELECT id, option_text, is_correct, match_pair FROM theory_options WHERE question_id = ?", (q[0],))
        options = list(cursor.fetchall())
        if q[2] == "match":
            options = sorted(options, key=lambda o: o[0])
        questions_with_options.append({"q": q, "options": options})
    content_slide_count = sum(1 for q in questions if q[2] in LESSON_SLIDE_TYPES)
    conn.close()
    return render_template(
        "take_test.html",
        test=test,
        questions=questions_with_options,
        attempt_number=0,
        is_preview=True,
        initial_slide=0,
        has_content_slides=content_slide_count > 0,
        content_slide_count=content_slide_count,
        question_count=len(questions) - content_slide_count,
        is_lesson_only=content_slide_count > 0 and len(questions) == content_slide_count,
        max_slide=0,
        lesson_completed=False
    )


@app.route("/tests/<int:test_id>/progress", methods=["POST"])
def save_theory_progress(test_id):
    username = session.get("username")
    if not username:
        return {"ok": False, "error": "not_logged_in"}, 401

    payload = request.get_json(silent=True) or {}
    current_slide = int(payload.get("current_slide", 0) or 0)
    max_slide = int(payload.get("max_slide", current_slide) or 0)
    time_delta = max(0, int(payload.get("time_delta", 0) or 0))
    completed = 1 if payload.get("completed") else 0

    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("SELECT 1 FROM theory_tests WHERE id = ?", (test_id,))
    if not cursor.fetchone():
        conn.close()
        return {"ok": False, "error": "not_found"}, 404

    cursor.execute("""
        INSERT INTO theory_progress (test_id, username, current_slide, max_slide, time_spent_seconds, completed, updated_at)
        VALUES (?, ?, ?, ?, ?, ?, ?)
        ON CONFLICT(test_id, username) DO UPDATE SET
            current_slide = excluded.current_slide,
            max_slide = MAX(theory_progress.max_slide, excluded.max_slide),
            time_spent_seconds = theory_progress.time_spent_seconds + excluded.time_spent_seconds,
            completed = CASE WHEN excluded.completed = 1 THEN 1 ELSE theory_progress.completed END,
            updated_at = excluded.updated_at
    """, (test_id, username, current_slide, max_slide, time_delta, completed, datetime.now().isoformat()))
    conn.commit()
    conn.close()
    return {"ok": True}


@app.route("/take_test/<int:test_id>", methods=["GET", "POST"])
def take_test(test_id):
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    conn = get_db()
    cursor = conn.cursor()

    cursor.execute("SELECT id, title, subject, time_limit, allow_multiple, max_attempts, show_answers, background_image, COALESCE(background_fit, 'cover') FROM theory_tests WHERE id = ? AND is_active = 1", (test_id,))
    test = cursor.fetchone()
    if not test:
        conn.close()
        return "Test not found or not available", 404

    cursor.execute("""
        SELECT id, question_text, question_type, marks
        FROM theory_questions WHERE test_id = ? ORDER BY order_index
    """, (test_id,))
    questions = cursor.fetchall()
    content_slide_count = sum(1 for q in questions if q[2] in LESSON_SLIDE_TYPES)
    question_count = len(questions) - content_slide_count
    has_content_slides = content_slide_count > 0
    is_lesson_only = content_slide_count > 0 and question_count == 0

    cursor.execute("SELECT COUNT(*) FROM theory_submissions WHERE test_id = ? AND username = ?", (test_id, username))
    attempt_count = cursor.fetchone()[0]
    allow_multiple = test[4]
    max_attempts = test[5]

    redirect_endpoint = "lesson_tests" if has_content_slides else "learner_tests"

    if not is_lesson_only and not allow_multiple and attempt_count > 0:
        conn.close()
        return redirect(url_for(redirect_endpoint))

    if not is_lesson_only and allow_multiple and attempt_count >= max_attempts:
        conn.close()
        return redirect(url_for(redirect_endpoint))

    if request.method == "GET":
        cursor.execute("SELECT current_slide, max_slide, completed FROM theory_progress WHERE test_id = ? AND username = ?", (test_id, username))
        progress_row = cursor.fetchone()
        initial_slide = progress_row[0] if progress_row and is_lesson_only and not progress_row[2] else 0
        max_slide = (len(questions) - 1) if progress_row and is_lesson_only and progress_row[2] else (progress_row[1] if progress_row and is_lesson_only else 0)
        questions_with_options = []
        session_order = {}
        for q in questions:
            cursor.execute("SELECT id, option_text, is_correct, match_pair FROM theory_options WHERE question_id = ?", (q[0],))
            options = list(cursor.fetchall())
            q_type = q[2]
            if q_type in ["mcq_single", "mcq_multi"]:
                random.shuffle(options)
            elif q_type == "match":
                # For match questions we only randomize Column B display order.
                # Keep Column A order stable so grading indices still match.
                # Structure of `options` for match questions is:
                #   (id, col_a_text, is_correct, match_pair_col_b_text)
                options = sorted(options, key=lambda o: o[0])
                b_opts = options[:]
                random.shuffle(b_opts)
                # Re-pack so Column A values remain in original order (options),
                # while Column B (match_pair) is shuffled for display.
                options = [
                    (o[0], o[1], o[2], b_opts[i][3])
                    for i, o in enumerate(options)
                ]
            # Store option IDs in shuffled order in session
            session_order[str(q[0])] = [o[0] for o in options]
            questions_with_options.append({"q": q, "options": options})
        session[f"test_order_{test_id}"] = session_order
        conn.close()
        return render_template(
            "take_test.html",
            test=test,
            questions=questions_with_options,
            attempt_number=attempt_count + 1,
            initial_slide=initial_slide,
            has_content_slides=has_content_slides,
            content_slide_count=content_slide_count,
            question_count=question_count,
            is_lesson_only=is_lesson_only,
            max_slide=max_slide,
            lesson_completed=is_lesson_only and attempt_count > 0
        )

    # POST — restore shuffled order from session
    session_order = session.get(f"test_order_{test_id}", {})
    questions_with_options = []
    for q in questions:
        q_id = q[0]
        cursor.execute("SELECT id, option_text, is_correct, match_pair FROM theory_options WHERE question_id = ?", (q_id,))
        all_options = {o[0]: o for o in cursor.fetchall()}
        stored_ids = session_order.get(str(q_id), [])
        if stored_ids:
            options = [all_options[oid] for oid in stored_ids if oid in all_options]
        else:
            options = list(all_options.values())
        questions_with_options.append({"q": q, "options": options})

    if is_lesson_only:
        time_spent = max(0, int(request.form.get("time_spent_seconds", 0) or 0))
        cursor.execute("""
            SELECT COALESCE(time_spent_seconds, 0)
            FROM theory_progress
            WHERE test_id = ? AND username = ?
        """, (test_id, username))
        saved_time_row = cursor.fetchone()
        saved_time = saved_time_row[0] if saved_time_row else 0
        total_time = saved_time + time_spent
        cursor.execute("""
            INSERT INTO theory_submissions (test_id, username, score, total, percentage, submitted_at, time_spent_seconds, submission_type)
            VALUES (?, ?, 1, 1, 100, ?, ?, 'lesson')
        """, (test_id, username, datetime.now().isoformat(), total_time))
        submission_id = cursor.lastrowid
        for item in questions_with_options:
            cursor.execute("""
                INSERT INTO theory_answers (submission_id, question_id, answer_text, is_correct, marks_awarded)
                VALUES (?, ?, 'Viewed', 1, 0)
            """, (submission_id, item["q"][0]))
        cursor.execute("""
            INSERT INTO theory_progress (test_id, username, current_slide, max_slide, time_spent_seconds, completed, updated_at)
            VALUES (?, ?, ?, ?, 0, 1, ?)
            ON CONFLICT(test_id, username) DO UPDATE SET
                current_slide = excluded.current_slide,
                max_slide = excluded.max_slide,
                completed = 1,
                updated_at = excluded.updated_at
        """, (test_id, username, len(questions_with_options) - 1, len(questions_with_options) - 1, datetime.now().isoformat()))
        conn.commit()
        conn.close()
        session.pop(f"test_order_{test_id}", None)
        log_activity(username, f"completed lesson {test_id} — 100%")
        return redirect(url_for("lesson_tests"))

    score = 0
    total = 0
    answers_to_save = []

    for item in questions_with_options:
        q = item["q"]
        q_id, q_text, q_type, marks = q
        options = item["options"]
        if q_type in ["content_slide", "title_slide", "heading_slide"]:
            answers_to_save.append((q_id, "Viewed", 1, 0))
            continue
        effective_marks = len([o for o in options if o[3] and o[3] != 'correction']) if q_type == "match" else marks
        total += effective_marks
        awarded = 0
        answer_text = ""

        if q_type == "mcq_single":
            selected = request.form.get(f"q_{q_id}")
            answer_text = selected or ""
            cursor.execute("SELECT option_text FROM theory_options WHERE question_id = ? AND is_correct = 1 LIMIT 1", (q_id,))
            correct_row = cursor.fetchone()
            correct_option = correct_row[0] if correct_row else None
            if selected and selected == correct_option:
                awarded = marks

        elif q_type == "mcq_multi":
            selected = set(request.form.getlist(f"q_{q_id}"))
            cursor.execute("SELECT option_text FROM theory_options WHERE question_id = ? AND is_correct = 1", (q_id,))
            correct = set(row[0] for row in cursor.fetchall())
            answer_text = ", ".join(sorted(selected))
            if selected == correct:
                awarded = marks

        elif q_type == "true_false":
            selected = request.form.get(f"q_{q_id}")
            correction_submitted = request.form.get(f"q_{q_id}_correction", "").strip()
            answer_text = selected or ""
            if correction_submitted:
                answer_text += f" (correction: {correction_submitted})"
            awarded = score_true_false_answer(selected, correction_submitted, options, effective_marks)

        elif q_type == "fill_in":
            answer_text = request.form.get(f"q_{q_id}", "").strip()
            awarded = score_fill_in_answer(answer_text, options, marks)

        elif q_type == "match":
            match_answers = []
            awarded = 0
            for idx, o in enumerate(options, start=1):
                col_a_item = o[1]
                col_b_correct = o[3]
                submitted = request.form.get(f"q_{q_id}_{idx}", "")
                match_answers.append(f"{col_a_item}={submitted}")
                if submitted == col_b_correct:
                    awarded += 1
            answer_text = "; ".join(match_answers)

        score += awarded
        answers_to_save.append((q_id, answer_text, 1 if awarded == effective_marks else 0, awarded))

    percentage = round((score / total) * 100) if total else 0
    time_spent = max(0, int(request.form.get("time_spent_seconds", 0) or 0))

    cursor.execute("""
        INSERT INTO theory_submissions (test_id, username, score, total, percentage, submitted_at, time_spent_seconds, submission_type)
        VALUES (?, ?, ?, ?, ?, ?, ?, ?)
    """, (test_id, username, score, total, percentage, datetime.now().isoformat(), time_spent, "test"))
    submission_id = cursor.lastrowid

    for q_id, answer_text, is_correct, awarded in answers_to_save:
        cursor.execute("""
            INSERT INTO theory_answers (submission_id, question_id, answer_text, is_correct, marks_awarded)
            VALUES (?, ?, ?, ?, ?)
        """, (submission_id, q_id, answer_text, is_correct, awarded))

    conn.commit()
    cursor.execute(
        "UPDATE theory_progress SET completed = 1, current_slide = ?, updated_at = ? WHERE test_id = ? AND username = ?",
        (len(questions_with_options) - 1, datetime.now().isoformat(), test_id, username)
    )
    conn.commit()
    conn.close()
    session.pop(f"test_order_{test_id}", None)
    log_activity(username, f"completed theory test {test_id} — {percentage}%")
    return redirect(url_for("test_results", submission_id=submission_id))


@app.route("/test_results/<int:submission_id>")
def test_results(submission_id):
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))

    conn = get_db()
    cursor = conn.cursor()

    cursor.execute("""
        SELECT s.id, s.score, s.total, s.percentage, s.submitted_at, t.title, t.subject,
               t.show_answers, t.allow_multiple, t.max_attempts, t.id as test_id
        FROM theory_submissions s
        JOIN theory_tests t ON s.test_id = t.id
        WHERE s.id = ? AND s.username = ?
    """, (submission_id, username))
    submission = cursor.fetchone()
    if not submission:
        conn.close()
        return "Results not found", 404

    test_id = submission[10]
    show_answers = submission[7]
    allow_multiple = submission[8]
    max_attempts = submission[9]

    # Count attempts and get best score
    cursor.execute("SELECT COUNT(*), MAX(percentage) FROM theory_submissions WHERE test_id = ? AND username = ?", (test_id, username))
    attempt_row = cursor.fetchone()
    attempts_used = attempt_row[0]
    best_percentage = attempt_row[1]
    can_retry = allow_multiple and attempts_used < max_attempts

    cursor.execute("""
        SELECT q.question_text, q.question_type, q.marks,
               a.answer_text, a.is_correct, a.marks_awarded, a.question_id
        FROM theory_answers a
        JOIN theory_questions q ON a.question_id = q.id
        WHERE a.submission_id = ?
        ORDER BY q.order_index
    """, (submission_id,))
    answers = cursor.fetchall()

    detailed = []
    for ans in answers:
        q_text, q_type, marks, answer_text, is_correct, marks_awarded, q_id = ans
        cursor.execute("""
            SELECT option_text, is_correct, match_pair
            FROM theory_options WHERE question_id = ?
        """, (q_id,))
        options = cursor.fetchall()
        detailed.append({
            "question": q_text,
            "type": q_type,
            "marks": marks,
            "answer": answer_text,
            "correct": is_correct,
            "awarded": marks_awarded,
            "options": options
        })

    conn.close()
    return render_template(
        "test_results.html",
        submission=submission,
        detailed=detailed,
        show_answers=show_answers,
        can_retry=can_retry,
        attempts_used=attempts_used,
        max_attempts=max_attempts,
        best_percentage=best_percentage,
        test_id=test_id
    )


@app.route("/manage_tests/<int:test_id>/reuse", methods=["POST"])
def reuse_test(test_id):
    username = session.get("username")
    if not username:
        return redirect(url_for("login"))
    if get_user_role(username) not in ["teacher", "admin"]:
        return "Access denied", 403

    title = request.form.get("title", "").strip()
    subject = request.form.get("subject", "").strip()
    assign_date = request.form.get("assign_date")
    time_limit = safe_int(request.form.get("time_limit"), 0)
    allow_multiple = 1 if request.form.get("allow_multiple") else 0
    max_attempts = safe_int(request.form.get("max_attempts"), 1)
    show_answers = 1 if request.form.get("show_answers") else 0
    groups = request.form.getlist("groups")

    if not title or not assign_date:
        return redirect(url_for("manage_tests"))

    conn = get_db()
    c = conn.cursor()
    c.execute("SELECT background_image, COALESCE(background_fit, 'cover') FROM theory_tests WHERE id = ?", (test_id,))
    source_bg_row = c.fetchone()
    background_image = source_bg_row[0] if source_bg_row else None
    background_fit = source_bg_row[1] if source_bg_row else "cover"

    # Create the new test
    c.execute("""
        INSERT INTO theory_tests
            (title, subject, assign_date, time_limit, allow_multiple, max_attempts,
             show_answers, background_image, background_fit, created_by, created_at)
        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
    """, (title, subject, assign_date, time_limit, allow_multiple, max_attempts,
           show_answers, background_image, background_fit, username, datetime.now().isoformat()))
    new_test_id = c.lastrowid

    for g in groups:
        if g.strip():
            c.execute("INSERT INTO theory_test_groups (test_id, group_name) VALUES (?, ?)", (new_test_id, g))

    # Fetch all questions and options from source test up front to avoid cursor conflicts
    c.execute("""
        SELECT id, question_text, question_type, marks, order_index
        FROM theory_questions WHERE test_id = ? ORDER BY order_index
    """, (test_id,))
    questions = c.fetchall()

    # Fetch all options for all questions at once
    if questions:
        q_ids = [q[0] for q in questions]
        placeholders = ','.join('?' * len(q_ids))
        c.execute(f"""
            SELECT question_id, option_text, is_correct, match_pair
            FROM theory_options WHERE question_id IN ({placeholders})
            ORDER BY question_id, id
        """, q_ids)
        all_options = c.fetchall()
    else:
        all_options = []

    # Group options by question_id
    from collections import defaultdict
    options_by_q = defaultdict(list)
    for opt in all_options:
        options_by_q[opt[0]].append(opt[1:])

    # Insert copied questions and options
    for q_id, q_text, q_type, marks, order_index in questions:
        c.execute("""
            INSERT INTO theory_questions (test_id, question_text, question_type, marks, order_index)
            VALUES (?, ?, ?, ?, ?)
        """, (new_test_id, q_text, q_type, marks, order_index))
        new_q_id = c.lastrowid
        for opt_text, is_correct, match_pair in options_by_q[q_id]:
            c.execute("""
                INSERT INTO theory_options (question_id, option_text, is_correct, match_pair)
                VALUES (?, ?, ?, ?)
            """, (new_q_id, opt_text, is_correct, match_pair))

    conn.commit()
    conn.close()
    log_activity(username, f"reused test {test_id} as '{title}'")
    return redirect(url_for("manage_test_questions", test_id=new_test_id))


if __name__ == "__main__":
    init_db()
    init_marking_db()
    cleanup = threading.Thread(target=cleanup_thread, daemon=True)
    cleanup.start()
    
    app.run(host=os.getenv("COMPUTERNAME", "127.0.0.1"), port=5000, debug=True)

