import os
import importlib
import base64
import mimetypes
import re
import sqlite3
import zipfile
import uuid
from datetime import datetime, timedelta
from io import BytesIO
from xml.etree import ElementTree as ET
from markupsafe import escape
from .database import get_db


INTERACTIVE_LEARNING_DIR = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "INTERACTIVE LEARNING",
)
LESSON_ASSET_DIR = os.path.join("static", "uploads", "lesson_assets")
MAX_LESSON_IMAGE_DIMENSION = 1600
DATA_URI_IMAGE_RE = re.compile(r"data:(image/[a-zA-Z0-9.+-]+);base64,([A-Za-z0-9+/=\r\n]+)")


def safe_int(value, default=0):
    try:
        return int(value)
    except (TypeError, ValueError):
        return default


def save_lesson_asset(raw_image, mime_type=None):
    os.makedirs(LESSON_ASSET_DIR, exist_ok=True)
    try:
        from PIL import Image

        image = Image.open(BytesIO(raw_image))
        image.thumbnail((MAX_LESSON_IMAGE_DIMENSION, MAX_LESSON_IMAGE_DIMENSION))
        filename = f"{uuid.uuid4().hex}.webp"
        filepath = os.path.join(LESSON_ASSET_DIR, filename)
        if image.mode not in ("RGB", "RGBA"):
            image = image.convert("RGBA" if "A" in image.getbands() else "RGB")
        save_kwargs = {"format": "WEBP", "quality": 82, "method": 6}
        if image.mode == "RGBA":
            save_kwargs["lossless"] = False
        image.save(filepath, **save_kwargs)
    except Exception:
        ext = mimetypes.guess_extension(mime_type or "") or ".png"
        if ext == ".jpe":
            ext = ".jpg"
        filename = f"{uuid.uuid4().hex}{ext}"
        filepath = os.path.join(LESSON_ASSET_DIR, filename)
        with open(filepath, "wb") as asset_file:
            asset_file.write(raw_image)
    return f"/static/uploads/lesson_assets/{filename}"


def save_data_uri_image(data_uri):
    match = DATA_URI_IMAGE_RE.fullmatch(data_uri.strip())
    if not match:
        return data_uri
    mime_type, encoded = match.groups()
    raw_image = base64.b64decode(encoded)
    return save_lesson_asset(raw_image, mime_type)


def externalize_data_uri_images(html_or_uri):
    if not html_or_uri:
        return html_or_uri
    return DATA_URI_IMAGE_RE.sub(lambda match: save_data_uri_image(match.group(0)), html_or_uri)


def parse_module_names(raw_value):
    raw_value = (raw_value or "").replace("\r", "\n")
    parts = []
    for chunk in raw_value.replace(";", ",").split(","):
        for line in chunk.split("\n"):
            cleaned = line.strip()
            if cleaned:
                parts.append(cleaned)
    seen = set()
    ordered = []
    for item in parts:
        key = item.lower()
        if key in seen:
            continue
        seen.add(key)
        ordered.append(item)
    return ordered


def normalize_question_bank_group_text(question_text):
    text = " ".join((question_text or "").strip().split())
    if not text:
        return ""
    text = re.sub(r"\s*\((?:case|scenario|batch)\s+[^)]*\)\s*$", "", text, flags=re.IGNORECASE)
    return text.strip().lower()


# ── Term helpers ──────────────────────────────────────────────────────────────

def get_term_dates():
    conn = get_db()
    cursor = conn.cursor()
    try:
        cursor.execute("SELECT term, start_date, end_date FROM term_dates ORDER BY term")
        rows = cursor.fetchall()
    except Exception:
        rows = []
    conn.close()
    terms = {i: {"start": None, "end": None} for i in range(1, 5)}
    for term, start, end in rows:
        terms[term] = {"start": start, "end": end}
    return terms


def get_active_term_range():
    today = datetime.now().date().isoformat()
    for t in get_term_dates().values():
        if t["start"] and t["end"] and t["start"] <= today <= t["end"]:
            return t["start"], t["end"]
    return None


def get_all_term_days():
    days = set()
    for t in get_term_dates().values():
        if t["start"] and t["end"]:
            current = datetime.strptime(t["start"], "%Y-%m-%d").date()
            end = min(datetime.strptime(t["end"], "%Y-%m-%d").date(), datetime.now().date())
            while current <= end:
                if current.weekday() < 5:
                    days.add(current.strftime("%Y-%m-%d"))
                current += timedelta(days=1)
    return days


def get_last_21_days():
    days = []
    current = datetime.now().date()
    term_days = get_all_term_days()
    limit = 0
    while len(days) < 21 and limit < 365:
        s = current.strftime("%Y-%m-%d")
        if current.weekday() < 5 and (not term_days or s in term_days):
            days.append(s)
        current -= timedelta(days=1)
        limit += 1
    return list(reversed(days))


def get_last_7_days():
    days = []
    current = datetime.now().date()
    term_days = get_all_term_days()
    limit = 0
    while len(days) < 7 and limit < 365:
        s = current.strftime("%Y-%m-%d")
        if current.weekday() < 5 and (not term_days or s in term_days):
            days.append(s)
        current -= timedelta(days=1)
        limit += 1
    return list(reversed(days))


def get_current_year_attendance_days():
    days = []
    current = datetime(datetime.now().year, 1, 1).date()
    today = datetime.now().date()
    while current <= today:
        if current.weekday() < 5:
            days.append(current.strftime("%Y-%m-%d"))
        current += timedelta(days=1)
    return days


def get_days_in_active_term():
    term_range = get_active_term_range()
    if not term_range:
        return []
    start_date, end_date = term_range
    days = []
    current = datetime.strptime(start_date, "%Y-%m-%d").date()
    end = datetime.strptime(end_date, "%Y-%m-%d").date()
    while current <= end:
        if current.weekday() < 5:
            days.append(current.strftime("%Y-%m-%d"))
        current += timedelta(days=1)
    return days


def normalize_review_text(value):
    return " ".join((value or "").strip().lower().split())


def parse_true_false_answer_text(answer_text):
    raw = (answer_text or "").strip()
    if "(correction:" in raw:
        selected, correction = raw.split("(correction:", 1)
        return selected.strip(), correction.rstrip(") ").strip()
    return raw, ""


def get_true_false_option_data(options):
    correct_choice = ""
    accepted_corrections = []
    for option in options:
        option_text = option[1] if len(option) > 1 else ""
        is_correct = option[2] if len(option) > 2 else 0
        match_pair = option[3] if len(option) > 3 else None
        if is_correct == 1 and match_pair != "correction":
            correct_choice = option_text
        if match_pair == "correction" and option_text:
            accepted_corrections.append(option_text)
    return correct_choice, accepted_corrections


def score_true_false_answer(selected, correction_submitted, options, effective_marks):
    correct_choice, accepted_corrections = get_true_false_option_data(options)
    if selected != correct_choice:
        return 0
    if selected == "False" and accepted_corrections:
        submitted_key = normalize_review_text(correction_submitted)
        accepted_keys = {normalize_review_text(item) for item in accepted_corrections if item}
        return effective_marks if submitted_key and submitted_key in accepted_keys else 0
    return effective_marks


def get_fill_in_accepted_answers(options):
    return [option[1] for option in options if len(option) > 2 and option[2] == 1 and option[1]]


def score_fill_in_answer(answer_text, options, marks):
    submitted_key = normalize_review_text(answer_text)
    accepted_keys = {normalize_review_text(item) for item in get_fill_in_accepted_answers(options)}
    return marks if submitted_key and submitted_key in accepted_keys else 0


def compute_theory_answer_award(question_type, marks, options, answer_text):
    if question_type == "mcq_single":
        correct_values = {option[1] for option in options if len(option) > 2 and option[2] == 1}
        return marks if answer_text in correct_values else 0
    if question_type == "mcq_multi":
        correct_values = sorted(option[1] for option in options if len(option) > 2 and option[2] == 1)
        selected_values = sorted([item.strip() for item in (answer_text or "").split(",") if item.strip()])
        return marks if selected_values == correct_values else 0
    if question_type == "true_false":
        selected, correction = parse_true_false_answer_text(answer_text)
        return score_true_false_answer(selected, correction, options, marks)
    if question_type == "fill_in":
        return score_fill_in_answer(answer_text, options, marks)
    if question_type == "match":
        learner_map = {}
        for pair in (answer_text or "").split(";"):
            if "=" not in pair:
                continue
            left, chosen = pair.split("=", 1)
            learner_map[left.strip()] = chosen.strip()
        indexed_awarded = 0
        indexed_present = False
        for idx, option in enumerate(options, start=1):
            if len(option) <= 3 or not option[3] or option[3] == "correction":
                continue
            if str(idx) in learner_map:
                indexed_present = True
                if learner_map.get(str(idx), "") == option[1]:
                    indexed_awarded += 1
        if indexed_present:
            return indexed_awarded
        legacy_map = {option[1]: option[3] for option in options if len(option) > 3 and option[3] and option[3] != "correction"}
        swapped_map = {option[3]: option[1] for option in options if len(option) > 3 and option[3] and option[3] != "correction"}
        legacy_awarded = sum(1 for left, accepted in legacy_map.items() if learner_map.get(left, "") == accepted)
        swapped_awarded = sum(1 for left, accepted in swapped_map.items() if learner_map.get(left, "") == accepted)
        return max(legacy_awarded, swapped_awarded)
    return 0


def build_bank_option_signature(question_type, options):
    normalized = []
    for option_text, is_correct, match_pair in options:
        normalized.append((
            (option_text or "").strip().lower(),
            safe_int(is_correct, 0),
            (match_pair or "").strip().lower(),
        ))
    if question_type == "match":
        normalized.sort()
    return tuple(normalized)


def bank_question_exists(cursor, question_text, question_type, subject, modules, options):
    normalized_text = normalize_question_bank_group_text(question_text)
    normalized_subject = (subject or "").strip().lower()
    normalized_modules = (modules or "").strip().lower()
    target_signature = build_bank_option_signature(question_type, options)

    cursor.execute("""
        SELECT id, question_text
        FROM question_bank_questions
        WHERE question_type = ?
          AND LOWER(TRIM(COALESCE(subject, ''))) = ?
          AND LOWER(TRIM(COALESCE(modules, ''))) = ?
    """, (question_type, normalized_subject, normalized_modules))
    candidate_rows = cursor.fetchall()
    for candidate_id, candidate_text in candidate_rows:
        if normalize_question_bank_group_text(candidate_text) != normalized_text:
            continue
        cursor.execute("""
            SELECT option_text, is_correct, match_pair
            FROM question_bank_options
            WHERE bank_question_id = ?
            ORDER BY id
        """, (candidate_id,))
        candidate_signature = build_bank_option_signature(question_type, cursor.fetchall())
        if candidate_signature == target_signature:
            return True
    return False


def normalize_relative_path(path):
    return path.replace("/", os.sep).replace("\\", os.sep).strip()


def get_interactive_learning_files():
    files = []
    if not os.path.isdir(INTERACTIVE_LEARNING_DIR):
        return files

    for root, _, filenames in os.walk(INTERACTIVE_LEARNING_DIR):
        for filename in filenames:
            lower = filename.lower()
            if lower.startswith("~$"):
                continue
            if lower.endswith((".ppt", ".pptx", ".pps", ".ppsx")):
                full_path = os.path.join(root, filename)
                rel_path = os.path.relpath(full_path, INTERACTIVE_LEARNING_DIR)
                files.append({
                    "relative_path": rel_path.replace("\\", "/"),
                    "display_name": os.path.splitext(filename)[0],
                    "folder": os.path.dirname(rel_path).replace("\\", "/"),
                })

    files.sort(key=lambda item: (item["folder"].lower(), item["display_name"].lower()))
    return files


def resolve_interactive_learning_path(relative_path):
    if not relative_path:
        return None

    normalized = os.path.normpath(os.path.join(INTERACTIVE_LEARNING_DIR, normalize_relative_path(relative_path)))
    base_dir = os.path.normpath(INTERACTIVE_LEARNING_DIR)

    if not normalized.startswith(base_dir):
        return None
    if not os.path.isfile(normalized):
        return None
    return normalized


def extract_pptx_slides(file_path):
    lower = file_path.lower()
    if not lower.endswith((".pptx", ".ppsx")):
        return []

    namespace = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    slides = []

    with zipfile.ZipFile(file_path) as archive:
        slide_names = [
            name for name in archive.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        ]
        slide_names.sort(key=lambda item: int(re.search(r"slide(\d+)\.xml$", item).group(1)))

        for index, slide_name in enumerate(slide_names, start=1):
            root = ET.fromstring(archive.read(slide_name))
            text_runs = [
                (node.text or "").strip()
                for node in root.findall(".//a:t", namespace)
                if (node.text or "").strip()
            ]
            slides.append({
                "number": index,
                "title": text_runs[0] if text_runs else f"Slide {index}",
                "text_runs": text_runs,
                "body": text_runs[1:] if len(text_runs) > 1 else [],
            })

    return slides


def get_low_attendance_learners(limit=10):
    days = get_last_21_days()
    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("SELECT date FROM excluded_dates WHERE group_name IS NULL")
    excluded = {row[0] for row in cursor.fetchall()}
    days = [d for d in days if d not in excluded]

    cursor.execute("SELECT username, full_name FROM users WHERE role = 'student'")
    learners = cursor.fetchall()
    results = []
    for username, full_name in learners:
        present = 0
        for day in days:
            cursor.execute("SELECT 1 FROM login_history WHERE username = ? AND date = ?", (username, day))
            if cursor.fetchone():
                present += 1
                continue
            cursor.execute("SELECT status FROM attendance_override WHERE username = ? AND date = ?", (username, day))
            override = cursor.fetchone()
            if override and override[0] == "present":
                present += 1
        results.append((full_name or username, len(days) - present))

    conn.close()
    results.sort(key=lambda x: x[1], reverse=True)
    return results[:limit]


def attendance_status_counts_as_present(status):
    return (status or "").strip().lower() in {"present", "late"}


def fetch_first_login_times(cursor, usernames, days):
    login_map = {}
    if not usernames or not days:
        return login_map
    user_placeholders = ",".join("?" for _ in usernames)
    day_placeholders = ",".join("?" for _ in days)
    cursor.execute(
        f"""
        SELECT username, date, MIN(login_time)
        FROM login_history
        WHERE username IN ({user_placeholders}) AND date IN ({day_placeholders})
        GROUP BY username, date
        """,
        list(usernames) + list(days),
    )
    for username, day, login_time in cursor.fetchall():
        login_map[(username, day)] = login_time
    return login_map


def fetch_attendance_override_statuses(cursor, usernames, days):
    override_map = {}
    if not usernames or not days:
        return override_map
    user_placeholders = ",".join("?" for _ in usernames)
    day_placeholders = ",".join("?" for _ in days)
    cursor.execute(
        f"""
        SELECT username, date, status
        FROM attendance_override
        WHERE username IN ({user_placeholders}) AND date IN ({day_placeholders})
        """,
        list(usernames) + list(days),
    )
    for username, day, status in cursor.fetchall():
        override_map[(username, day)] = status
    return override_map


def fetch_group_late_thresholds(cursor, group, days):
    late_cutoffs = {day: None for day in days}
    if not group or not days:
        return late_cutoffs
    day_placeholders = ",".join("?" for _ in days)
    cursor.execute(
        f"""
        SELECT lh.date, MIN(lh.login_time)
        FROM login_history lh
        JOIN users u ON u.username = lh.username
        WHERE u.group_name = ? AND lh.date IN ({day_placeholders})
        GROUP BY lh.date
        """,
        [group, *days],
    )
    for day, min_login in cursor.fetchall():
        if not min_login:
            continue
        try:
            min_dt = datetime.fromisoformat(min_login)
        except ValueError:
            min_dt = datetime.strptime(min_login, "%Y-%m-%d %H:%M:%S")
        late_cutoffs[day] = (min_dt + timedelta(minutes=4)).strftime("%H:%M")
    return late_cutoffs


def fetch_group_excluded_dates(cursor, groups):
    excluded_by_group = {group: set() for group in groups if group}
    if not excluded_by_group:
        return excluded_by_group
    placeholders = ",".join("?" for _ in excluded_by_group)
    cursor.execute(
        f"""
        SELECT date, group_name
        FROM excluded_dates
        WHERE group_name IS NULL OR group_name IN ({placeholders})
        """,
        list(excluded_by_group.keys()),
    )
    rows = cursor.fetchall()
    global_dates = {date_str for date_str, group_name in rows if group_name is None}
    for group in excluded_by_group:
        excluded_by_group[group] = set(global_dates)
    for date_str, group_name in rows:
        if group_name in excluded_by_group:
            excluded_by_group[group_name].add(date_str)
    return excluded_by_group


def fetch_present_day_map(cursor, usernames, days):
    present_map = {username: set() for username in usernames}
    if not present_map or not days:
        return present_map
    user_placeholders = ",".join("?" for _ in present_map)
    day_placeholders = ",".join("?" for _ in days)
    cursor.execute(
        f"""
        SELECT DISTINCT username, date
        FROM login_history
        WHERE username IN ({user_placeholders}) AND date IN ({day_placeholders})
        """,
        list(present_map.keys()) + list(days),
    )
    for username, day in cursor.fetchall():
        present_map.setdefault(username, set()).add(day)
    cursor.execute(
        f"""
        SELECT username, date, status
        FROM attendance_override
        WHERE username IN ({user_placeholders}) AND date IN ({day_placeholders})
        """,
        list(present_map.keys()) + list(days),
    )
    for username, day, status in cursor.fetchall():
        if attendance_status_counts_as_present(status):
            present_map.setdefault(username, set()).add(day)
        else:
            present_map.setdefault(username, set()).discard(day)
    return present_map


def build_attendance_history(cursor, username, group_name, days, login_map=None, override_map=None, late_cutoffs=None, class_checked_dates=None, excluded_dates=None):
    history = []
    login_map = login_map or {}
    override_map = override_map or {}
    late_cutoffs = late_cutoffs or {}
    class_checked_dates = class_checked_dates or set()
    excluded_dates = excluded_dates or set()

    for day in days:
        weekday = datetime.strptime(day, "%Y-%m-%d").weekday()
        login_time = login_map.get((username, day))
        if day in excluded_dates:
            history.append({"date": day, "status": "Normal", "time": "", "late": False, "note": "", "weekday": weekday})
            continue

        if login_time:
            time_str = login_time.split(" ")[1][:5]
            cutoff = late_cutoffs.get(day)
            late = cutoff is not None and time_str > cutoff
            history.append({"date": day, "status": "Present", "time": time_str, "late": late, "note": "Auto", "weekday": weekday})
            continue

        override_status = (override_map.get((username, day)) or "").strip().lower()
        if attendance_status_counts_as_present(override_status):
            manual_time = "12:00" if override_status == "present" else ""
            history.append({"date": day, "status": "Present", "time": manual_time, "late": override_status == "late", "note": "Manual", "weekday": weekday})
        elif override_status == "absent" or day in class_checked_dates:
            note = "Manual" if override_status == "absent" else "Class checked in"
            history.append({"date": day, "status": "Absent", "time": "", "late": False, "note": note, "weekday": weekday})
        else:
            history.append({"date": day, "status": "Normal", "time": "", "late": False, "note": "", "weekday": weekday})

    return history


def summarize_attendance_history(history):
    counted_history = [item for item in history if item["status"] in ("Present", "Absent")]
    total_days = len(counted_history)
    present_days = sum(1 for item in history if item["status"] == "Present")
    absent_days = sum(1 for item in history if item["status"] == "Absent")
    late_days = sum(1 for item in history if item["late"])
    attendance_pct = round((present_days / total_days) * 100) if total_days else 0
    return {
        "total_days": total_days,
        "present_days": present_days,
        "absent_days": absent_days,
        "late_days": late_days,
        "attendance_pct": attendance_pct,
    }


def build_attendance_group_summary(cursor, groups, teacher_username=None, days=None):
    summaries = []
    days = days or get_last_21_days()
    for group_name in groups:
        if not group_name:
            continue
        if teacher_username:
            cursor.execute("""
                SELECT username, full_name
                FROM users
                WHERE group_name = ? AND role = 'student' AND teacher_username = ?
                ORDER BY full_name
            """, (group_name, teacher_username))
        else:
            cursor.execute("""
                SELECT username, full_name
                FROM users
                WHERE group_name = ? AND role = 'student'
                ORDER BY full_name
            """, (group_name,))
        learners = cursor.fetchall()
        usernames = [row[0] for row in learners]
        auto_exclude_empty_attendance_days(cursor, [group_name], created_by=teacher_username or "system", days=get_current_year_attendance_days())
        excluded_dates = fetch_group_excluded_dates(cursor, [group_name]).get(group_name, set())
        filtered_days = [day for day in days if day not in excluded_dates]
        login_map = fetch_first_login_times(cursor, usernames, filtered_days)
        override_map = fetch_attendance_override_statuses(cursor, usernames, filtered_days)
        late_cutoffs = fetch_group_late_thresholds(cursor, group_name, filtered_days)
        cursor.execute("""
            SELECT DISTINCT lh.date
            FROM login_history lh
            JOIN users u ON u.username = lh.username
            WHERE u.group_name = ? AND u.role = 'student'
        """, (group_name,))
        class_checked_dates = {row[0] for row in cursor.fetchall()}

        present_total = 0
        absent_total = 0
        late_total = 0
        counted_total = 0
        for uname in usernames:
            history = build_attendance_history(
                cursor,
                uname,
                group_name,
                filtered_days,
                login_map=login_map,
                override_map=override_map,
                late_cutoffs=late_cutoffs,
                class_checked_dates=class_checked_dates,
                excluded_dates=excluded_dates,
            )
            summary = summarize_attendance_history(history)
            present_total += summary["present_days"]
            absent_total += summary["absent_days"]
            late_total += summary["late_days"]
            counted_total += summary["total_days"]

        att_pct = round((present_total / counted_total) * 100) if counted_total else 0
        summaries.append({
            "group": group_name,
            "students": len(learners),
            "attendance_pct": att_pct,
            "present_days": present_total,
            "absent_days": absent_total,
            "late_days": late_total,
        })
    return summaries


def build_attendance_months(history):
    attendance_months = []
    for item in history:
        month_key = item["date"][:7]
        if not attendance_months or attendance_months[-1]["key"] != month_key:
            month_date = datetime.strptime(item["date"], "%Y-%m-%d")
            attendance_months.append({
                "key": month_key,
                "name": month_date.strftime("%B"),
                "days": []
            })
        attendance_months[-1]["days"].append(item)
    return attendance_months


def auto_exclude_empty_attendance_days(cursor, groups, created_by="system", days=None):
    groups = [group for group in groups if group]
    if not groups:
        return 0

    if days is None:
        days = get_current_year_attendance_days()
    today = datetime.now().date()
    filtered_days = []
    for day in days:
        try:
            day_obj = datetime.strptime(day, "%Y-%m-%d").date()
        except ValueError:
            continue
        if day_obj >= today:
            continue
        if day_obj.weekday() >= 5:
            continue
        filtered_days.append(day)

    days = filtered_days
    if not days:
        return 0

    group_placeholders = ",".join("?" for _ in groups)
    day_placeholders = ",".join("?" for _ in days)

    cursor.execute(
        f"""
        SELECT u.group_name, lh.date, COUNT(DISTINCT lh.username)
        FROM login_history lh
        JOIN users u ON u.username = lh.username
        WHERE u.group_name IN ({group_placeholders})
          AND u.role = 'student'
          AND lh.date IN ({day_placeholders})
        GROUP BY u.group_name, lh.date
        """,
        [*groups, *days],
    )
    login_counts = {(group_name, day): count for group_name, day, count in cursor.fetchall()}

    cursor.execute(
        f"""
        SELECT u.group_name, ao.date, COUNT(DISTINCT ao.username)
        FROM attendance_override ao
        JOIN users u ON u.username = ao.username
        WHERE u.group_name IN ({group_placeholders})
          AND u.role = 'student'
          AND ao.date IN ({day_placeholders})
          AND LOWER(COALESCE(ao.status, '')) IN ('present', 'late')
        GROUP BY u.group_name, ao.date
        """,
        [*groups, *days],
    )
    override_counts = {(group_name, day): count for group_name, day, count in cursor.fetchall()}

    cursor.execute(
        f"""
        SELECT date, group_name
        FROM excluded_dates
        WHERE (group_name IN ({group_placeholders}) OR group_name IS NULL)
          AND date IN ({day_placeholders})
        """,
        [*groups, *days],
    )
    existing = {(date_str, group_name) for date_str, group_name in cursor.fetchall()}

    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    inserted = 0
    for group_name in groups:
        for day in days:
            if (day, group_name) in existing:
                continue
            if (day, None) in existing:
                continue
            if login_counts.get((group_name, day), 0) > 0:
                continue
            if override_counts.get((group_name, day), 0) > 0:
                continue
            cursor.execute(
                """
                INSERT OR IGNORE INTO excluded_dates (date, group_name, reason, created_by, created_at)
                VALUES (?, ?, ?, ?, ?)
                """,
                (day, group_name, "Auto excluded: no class attendance recorded", created_by, timestamp),
            )
            if cursor.rowcount:
                inserted += 1
    return inserted


def fetch_student_practical_averages(cursor, teacher_username=None):
    params = []
    teacher_filter = ""
    if teacher_username:
        teacher_filter = " AND u.teacher_username = ?"
        params.append(teacher_username)
    cursor.execute(
        f"""
        SELECT best_scores.username, ROUND(AVG(best_scores.best_score), 1)
        FROM (
            SELECT username, subject, task, MAX(score) AS best_score
            FROM results
            GROUP BY username, subject, task
        ) best_scores
        JOIN users u ON u.username = best_scores.username
        WHERE u.role = 'student'{teacher_filter}
        GROUP BY best_scores.username
        """,
        params,
    )
    return {username: average for username, average in cursor.fetchall()}


def fetch_student_theory_averages(cursor, teacher_username=None):
    params = []
    teacher_filter = ""
    if teacher_username:
        teacher_filter = " AND u.teacher_username = ?"
        params.append(teacher_username)
    cursor.execute(
        f"""
        SELECT best_scores.username, ROUND(AVG(best_scores.best_pct), 1)
        FROM (
            SELECT username, test_id, MAX(percentage) AS best_pct
            FROM theory_submissions
            GROUP BY username, test_id
        ) best_scores
        JOIN users u ON u.username = best_scores.username
        WHERE u.role = 'student'{teacher_filter}
        GROUP BY best_scores.username
        """,
        params,
    )
    return {username: average for username, average in cursor.fetchall()}


def fetch_group_practical_averages(cursor, groups):
    if not groups:
        return {}
    placeholders = ",".join("?" for _ in groups)
    cursor.execute(
        f"""
        SELECT u.group_name, ROUND(AVG(best_scores.best_score), 1)
        FROM (
            SELECT username, subject, task, MAX(score) AS best_score
            FROM results
            GROUP BY username, subject, task
        ) best_scores
        JOIN users u ON u.username = best_scores.username
        WHERE u.role = 'student' AND u.group_name IN ({placeholders})
        GROUP BY u.group_name
        """,
        groups,
    )
    return {group: average for group, average in cursor.fetchall()}


def fetch_group_theory_averages(cursor, groups):
    if not groups:
        return {}
    placeholders = ",".join("?" for _ in groups)
    cursor.execute(
        f"""
        SELECT u.group_name, ROUND(AVG(best_scores.best_pct), 1)
        FROM (
            SELECT username, test_id, MAX(percentage) AS best_pct
            FROM theory_submissions
            GROUP BY username, test_id
        ) best_scores
        JOIN users u ON u.username = best_scores.username
        WHERE u.role = 'student' AND u.group_name IN ({placeholders})
        GROUP BY u.group_name
        """,
        groups,
    )
    return {group: average for group, average in cursor.fetchall()}


def fetch_theory_module_weaknesses(cursor, username=None, group_name=None, limit=10):
    where_parts = ["latest.rn = 1"]
    params = []
    if username:
        where_parts.append("latest.username = ?")
        params.append(username)
    if group_name:
        where_parts.append("u.group_name = ?")
        params.append(group_name)
    where_clause = " AND ".join(where_parts)
    cursor.execute(f"""
        SELECT latest.username, tq.source_modules, tq.question_type, ta.is_correct
        FROM (
            SELECT ts.id, ts.username, ts.test_id,
                   ROW_NUMBER() OVER (
                       PARTITION BY ts.username, ts.test_id
                       ORDER BY COALESCE(ts.submitted_at, '') DESC, ts.id DESC
                   ) AS rn
            FROM theory_submissions ts
        ) latest
        JOIN theory_answers ta ON ta.submission_id = latest.id
        JOIN theory_questions tq ON tq.id = ta.question_id
        JOIN users u ON u.username = latest.username
        WHERE {where_clause}
    """, params)
    summary = {}
    for learner_username, source_modules, question_type, is_correct in cursor.fetchall():
        modules = parse_module_names(source_modules or "") or ["Unmapped"]
        for module_name in modules:
            item = summary.setdefault(module_name, {
                "asked": 0,
                "wrong": 0,
                "learners": set(),
                "question_types": set(),
            })
            item["asked"] += 1
            item["wrong"] += 0 if is_correct else 1
            item["learners"].add(learner_username)
            if question_type:
                item["question_types"].add(question_type)
    rows = []
    for module_name, item in summary.items():
        if not item["asked"]:
            continue
        rows.append({
            "module": module_name,
            "asked": item["asked"],
            "wrong": item["wrong"],
            "wrong_pct": round((item["wrong"] / item["asked"]) * 100),
            "learners": len(item["learners"]),
            "question_types": ", ".join(sorted(item["question_types"])),
        })
    rows.sort(key=lambda row: (row["wrong_pct"], row["wrong"], row["asked"]), reverse=True)
    return rows[:limit]


def get_low_attendance_learners_filtered(limit=10, groups=None, teacher_username=None):
    days = get_last_21_days()
    conn = get_db()
    cursor = conn.cursor()

    if teacher_username:
        cursor.execute("""
            SELECT username, full_name, group_name
            FROM users
            WHERE role = 'student' AND teacher_username = ?
            ORDER BY group_name, full_name, username
        """, (teacher_username,))
    elif groups:
        placeholders = ",".join("?" for _ in groups)
        cursor.execute(f"""
            SELECT username, full_name, group_name
            FROM users
            WHERE role = 'student' AND group_name IN ({placeholders})
            ORDER BY group_name, full_name, username
        """, groups)
    else:
        cursor.execute("""
            SELECT username, full_name, group_name
            FROM users
            WHERE role = 'student'
            ORDER BY group_name, full_name, username
        """)
    learners = cursor.fetchall()

    learner_groups = sorted({row[2] for row in learners if row[2]})
    excluded_by_group = fetch_group_excluded_dates(cursor, learner_groups)
    login_map = fetch_first_login_times(cursor, [row[0] for row in learners], days)
    override_map = fetch_attendance_override_statuses(cursor, [row[0] for row in learners], days)
    late_cutoffs_by_group = {
        group_name: fetch_group_late_thresholds(cursor, group_name, days)
        for group_name in learner_groups
    }
    class_checked_by_group = {}
    for group_name in learner_groups:
        cursor.execute("""
            SELECT DISTINCT lh.date
            FROM login_history lh
            JOIN users u ON u.username = lh.username
            WHERE u.group_name = ? AND u.role = 'student'
        """, (group_name,))
        class_checked_by_group[group_name] = {row[0] for row in cursor.fetchall()}

    results = []
    for username, full_name, group_name in learners:
        excluded_dates = excluded_by_group.get(group_name, set())
        group_days = [day for day in days if day not in excluded_dates]
        history = build_attendance_history(
            cursor,
            username,
            group_name,
            group_days,
            login_map=login_map,
            override_map=override_map,
            late_cutoffs=late_cutoffs_by_group.get(group_name, {}),
            class_checked_dates=class_checked_by_group.get(group_name, set()),
            excluded_dates=excluded_dates,
        )
        absent = summarize_attendance_history(history)["absent_days"]
        results.append((full_name or username, username, absent))

    conn.close()
    results.sort(key=lambda x: x[2], reverse=True)
    return results[:limit]


def get_teacher_quick_action_catalog():
    return [
        {"key": "manage_subjects", "label": "Practical", "icon": "📁", "href": "/manage_subjects", "kind": "link"},
        {"key": "manage_lessons", "label": "Lesson Setup", "icon": "📘", "href": "/manage_lessons", "kind": "link"},
        {"key": "manage_tests", "label": "Theory Tests", "icon": "📝", "href": "/manage_tests", "kind": "link"},
        {"key": "response_review", "label": "Review Responses", "icon": "🧾", "href": "/response_review", "kind": "link"},
        {"key": "marking_setup", "label": "Marking Setup", "icon": "🛠️", "href": "/marking_setup", "kind": "link"},
        {"key": "attendance", "label": "Attendance", "icon": "📅", "href": "/attendance", "kind": "link"},
        {"key": "group_results", "label": "Results", "icon": "📊", "href": "/group_results", "kind": "link"},
        {"key": "communications", "label": "Messages", "icon": "💬", "href": "/communications", "kind": "link"},
        {"key": "export", "label": "Export", "icon": "⬇️", "kind": "export"},
        {"key": "view_group", "label": "View Group", "icon": "👁", "kind": "view_group"},
    ]


def get_teacher_selected_quick_actions(username):
    catalog = get_teacher_quick_action_catalog()
    catalog_map = {item["key"]: item for item in catalog}
    default_keys = ["manage_subjects", "manage_tests", "attendance", "communications", "response_review", "export", "view_group"]
    conn = get_db()
    cursor = conn.cursor()
    cursor.execute(
        """
        SELECT action_key
        FROM teacher_quick_actions
        WHERE username = ?
        ORDER BY rowid
        """,
        (username,),
    )
    stored_keys = [row[0] for row in cursor.fetchall()]
    conn.close()
    selected = [catalog_map[key] for key in stored_keys if key in catalog_map]
    if not selected:
        selected = [catalog_map[key] for key in default_keys if key in catalog_map]
    selected_keys = {item["key"] for item in selected}
    available = [dict(item, selected=item["key"] in selected_keys) for item in catalog]
    return selected, available


def get_teacher_unread_message_count(username):
    conn = get_db()
    cursor = conn.cursor()
    cursor.execute(
        """
        SELECT COUNT(*)
        FROM communication_threads t
        WHERE t.teacher_username = ?
          AND EXISTS (
              SELECT 1
              FROM communication_messages m
              WHERE m.thread_id = t.id
                AND COALESCE(m.sender_role, '') = 'student'
                AND (t.teacher_read_at IS NULL OR COALESCE(m.created_at, '') > t.teacher_read_at)
          )
        """,
        (username,),
    )
    count = cursor.fetchone()[0] or 0
    conn.close()
    return count


def get_student_unread_message_count(username):
    conn = get_db()
    cursor = conn.cursor()
    cursor.execute(
        """
        SELECT COUNT(*)
        FROM communication_threads t
        WHERE t.student_username = ?
          AND EXISTS (
              SELECT 1
              FROM communication_messages m
              WHERE m.thread_id = t.id
                AND COALESCE(m.sender_role, '') IN ('teacher', 'admin')
                AND (t.student_read_at IS NULL OR COALESCE(m.created_at, '') > t.student_read_at)
          )
        """,
        (username,),
    )
    count = cursor.fetchone()[0] or 0
    conn.close()
    return count


def student_has_fresh_teacher_reply(username):
    return get_student_unread_message_count(username) > 0


def get_student_message_threads(username):
    conn = get_db()
    conn.row_factory = sqlite3.Row
    cursor = conn.cursor()
    cursor.execute("""
        SELECT *
        FROM communication_threads
        WHERE student_username = ?
        ORDER BY COALESCE(updated_at, created_at) DESC, id DESC
    """, (username,))
    threads = []
    for row in cursor.fetchall():
        thread = dict(row)
        cursor.execute("""
            SELECT *
            FROM communication_messages
            WHERE thread_id = ?
            ORDER BY COALESCE(created_at, '') ASC, id ASC
        """, (thread["id"],))
        thread["messages"] = [dict(message_row) for message_row in cursor.fetchall()]
        threads.append(thread)
    conn.close()
    return threads


def mark_student_threads_read(username):
    now_text = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("""
        UPDATE communication_threads
        SET student_read_at = ?
        WHERE student_username = ?
    """, (now_text, username))
    conn.commit()
    conn.close()


def create_communication_thread(cursor, student_username, topic, subject_line="", attendance_date="", initial_message="", chat_session_id=""):
    created_at = datetime.now().isoformat()
    cursor.execute("SELECT full_name, group_name, teacher_username FROM users WHERE username = ?", (student_username,))
    user_row = cursor.fetchone()
    group_name = user_row[1] if user_row else ""
    teacher_username = user_row[2] if user_row else ""
    cursor.execute(
        """
        INSERT INTO communication_threads (
            student_username, teacher_username, group_name, topic, subject_line,
            attendance_date, status, chat_session_id, created_at, updated_at
        )
        VALUES (?, ?, ?, ?, ?, ?, 'open', ?, ?, ?)
        """,
        (student_username, teacher_username, group_name, topic, subject_line, attendance_date, chat_session_id, created_at, created_at),
    )
    thread_id = cursor.lastrowid
    if initial_message:
        add_communication_message(cursor, thread_id, student_username, "student", initial_message)
    return thread_id


def add_communication_message(cursor, thread_id, sender_username, sender_role, message):
    created_at = datetime.now().isoformat()
    cursor.execute(
        """
        INSERT INTO communication_messages (thread_id, sender_username, sender_role, message, created_at)
        VALUES (?, ?, ?, ?, ?)
        """,
        (thread_id, sender_username, sender_role, message, created_at),
    )
    cursor.execute(
        """
        UPDATE communication_threads
        SET updated_at = ?,
            teacher_read_at = CASE WHEN ? IN ('teacher', 'admin') THEN ? ELSE teacher_read_at END,
            student_read_at = CASE WHEN ? = 'student' THEN ? ELSE student_read_at END
        WHERE id = ?
        """,
        (created_at, sender_role, created_at, sender_role, created_at, thread_id),
    )


QUESTION_BANK_SUPPORTED_TYPES = ["mcq_single", "fill_in", "true_false", "match"]


def get_question_bank_counts(cursor, modules=None, subjects=None):
    modules = parse_module_names(",".join(modules or []))
    subjects = [item.strip() for item in (subjects or []) if item and item.strip()]
    counts = {}
    for q_type in QUESTION_BANK_SUPPORTED_TYPES:
        where_parts = ["question_type = ?"]
        params = [q_type]
        if modules:
            module_filters = " OR ".join(["LOWER(COALESCE(modules, '')) LIKE ?" for _ in modules])
            where_parts.append(f"({module_filters})")
            params.extend(f"%{module.lower()}%" for module in modules)
        if subjects:
            subject_filters = " OR ".join(["LOWER(COALESCE(subject, '')) = ?" for _ in subjects])
            where_parts.append(f"({subject_filters})")
            params.extend(subject.lower() for subject in subjects)
        cursor.execute(f"""
            SELECT COUNT(*)
            FROM question_bank_questions
            WHERE {' AND '.join(where_parts)}
        """, params)
        counts[q_type] = cursor.fetchone()[0] or 0
    return counts


def build_match_review_rows(options, answer_text):
    learner_map = {}
    for pair in (answer_text or "").split(";"):
        if "=" in pair:
            left, chosen = pair.split("=", 1)
            learner_map[left.strip()] = chosen.strip()
    indexed_rows = []
    for idx, option in enumerate(options, start=1):
        if len(option) <= 3 or not option[3] or option[3] == "correction":
            continue
        if str(idx) in learner_map:
            indexed_rows.append({
                "left": option[3],
                "learner_match": learner_map.get(str(idx), "") or "No answer",
                "correct_match": option[1],
                "is_correct": learner_map.get(str(idx), "") == option[1]
            })
    if indexed_rows:
        return indexed_rows, {row["left"]: row["correct_match"] for row in indexed_rows}
    legacy_map = {option[1]: option[3] for option in options if len(option) > 3 and option[3] and option[3] != "correction"}
    swapped_map = {option[3]: option[1] for option in options if len(option) > 3 and option[3] and option[3] != "correction"}
    legacy_score = sum(1 for left, accepted in legacy_map.items() if learner_map.get(left, "") == accepted)
    swapped_score = sum(1 for left, accepted in swapped_map.items() if learner_map.get(left, "") == accepted)
    active_map = swapped_map if swapped_score > legacy_score else legacy_map
    rows = []
    for left, accepted in active_map.items():
        chosen = learner_map.get(left, "")
        rows.append({
            "left": left,
            "learner_match": chosen or "No answer",
            "correct_match": accepted,
            "is_correct": chosen == accepted
        })
    return rows, active_map


def pick_unique_bank_question_ids(cursor, question_type, needed, modules=None, subjects=None, used_question_texts=None):
    modules = parse_module_names(",".join(modules or []))
    subjects = [item.strip() for item in (subjects or []) if item and item.strip()]
    used_question_texts = {item.strip().lower() for item in (used_question_texts or set()) if item}
    where_parts = ["question_type = ?"]
    params = [question_type]
    if modules:
        module_filters = " OR ".join(["LOWER(COALESCE(modules, '')) LIKE ?" for _ in modules])
        where_parts.append(f"({module_filters})")
        params.extend(f"%{module.lower()}%" for module in modules)
    if subjects:
        subject_filters = " OR ".join(["LOWER(COALESCE(subject, '')) = ?" for _ in subjects])
        where_parts.append(f"({subject_filters})")
        params.extend(subject.lower() for subject in subjects)
    cursor.execute(f"""
        SELECT id, question_text
        FROM question_bank_questions
        WHERE {' AND '.join(where_parts)}
        ORDER BY RANDOM()
    """, params)
    picked_ids = []
    seen_texts = set(used_question_texts)
    seen_match_pairs = set()
    for bank_question_id, question_text in cursor.fetchall():
        normalized_text = (question_text or "").strip().lower()
        if question_type != "match" and normalized_text in seen_texts:
            continue
        if question_type == "match":
            cursor.execute("""
                SELECT option_text, match_pair
                FROM question_bank_options
                WHERE bank_question_id = ?
                ORDER BY id
                LIMIT 1
            """, (bank_question_id,))
            option_row = cursor.fetchone()
            if not option_row:
                continue
            pair_signature = (
                (option_row[0] or "").strip().lower(),
                (option_row[1] or "").strip().lower(),
            )
            if pair_signature in seen_match_pairs:
                continue
            seen_match_pairs.add(pair_signature)
        else:
            seen_texts.add(normalized_text)
        picked_ids.append(bank_question_id)
        if len(picked_ids) >= needed:
            break
    return picked_ids, seen_texts


def clone_bank_question_to_test(cursor, bank_question_id, test_id, order_index):
    cursor.execute("""
        SELECT question_text, question_type, marks, subject, modules
        FROM question_bank_questions
        WHERE id = ?
    """, (bank_question_id,))
    row = cursor.fetchone()
    if not row:
        return None
    question_text, question_type, marks, subject, modules = row
    cursor.execute("""
        INSERT INTO theory_questions (test_id, question_text, question_type, marks, order_index, bank_question_id, source_subject, source_modules)
        VALUES (?, ?, ?, ?, ?, ?, ?, ?)
    """, (test_id, question_text, question_type, marks, order_index, bank_question_id, subject, modules))
    new_question_id = cursor.lastrowid
    cursor.execute("""
        SELECT option_text, is_correct, match_pair
        FROM question_bank_options
        WHERE bank_question_id = ?
        ORDER BY id
    """, (bank_question_id,))
    for option_text, is_correct, match_pair in cursor.fetchall():
        cursor.execute("""
            INSERT INTO theory_options (question_id, option_text, is_correct, match_pair)
            VALUES (?, ?, ?, ?)
        """, (new_question_id, option_text, is_correct, match_pair))
    return new_question_id


def create_generated_match_question(cursor, bank_question_ids, test_id, order_index):
    merged_count, _ = merge_bank_match_rows_into_test(cursor, bank_question_ids, test_id, order_index)
    return merged_count


def merge_bank_match_rows_into_test(cursor, bank_question_ids, test_id, order_index):
    if not bank_question_ids:
        return 0, order_index

    placeholders = ",".join("?" for _ in bank_question_ids)
    cursor.execute(f"""
        SELECT qbq.id, qbq.question_text, qbo.option_text, qbo.is_correct, qbo.match_pair
        FROM question_bank_questions qbq
        JOIN question_bank_options qbo ON qbo.bank_question_id = qbq.id
        WHERE qbq.id IN ({placeholders})
          AND qbq.question_type = 'match'
          AND COALESCE(qbo.match_pair, '') != ''
        ORDER BY qbq.id, qbo.id
    """, bank_question_ids)
    rows = cursor.fetchall()
    if not rows:
        return 0, order_index

    all_pairs = []
    seen_pairs = set()
    question_labels = []
    source_subjects = set()
    source_modules_set = set()
    for bank_question_id, question_text, option_text, is_correct, match_pair in rows:
        if question_text:
            question_labels.append(question_text)
        cursor.execute("SELECT subject, modules FROM question_bank_questions WHERE id = ?", (bank_question_id,))
        source_row = cursor.fetchone()
        if source_row:
            if source_row[0]:
                source_subjects.add(source_row[0])
            source_modules_set.update(parse_module_names(source_row[1] or ""))
        pair_signature = ((option_text or "").strip().lower(), (match_pair or "").strip().lower())
        if not pair_signature[0] or not pair_signature[1] or pair_signature in seen_pairs:
            continue
        seen_pairs.add(pair_signature)
        all_pairs.append((option_text, is_correct, match_pair))

    if not all_pairs:
        return 0, order_index

    cursor.execute("""
        SELECT id, COALESCE(source_modules, '')
        FROM theory_questions
        WHERE test_id = ?
          AND question_type = 'match'
        ORDER BY order_index, id
        LIMIT 1
    """, (test_id,))
    existing_question = cursor.fetchone()

    if existing_question:
        question_id = existing_question[0]
        source_modules_set.update(parse_module_names(existing_question[1] or ""))
    else:
        source_subject = next(iter(source_subjects), "")
        source_modules = ", ".join(sorted(source_modules_set, key=str.lower))
        question_text = question_labels[0] if question_labels else "Match Column A to B"
        cursor.execute("""
            INSERT INTO theory_questions (test_id, question_text, question_type, marks, order_index, source_subject, source_modules)
            VALUES (?, ?, 'match', 0, ?, ?, ?)
        """, (test_id, question_text, order_index, source_subject, source_modules))
        question_id = cursor.lastrowid
        order_index += 1

    cursor.execute("""
        SELECT option_text, match_pair
        FROM theory_options
        WHERE question_id = ?
          AND COALESCE(match_pair, '') != ''
    """, (question_id,))
    existing_pairs = {
        ((option_text or "").strip().lower(), (match_pair or "").strip().lower())
        for option_text, match_pair in cursor.fetchall()
    }

    inserted = 0
    for option_text, is_correct, match_pair in all_pairs:
        pair_signature = ((option_text or "").strip().lower(), (match_pair or "").strip().lower())
        if pair_signature in existing_pairs:
            continue
        cursor.execute("""
            INSERT INTO theory_options (question_id, option_text, is_correct, match_pair)
            VALUES (?, ?, ?, ?)
        """, (question_id, option_text, is_correct, match_pair))
        existing_pairs.add(pair_signature)
        inserted += 1

    cursor.execute("""
        SELECT COUNT(*)
        FROM theory_options
        WHERE question_id = ?
          AND COALESCE(match_pair, '') != ''
    """, (question_id,))
    total_pairs = cursor.fetchone()[0] or 0
    cursor.execute("""
        UPDATE theory_questions
        SET marks = ?, source_modules = ?
        WHERE id = ?
    """, (total_pairs, ", ".join(sorted(source_modules_set, key=str.lower)), question_id))

    return inserted, order_index


def regrade_theory_question_answers(test_id, question_id, selected_group=None, learner_username=None):
    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("""
        SELECT question_type, marks
        FROM theory_questions
        WHERE id = ? AND test_id = ?
    """, (question_id, test_id))
    question_row = cursor.fetchone()
    if not question_row:
        conn.close()
        return 0

    q_type, marks = question_row
    cursor.execute("""
        SELECT id, option_text, is_correct, match_pair
        FROM theory_options
        WHERE question_id = ?
        ORDER BY id
    """, (question_id,))
    options = cursor.fetchall()

    query = """
        SELECT a.id, a.answer_text, a.submission_id
        FROM theory_answers a
        JOIN theory_submissions s ON s.id = a.submission_id
        JOIN users u ON u.username = s.username
        WHERE s.test_id = ? AND a.question_id = ?
    """
    params = [test_id, question_id]
    if selected_group:
        query += " AND u.group_name = ?"
        params.append(selected_group)
    if learner_username:
        query += " AND s.username = ?"
        params.append(learner_username)
    cursor.execute(query, params)
    answer_rows = cursor.fetchall()

    updated = 0
    submission_ids = set()
    for answer_id, answer_text, submission_id in answer_rows:
        awarded = compute_theory_answer_award(q_type, marks, options, answer_text)
        is_correct = 1 if awarded == marks else 0
        if q_type == "match":
            expected_total = len([option for option in options if len(option) > 3 and option[3] and option[3] != "correction"])
            is_correct = 1 if awarded == expected_total else 0
        cursor.execute("""
            UPDATE theory_answers
            SET marks_awarded = ?, is_correct = ?
            WHERE id = ?
        """, (awarded, is_correct, answer_id))
        updated += cursor.rowcount
        submission_ids.add(submission_id)

    for submission_id in submission_ids:
        cursor.execute("""
            SELECT COALESCE(SUM(marks_awarded), 0)
            FROM theory_answers
            WHERE submission_id = ?
        """, (submission_id,))
        score = cursor.fetchone()[0] or 0
        cursor.execute("SELECT total FROM theory_submissions WHERE id = ?", (submission_id,))
        total_row = cursor.fetchone()
        total = total_row[0] if total_row and total_row[0] else 0
        percentage = round((score / total) * 100) if total else 0
        cursor.execute("""
            UPDATE theory_submissions
            SET score = ?, percentage = ?
            WHERE id = ?
        """, (score, percentage, submission_id))

    conn.commit()
    conn.close()
    return updated


def add_learner_note_entry(cursor, username, note, created_by, flag=""):
    cursor.execute("""
        INSERT INTO learner_notes (username, note, flag, created_by, created_at)
        VALUES (?, ?, ?, ?, ?)
    """, (username, note, flag, created_by, datetime.now().isoformat()))


def get_group_late_threshold(group, date):
    conn = get_db()
    cursor = conn.cursor()

    cursor.execute("""
    SELECT MIN(login_time)
    FROM login_history
    WHERE date = ?
      AND username IN (
          SELECT username FROM users WHERE group_name = ?
      )
    """, (date, group))
    min_login = cursor.fetchone()[0]
    conn.close()

    if not min_login:
        return None

    try:
        min_dt = datetime.fromisoformat(min_login)
    except ValueError:
        min_dt = datetime.strptime(min_login, "%Y-%m-%d %H:%M:%S")

    late_cutoff = min_dt + timedelta(minutes=4)
    return late_cutoff.strftime("%H:%M")


def is_attendance_editable(date_str, role="teacher"):
    if role == "admin":
        return True
    try:
        date_obj = datetime.strptime(date_str, "%Y-%m-%d").date()
        days_diff = (datetime.now().date() - date_obj).days
        return days_diff <= 10
    except ValueError:
        return False


def get_attendance_data(group, start_date=None, end_date=None, teacher_username=None):
    conn = get_db()
    cursor = conn.cursor()

    if not start_date or not end_date:
        term_range = get_active_term_range()
        if term_range:
            start_date, end_date = term_range
            current = datetime.strptime(start_date, "%Y-%m-%d").date()
            end = datetime.strptime(end_date, "%Y-%m-%d").date()
            days = []
            while current <= end:
                if current.weekday() < 5:
                    days.append(current.strftime("%Y-%m-%d"))
                current += timedelta(days=1)
        else:
            days = get_last_7_days()
    else:
        days = []
        current = datetime.strptime(start_date, "%Y-%m-%d").date()
        end = datetime.strptime(end_date, "%Y-%m-%d").date()
        while current <= end:
            if current.weekday() < 5:
                days.append(current.strftime("%Y-%m-%d"))
            current += timedelta(days=1)

    today_str = datetime.now().date().isoformat()
    days = [day for day in days if day <= today_str]

    auto_exclude_empty_attendance_days(cursor, [group], created_by=teacher_username or "system", days=days)
    cursor.execute("""
    SELECT date FROM excluded_dates 
    WHERE group_name IS NULL OR group_name = ?
    """, (group,))
    excluded_dates = {row[0] for row in cursor.fetchall()}
    days = [day for day in days if day not in excluded_dates]

    year_days = [day for day in get_current_year_attendance_days() if day <= today_str and day not in excluded_dates]
    auto_exclude_empty_attendance_days(cursor, [group], created_by=teacher_username or "system", days=year_days)
    cursor.execute("""
    SELECT date FROM excluded_dates 
    WHERE group_name IS NULL OR group_name = ?
    """, (group,))
    excluded_dates = {row[0] for row in cursor.fetchall()}
    days = [day for day in days if day not in excluded_dates]
    year_days = [day for day in get_current_year_attendance_days() if day <= today_str and day not in excluded_dates]

    query = """
    SELECT username, full_name, group_name
    FROM users
    WHERE group_name = ? AND role = 'student'
    """
    params = [group]
    if teacher_username:
        query += " AND teacher_username = ?"
        params.append(teacher_username)

    cursor.execute(query, params)
    learners = cursor.fetchall()

    attendance = []
    late_cutoffs = fetch_group_late_thresholds(cursor, group, days)
    year_late_cutoffs = fetch_group_late_thresholds(cursor, group, year_days)
    usernames = [user for user, _, _ in learners]
    login_map = fetch_first_login_times(cursor, usernames, days)
    override_map = fetch_attendance_override_statuses(cursor, usernames, days)
    year_login_map = fetch_first_login_times(cursor, usernames, year_days)
    year_override_map = fetch_attendance_override_statuses(cursor, usernames, year_days)
    cursor.execute("""
        SELECT DISTINCT lh.date
        FROM login_history lh
        JOIN users u ON u.username = lh.username
        WHERE u.group_name = ? AND u.role = 'student'
    """, (group,))
    class_checked_dates = {row[0] for row in cursor.fetchall()}

    for user, name, user_group_name in learners:
        history = build_attendance_history(
            cursor,
            user,
            group,
            days,
            login_map=login_map,
            override_map=override_map,
            late_cutoffs=late_cutoffs,
            class_checked_dates=class_checked_dates,
        )
        year_history = build_attendance_history(
            cursor,
            user,
            group,
            year_days,
            login_map=year_login_map,
            override_map=year_override_map,
            late_cutoffs=year_late_cutoffs,
            class_checked_dates=class_checked_dates,
        )
        summary = summarize_attendance_history(history)
        year_summary = summarize_attendance_history(year_history)
        row = {
            "username": user,
            "name": name,
            "group": user_group_name,
            "days": {},
            "attendance_pct": summary["attendance_pct"],
            "present_days": summary["present_days"],
            "absent_days": year_summary["absent_days"],
            "year_present_days": year_summary["present_days"],
            "year_attendance_pct": year_summary["attendance_pct"],
            "year_total_days": year_summary["total_days"],
            "late_days": summary["late_days"],
        }
        for item in history:
            if item["status"] == "Present":
                row["days"][item["date"]] = {
                    "time": item["time"],
                    "late": item["late"],
                    "manual": item["note"] == "Manual"
                }
            elif item["status"] == "Absent":
                row["days"][item["date"]] = None
            else:
                row["days"][item["date"]] = None

        attendance.append(row)

    conn.commit()
    conn.close()
    return days, attendance


def calculate_attendance_percentage(data, days):
    total_cells = len(data) * len(days)
    present = 0

    for row in data:
        for d in days:
            if row["days"].get(d):
                present += 1

    return round((present / total_cells) * 100) if total_cells else 0


def pptx_to_content_slide_html(uploaded_file):
    from pptx import Presentation

    prs = Presentation(uploaded_file)
    slide_w = float(prs.slide_width)
    slide_h = float(prs.slide_height)
    html_slides = []

    for slide in prs.slides:
        boxes = []
        for shape in slide.shapes:
            left = max(0, (float(shape.left) / slide_w) * 100)
            top = max(0, (float(shape.top) / slide_h) * 100)
            width = max(5, (float(shape.width) / slide_w) * 100)
            height = max(5, (float(shape.height) / slide_h) * 100)
            style = f"left:{left:.2f}%;top:{top:.2f}%;width:{width:.2f}%;height:{height:.2f}%;"

            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if not text:
                    continue
                paragraphs = "".join(
                    f"<p>{escape(line)}</p>"
                    for line in text.splitlines()
                    if line.strip()
                )
                boxes.append(f'<div class="slide-box text-box" style="{style}">{paragraphs}</div>')
                continue

            if hasattr(shape, "image"):
                image = shape.image
                ext = (image.ext or "png").lower()
                mime = mimetypes.types_map.get(f".{ext}", "image/png")
                encoded = base64.b64encode(image.blob).decode("ascii")
                src = f"data:{mime};base64,{encoded}"
                boxes.append(f'<div class="slide-box image-box" style="{style}"><img src="{src}" alt=""></div>')

        if boxes:
            html_slides.append("".join(boxes))

    return html_slides


# ── Marking helper ────────────────────────────────────────────────────────────

def mark_file(filepath, marking_script, marking_setup_id=None):
    if not marking_script:
        return {"task_name": "Unknown Task", "score": 0, "total": 0, "percentage": 0, "results": [],
                "error": "No marking script assigned to this task. Please contact your teacher."}
    try:
        module = importlib.import_module(f"marking.tasks.{marking_script}")
        if marking_setup_id is not None and hasattr(module, "mark_with_setup"):
            return module.mark_with_setup(filepath, int(marking_setup_id))
        return module.mark(filepath)
    except ModuleNotFoundError:
        return {"task_name": marking_script, "score": 0, "total": 0, "percentage": 0, "results": [],
                "error": f"Marking script '{marking_script}' not found. Please contact your teacher."}


def get_marking_scripts():
    tasks_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "marking", "tasks")
    scripts = []
    if os.path.exists(tasks_dir):
        for f in sorted(os.listdir(tasks_dir)):
            if f.endswith(".py") and f != "__init__.py":
                scripts.append(f[:-3])
    return scripts
